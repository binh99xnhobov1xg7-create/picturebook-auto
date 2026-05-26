"""快速检查生成的 pptx 是否符合规范。"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

p = Path(sys.argv[1] if len(sys.argv) > 1 else
         Path(__file__).resolve().parents[1] / "outputs/Visiting_Scotland/Visiting_Scotland.pptx")

prs = Presentation(str(p))
print(f"file: {p.name}")
print(f"size: {p.stat().st_size // 1024} KB")
print(f"slides: {len(prs.slides)}  ({prs.slide_width.inches:.1f} x {prs.slide_height.inches:.1f} in)")

for i, slide in enumerate(prs.slides, start=1):
    pic = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    shapes = list(slide.shapes)
    texts = []
    fonts = set()
    for s in shapes:
        if not s.has_text_frame:
            continue
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts.add(run.font.name)
                if run.text.strip():
                    texts.append(run.text.strip()[:55])
    print(f"  p{i}: shapes={len(shapes)} pics={pic} fonts={fonts or '-'}")
    for t in texts[:4]:
        print(f"       · {t}")
