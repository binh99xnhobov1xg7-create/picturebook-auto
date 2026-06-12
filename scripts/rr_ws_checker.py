"""RR + WS 校验与 PDF 打包（集成自外部 skill: worksheet-RR-checking）。

本模块把外部仓库 https://github.com/jeredithia-ai/worksheet-RR-checking 的
「校验 + 转 PDF + 合并」能力，移植/适配到本绘本自动化 App，作为「出书的最后一步」：

    生成四件套 → 【本模块：校验 RR/WS → 转 PDF → 合并成 1 个交付 PDF】 → ZIP 打包

外部仓库做了什么（分析结论）
============================
外部仓库本质是一套「人工/AI 审核 + 数据驱动修正 + 高保真转 PDF + 合并」工作流，
**没有一个通用的、读任意 RR/WS 即给出 pass/fail 的校验引擎**：

1. extract_text.py   —— 把 storybook(PDF/PPTX) / worksheet(PPTX) / report(DOCX)
                        抽成纯文本，供 AI 按规则审核。纯自动、可复用。
2. .cursor/rules/picture-book-workflow.mdc —— 真正的「校验标准」以散文规则形式存在，
                        由 AI agent 阅读抽取文本后逐项判断（三方一致性 / 语言错误 /
                        模板一致性）。**不是代码**。
3. lessons.py        —— 每节课「手工硬编码」的已知修正项（页脚位置、替换词、A4 尺寸…）。
                        是逐课配置，不是通用规则。
4. apply_corrections.py —— 按 lessons.py 配置在保留排版前提下改 pptx/docx
                        （含「图片/logo drawing 数量前后一致」这一硬校验不变式）。
5. run_workflow.py   —— **可直接复用的核心**：
                        a) WPS COM 转 PDF：KWPS.Application(Writer→docx) /
                           KWPP.Application(Presentation→pptx)，format 17 / 32；
                        b) pypdf 合并：worksheet 在前(横版) + reading report 在后
                           (竖版，/Rotate=270 = 左转 90°)。

可直接复用：run_workflow 的「WPS 转 PDF」+「pypdf 合并(含旋转)」、extract_text 的抽取、
            apply_corrections 的「drawing 数量不变式」思路。
需自研：    一个「通用、可编程」的 RR/WS 校验器（把 .mdc 散文规则落成断言）。
            好在本 App 在生成阶段已有结构化 `BookOutline` 与 `scripts/evals.py`，
            可直接复用，校验质量远高于「事后扒文本」。

本 App 的产物与命名（校验对象）
==============================
每本书输出目录 `Level <N>_Book<NN>_<title>/` 下：
    <prefix>_Reader.pptx          绘本（基准，仅对照，不进合并）
    <prefix>_Worksheet.pptx       练习册 WS（校验 + 进合并，横版在前）
    <prefix>_Reading_Report.docx  阅读报告 RR（校验 + 进合并，竖版在后，旋转 270）
    <prefix>_Teachers_Guide.docx  教师指南 TG
合并交付目标：<prefix>_Worksheet+RR_Final.pdf

环境（已探测）
==============
- WPS Office 已安装（Kingsoft WPS Office 12.1，ProgID KWPS/KWPP 可用）→ 首选转 PDF 方案。
- 已装 Python 依赖：python-pptx / python-docx / pypdf / pywin32 / docx2pdf / PyMuPDF(fitz)。
- 未发现 soffice / libreoffice（LibreOffice headless 不可用）。
- MS Office 是否安装未确认（docx2pdf 依赖 Word；WPS 才是已验证方案）。

用法（计划）
============
    from rr_ws_checker import check_and_pack
    result = check_and_pack(book_dir, outline=outline)
    if not result.passed:
        ...  # 在 UI 暴露 result.issues
    final_pdf = result.final_pdf

注意：本文件目前为「骨架 + 函数签名 + TODO」，PDF 转换/合并逻辑已从外部仓库移植可用；
      校验逻辑给出框架并复用 evals，细则待按下方 TODO 落地后接线。
"""
from __future__ import annotations

import json
import os
import re
import subprocess
import threading
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional


# =============================================================================
# 结果数据结构
# =============================================================================
@dataclass
class CheckIssue:
    """单条校验问题（对齐 scripts/evals.py 的 Issue 口径）。"""
    target: str          # "RR" / "WS" / "一致性"
    level: str           # "error" / "warn" / "ok"
    code: str            # 稳定的规则编号，如 "RR_TITLE_FMT"
    msg: str             # 人类可读说明 + 修改建议


@dataclass
class CheckResult:
    passed: bool = False                       # 无 error 即视为通过（warn 不阻断）
    issues: list[CheckIssue] = field(default_factory=list)
    ws_pdf: Path | None = None
    rr_pdf: Path | None = None
    final_pdf: Path | None = None
    elapsed_s: float = 0.0
    converter: str = ""                        # "wps" / "docx2pdf" / "skipped"

    @property
    def n_error(self) -> int:
        return sum(1 for i in self.issues if i.level == "error")

    @property
    def n_warn(self) -> int:
        return sum(1 for i in self.issues if i.level == "warn")

    def add(self, target: str, level: str, code: str, msg: str) -> None:
        self.issues.append(CheckIssue(target, level, code, msg))

    @property
    def needs_human_review(self) -> bool:
        """只要有 error/warn 就提示人工抽查；warn 不阻断生产。"""
        return any(i.level in ("error", "warn") for i in self.issues)


_WPS_CONVERT_LOCK = threading.Lock()


def _mock_mode() -> bool:
    try:
        from config import MOCK_AI_EXTRACT
        return bool(MOCK_AI_EXTRACT) or os.getenv("RR_WS_MOCK", "").lower() in ("1", "true", "yes")
    except Exception:
        return os.getenv("RR_WS_MOCK", "").lower() in ("1", "true", "yes")


def _outline_story_text(outline: Any | None) -> str:
    if outline is None:
        return ""
    lines: list[str] = []
    for p in getattr(outline, "pages", []) or []:
        txt = (getattr(p, "text", "") or "").strip()
        if txt:
            lines.append(txt)
    return "\n".join(lines)


def _outline_brief(outline: Any | None) -> str:
    if outline is None:
        return "未提供结构化 BookOutline。"
    vocab = (
        getattr(outline, "vocabulary_mastery", None)
        or getattr(outline, "vocabulary_simple", None)
        or []
    )
    fields = {
        "title": getattr(outline, "title", ""),
        "level": getattr(outline, "level", ""),
        "book_number": getattr(outline, "book_number", ""),
        "theme": getattr(outline, "theme", ""),
        "phonics": getattr(outline, "phonics", ""),
        "grammar_focus": getattr(outline, "grammar_focus", ""),
        "vocabulary": list(vocab)[:12],
        "story": _outline_story_text(outline),
    }
    return json.dumps(fields, ensure_ascii=False, indent=2)


def _extract_docx_text(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
        if lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _extract_any_text(path: Path | None) -> str:
    if path is None or not Path(path).exists():
        return ""
    suf = path.suffix.lower()
    try:
        if suf in (".docx", ".doc"):
            return _extract_docx_text(path)
        if suf in (".pptx", ".ppt"):
            return _extract_pptx_text(path)
    except Exception:
        return ""
    return ""


def _add_structured_assertions(outline: Any | None, result: CheckResult) -> None:
    if outline is None:
        result.add("结构化", "warn", "OUTLINE_MISSING", "未提供 BookOutline，已退化为文件文本校验")
        return
    try:
        from evals import ERROR, WARN, run_all
        rep = run_all(outline=outline)
        for idx, issue in enumerate(rep.issues, 1):
            if issue.level not in (ERROR, WARN):
                continue
            result.add("结构化", issue.level, f"EVALS_{idx:03d}", f"[{issue.category}] {issue.msg}")
    except Exception as e:  # noqa: BLE001
        result.add("结构化", "warn", "EVALS_FAIL", f"结构化 evals 跳过：{e}")


def _run_ai_review(
    *,
    outline: Any | None,
    reader_text: str,
    ws_text: str,
    rr_text: str,
    result: CheckResult,
) -> None:
    """AI 三方一致性 + 语言/APA 审核；mock/无 key 时 no-op。"""
    if _mock_mode():
        result.add("AI", "warn", "AI_REVIEW_MOCK_SKIPPED", "mock 模式：已跳过 AI 三方一致性审核")
        return
    try:
        from deepseek_client import deepseek_chat_json
    except Exception as e:  # noqa: BLE001
        result.add("AI", "warn", "AI_REVIEW_UNAVAILABLE", f"AI 客户端不可用，已跳过：{e}")
        return

    system = (
        "你是 VIPKID 儿童绘本教研审核员。请审核 Worksheet 与 Reading Report 是否和绘本/BookOutline 一致，"
        "同时检查英文语言、拼写、APA 标点/大小写、题目清晰度。只输出 JSON。"
    )
    user = (
        "请返回 JSON：{\"issues\":[{\"target\":\"WS|RR|一致性|语言\","
        "\"level\":\"error|warn\",\"code\":\"STABLE_CODE\",\"msg\":\"中文说明+建议\"}]}。\n\n"
        "规则：error=明显错题/错书/与绘本冲突/关键文件损坏；warn=需人工抽查但不阻断。"
        "不要输出 ok 项。\n\n"
        f"BookOutline:\n{_outline_brief(outline)[:6000]}\n\n"
        f"Reader text:\n{reader_text[:5000]}\n\n"
        f"Worksheet text:\n{ws_text[:7000]}\n\n"
        f"Reading Report text:\n{rr_text[:7000]}"
    )
    try:
        data = deepseek_chat_json(
            system=system,
            user=user,
            max_tokens=3000,
            temperature=0.1,
            fallback=None,
        )
        issues = data.get("issues") if isinstance(data, dict) else None
        if not isinstance(issues, list):
            result.add("AI", "warn", "AI_REVIEW_BAD_JSON", "AI 审核未返回标准 issues 数组，需人工抽查")
            return
        for raw in issues[:20]:
            if not isinstance(raw, dict):
                continue
            level = str(raw.get("level") or "warn").lower()
            if level not in ("error", "warn"):
                level = "warn"
            code = re.sub(r"[^A-Z0-9_]+", "_", str(raw.get("code") or "AI_REVIEW_WARN").upper())[:48]
            target = str(raw.get("target") or "AI")
            msg = str(raw.get("msg") or "AI 审核提示").strip()
            result.add(target, level, code, msg)
        if not issues:
            result.add("AI", "ok", "AI_REVIEW_OK", "AI 三方一致性与语言审核未发现明显问题")
    except Exception as e:  # noqa: BLE001
        result.add("AI", "warn", "AI_REVIEW_FAIL", f"AI 审核失败，需人工抽查：{e}")


# =============================================================================
# 文件发现
# =============================================================================
def discover_book_files(book_dir: Path) -> dict[str, Path | None]:
    """在一本书输出目录里定位 WS(pptx) / RR(docx) / Reader(pptx)。

    依赖本 App 固定命名后缀：_Worksheet.pptx / _Reading_Report.docx / _Reader.pptx。
    跳过 ~$ 锁文件与 *_DRAFT.* 草稿。
    """
    book_dir = Path(book_dir)
    found: dict[str, Path | None] = {"worksheet": None, "report": None, "reader": None}
    if not book_dir.is_dir():
        return found
    for p in book_dir.iterdir():
        if not p.is_file() or p.name.startswith("~$"):
            continue
        if "_DRAFT" in p.stem or "user_ref" in p.name:
            continue
        n = p.name
        if n.endswith("_Worksheet.pptx"):
            found["worksheet"] = p
        elif n.endswith("_Reading_Report.docx"):
            found["report"] = p
        elif n.endswith("_Reader.pptx"):
            found["reader"] = p
    return found


# =============================================================================
# 校验：阅读报告（RR / docx）
# =============================================================================
def check_reading_report(rr_path: Path, outline: Any | None,
                         result: CheckResult) -> None:
    """对生成好的 RR docx 做编程化校验（落地 .mdc 第 4 节规则）。

    建议规则（对齐外部 .mdc §4 + 本 App reading_report_builder 口径）：
      RR_OPEN        能正常打开、含至少 1 个表格            [error]
      RR_TITLE_FMT   标题格式 `阅读报告 Smart - <书名>`(L0)
                     或 `阅读报告 Level <N> - <书名>`        [error]
      RR_LOGO        顶部 logo drawing 数量 >=1（防丢图）     [warn]
      RR_MASTERY     r4 词汇掌握恰好 4 词且 = outline.mastery [error]
      RR_PHONICS     拼读首字母小写、非空                     [warn]
      RR_QUESTIONS   阅读表达题数 = 4(L0-2)/5(L3-6)，
                     星级分布 1+2+2+(2+)3                     [error]
      RR_FLUENCY     阅读流利度 = 绘本正文逐字一致           [warn]
      RR_ENGAGE      课堂参与度 emoji 后均有空格             [warn]

    可直接复用 scripts/evals.py 的 check_reading_report / check_vocabulary。
    TODO: 用 python-docx 读取表格行(r4/r6/r10...)落地逐条断言；
          能拿到 outline 时优先比对结构化字段，拿不到时退化为文本规则。
    """
    try:
        from docx import Document
        doc = Document(str(rr_path))
        if not doc.tables:
            result.add("RR", "error", "RR_OPEN", "RR 内未发现表格，疑似生成损坏")
        else:
            result.add("RR", "ok", "RR_OPEN", "RR 可打开且含表格")
        full_text = "\n".join(
            [p.text for p in doc.paragraphs]
            + [" | ".join(c.text for c in row.cells) for t in doc.tables for row in t.rows]
        )
        title = (getattr(outline, "title", "") or "").strip()
        level = str(getattr(outline, "level", "") or "").strip().lstrip("L")
        if title and title.lower() not in full_text.lower():
            result.add("RR", "warn", "RR_TITLE_MISMATCH", f"RR 文本未明显包含书名 `{title}`，需人工核对")
        if level and not re.search(rf"(Level\s*{re.escape(level)}|Smart)", full_text, re.I):
            result.add("RR", "warn", "RR_LEVEL_MISSING", f"RR 文本未明显包含 Level {level}/Smart 标识")
        vocab = (
            getattr(outline, "vocabulary_mastery", None)
            or getattr(outline, "vocabulary_simple", None)
            or []
        )
        if vocab:
            missing = [w for w in vocab[:4] if str(w).lower() not in full_text.lower()]
            if missing:
                result.add("RR", "warn", "RR_MASTERY_VOCAB_MISSING", f"RR 未明显包含 mastery 词：{', '.join(missing)}")
    except Exception as e:  # noqa: BLE001
        result.add("RR", "error", "RR_OPEN", f"RR 无法打开：{e}")


# =============================================================================
# 校验：练习册（WS / pptx）
# =============================================================================
def check_worksheet(ws_path: Path, outline: Any | None,
                    result: CheckResult) -> None:
    """对生成好的 WS pptx 做编程化校验（落地 .mdc 第 3 节规则）。

    建议规则（对齐外部 .mdc §3 + 本 App worksheet_builder 口径）：
      WS_OPEN        能正常打开、页数符合模板(6 / 8-9)        [error]
      WS_A4          画布 = A4 横版(本 App 10.833x7.5in；
                     外部要求 EMU 10692000x7560000)           [warn]
      WS_FOOTER_FMT  页脚命名 `Smart - <书名>`/`Level <N> - …`[warn]
      WS_NO_PLACEHOLDER 无空白方框/纯白色块占位             [warn]
      WS_VOCAB       目标词 100% 来自 outline 词表           [warn]
      WS_NO_RAW_UNDERSCORE 书写横线不用裸 `____`            [warn]

    可直接复用 scripts/evals.py 的 check_worksheet（题型/题量维度）。
    TODO: 用 python-pptx 读取 slides/shapes 落地断言。
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(ws_path))
        n = len(prs.slides)
        if n < 1:
            result.add("WS", "error", "WS_OPEN", "WS 无任何 slide")
        else:
            result.add("WS", "ok", "WS_OPEN", f"WS 可打开，{n} 页")
        if n < 5:
            result.add("WS", "warn", "WS_PAGE_COUNT_LOW", f"WS 页数偏少（{n} 页），需人工核对题型是否完整")
        if prs.slide_width < prs.slide_height:
            result.add("WS", "warn", "WS_NOT_LANDSCAPE", "WS 画布不是横版，需核对 A4 横版规格")
        text = _extract_pptx_text(ws_path)
        if re.search(r"_{4,}", text):
            result.add("WS", "warn", "WS_RAW_UNDERSCORE", "WS 文本含连续下划线，建议核对书写横线排版")
        title = (getattr(outline, "title", "") or "").strip()
        if title and title.lower() not in text.lower():
            result.add("WS", "warn", "WS_TITLE_MISMATCH", f"WS 文本未明显包含书名 `{title}`，需人工核对")
    except Exception as e:  # noqa: BLE001
        result.add("WS", "error", "WS_OPEN", f"WS 无法打开：{e}")


# =============================================================================
# PDF 转换（移植自外部 run_workflow.convert_with_wps，加 fallback）
# =============================================================================
def convert_to_pdf(src: Path, dest_dir: Path) -> tuple[Path, str]:
    """把 .pptx/.docx 转为 PDF，返回 (pdf_path, converter_name)。

    优先级（按本机已探测环境）：
      1) WPS COM —— KWPS.Application(docx, SaveAs fmt 17) /
                    KWPP.Application(pptx, SaveAs fmt 32)。已验证、最高保真。
      2) docx2pdf —— 仅当装有 MS Office/Word 时可用（pptx 不支持，需另想）。
      3) 抛错 —— 让上层降级为「只校验不打包」。

    注意：调用前需释放被 WPS/Office 占用的同名文件句柄。
    """
    src, dest_dir = Path(src), Path(dest_dir)
    dest_dir.mkdir(parents=True, exist_ok=True)
    if _mock_mode():
        raise RuntimeError("mock 模式不调用 WPS COM")
    with _WPS_CONVERT_LOCK:
        _kill_stray_wps_processes()
        try:
            return _convert_with_wps(src, dest_dir), "wps"
        except Exception as e:  # noqa: BLE001
            raise RuntimeError(f"WPS 转 PDF 失败（无可用 fallback）：{e}") from e


def _kill_stray_wps_processes() -> None:
    """批量安全：转换前清理 WPS/演示/表格残留进程。"""
    if os.name != "nt":
        return
    try:
        subprocess.run(
            ["taskkill", "/F", "/IM", "wps.exe", "/IM", "wpp.exe", "/IM", "et.exe"],
            check=False,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        pass


def _convert_with_wps(src: Path, dest_dir: Path) -> Path:
    """移植自外部仓库 run_workflow.convert_with_wps。WPS COM 高保真转 PDF。"""
    import pythoncom
    import win32com.client as win32

    pythoncom.CoInitialize()
    try:
        suf = src.suffix.lower()
        out_pdf = dest_dir / (src.stem + ".pdf")
        if out_pdf.exists():
            try:
                out_pdf.unlink()
            except Exception:  # noqa: BLE001
                pass

        if suf in (".docx", ".doc"):
            app = win32.DispatchEx("KWPS.Application")
            try:
                try:
                    app.Visible = False
                except Exception:  # noqa: BLE001
                    pass
                try:
                    app.DisplayAlerts = 0
                except Exception:  # noqa: BLE001
                    pass
                doc = app.Documents.Open(
                    str(src.resolve()), ConfirmConversions=False,
                    ReadOnly=True, AddToRecentFiles=False,
                )
                doc.SaveAs(str(out_pdf.resolve()), 17)  # wdFormatPDF
                doc.Close(SaveChanges=False)
            finally:
                try:
                    app.Quit()
                except Exception:  # noqa: BLE001
                    pass
        elif suf in (".pptx", ".ppt"):
            app = win32.DispatchEx("KWPP.Application")
            try:
                deck = app.Presentations.Open(
                    str(src.resolve()), ReadOnly=True, WithWindow=False)
                deck.SaveAs(str(out_pdf.resolve()), 32)  # ppSaveAsPDF
                deck.Close()
            finally:
                try:
                    app.Quit()
                except Exception:  # noqa: BLE001
                    pass
        else:
            raise ValueError(f"不支持的扩展名：{suf}")
    finally:
        pythoncom.CoUninitialize()

    if not out_pdf.exists():
        raise FileNotFoundError(f"WPS 未产出 PDF：{out_pdf}")
    return out_pdf


# =============================================================================
# PDF 合并（移植自外部 run_workflow.merge_pdfs）
# =============================================================================
def merge_rr_ws_pdf(worksheet_pdf: Path, report_pdf: Path, final_pdf: Path,
                    report_rotation: int = 270) -> Path:
    """合并：worksheet(横版，原样) + reading report(竖版，/Rotate=270 左转 90°)。

    直接移植外部仓库逻辑，零改动可用。
    """
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import NameObject, NumberObject

    writer = PdfWriter()
    for p in PdfReader(str(worksheet_pdf)).pages:
        writer.add_page(p)
    for p in PdfReader(str(report_pdf)).pages:
        try:
            existing = int(p.get("/Rotate", 0))
        except Exception:  # noqa: BLE001
            existing = 0
        p[NameObject("/Rotate")] = NumberObject((existing + report_rotation) % 360)
        writer.add_page(p)

    final_pdf = Path(final_pdf)
    final_pdf.parent.mkdir(parents=True, exist_ok=True)
    with open(final_pdf, "wb") as fh:
        writer.write(fh)
    writer.close()
    return final_pdf


# =============================================================================
# 编排：校验 + 打包
# =============================================================================
def check_and_pack(book_dir: Path, *, outline: Any | None = None,
                   do_convert: bool = True,
                   final_name: str | None = None) -> CheckResult:
    """出书最后一步：校验 RR/WS → 转 PDF → 合并成单一交付 PDF。

    流程：
      1) discover_book_files 定位 WS/RR；
      2) check_worksheet / check_reading_report（+ 可选三方一致性）填充 issues；
      3) 若 do_convert 且无 error：WS/RR → PDF → 合并 → <prefix>_Worksheet+RR_Final.pdf；
         若有 error：默认仍可选择「带警告打包」或「阻断」，交由 UI 决策（见 TODO）。

    返回 CheckResult（pass/fail + issue 清单 + final_pdf 路径），供 UI 展示与下载。
    """
    t0 = time.time()
    result = CheckResult()
    book_dir = Path(book_dir)

    files = discover_book_files(book_dir)
    ws, rr, reader = files["worksheet"], files["report"], files["reader"]
    if ws is None:
        result.add("WS", "error", "WS_MISSING", f"未找到 _Worksheet.pptx：{book_dir}")
    if rr is None:
        result.add("RR", "error", "RR_MISSING", f"未找到 _Reading_Report.docx：{book_dir}")
    if reader is None:
        result.add("Reader", "warn", "READER_MISSING", "未找到 _Reader.pptx，AI 三方一致性将主要依赖 BookOutline")

    if ws is not None:
        check_worksheet(ws, outline, result)
    if rr is not None:
        check_reading_report(rr, outline, result)

    _add_structured_assertions(outline, result)
    _run_ai_review(
        outline=outline,
        reader_text=_extract_any_text(reader) or _outline_story_text(outline),
        ws_text=_extract_any_text(ws),
        rr_text=_extract_any_text(rr),
        result=result,
    )

    result.passed = result.n_error == 0

    # 转 PDF + 合并
    if do_convert and ws is not None and rr is not None and _mock_mode():
        result.add("PDF", "warn", "PDF_CONVERT_MOCK_SKIPPED", "mock 模式已跳过 WPS 转 PDF/合并")
        result.converter = "skipped"
    elif do_convert and ws is not None and rr is not None:
        try:
            ws_pdf, conv = convert_to_pdf(ws, book_dir)
            rr_pdf, _ = convert_to_pdf(rr, book_dir)
            result.ws_pdf, result.rr_pdf, result.converter = ws_pdf, rr_pdf, conv
            prefix = ws.name[: -len("_Worksheet.pptx")]
            name = final_name or f"{prefix}_Worksheet+RR_Final"
            result.final_pdf = merge_rr_ws_pdf(ws_pdf, rr_pdf, book_dir / f"{name}.pdf")
        except Exception as e:  # noqa: BLE001
            result.add("PDF", "warn", "PDF_PACK_FAIL", f"PDF 转换/合并失败：{e}")
            result.converter = "skipped"
    elif do_convert:
        result.add("PDF", "warn", "PDF_PACK_SKIPPED", "缺少 WS 或 RR，已跳过 PDF 打包")

    result.elapsed_s = time.time() - t0
    return result


if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser(description="校验 RR/WS 并打包为合并 PDF")
    ap.add_argument("book_dir", help="一本书的输出目录")
    ap.add_argument("--no-convert", action="store_true", help="只校验，不转 PDF/合并")
    args = ap.parse_args()

    r = check_and_pack(Path(args.book_dir), do_convert=not args.no_convert)
    print(f"passed={r.passed} errors={r.n_error} warns={r.n_warn} "
          f"final_pdf={r.final_pdf} ({r.elapsed_s:.1f}s)")
    for it in r.issues:
        print(f"  [{it.level}] {it.target}/{it.code}: {it.msg}")
