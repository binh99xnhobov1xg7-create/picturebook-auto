"""交付物文档预览：把 docx / pptx 渲染成 PNG 页图，供网页内直接预览。

管线：doc/pptx --(LibreOffice headless)--> PDF --(PyMuPDF)--> 每页 PNG。
- 结果缓存在 <文档同级>/_preview/<stem>/，按源文件 mtime 失效，二次预览秒开。
- soffice 不可用时，提供 docx/pptx 的纯文字兜底抽取。
"""
from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path

_SOFFICE_CANDIDATES = [
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/soffice",
    "/usr/local/bin/soffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
]


def find_soffice() -> str | None:
    """定位 LibreOffice soffice 可执行文件。"""
    for p in _SOFFICE_CANDIDATES:
        if Path(p).exists():
            return p
    return shutil.which("soffice") or shutil.which("soffice.exe")


def has_visual_preview() -> bool:
    """是否具备真·视觉预览能力（soffice + fitz 都在）。"""
    if not find_soffice():
        return False
    try:
        import fitz  # noqa: F401
        return True
    except Exception:
        return False


def _preview_dir(doc_path: Path) -> Path:
    d = doc_path.parent / "_preview" / doc_path.stem
    d.mkdir(parents=True, exist_ok=True)
    return d


def render_to_images(
    doc_path: str | Path,
    *,
    dpi: int = 110,
    max_pages: int = 20,
    timeout: int = 150,
) -> list[Path]:
    """把 docx/pptx/pdf 渲染成 PNG 页图列表（带缓存）。失败返回 []。"""
    doc_path = Path(doc_path)
    if not doc_path.exists():
        return []

    out_dir = _preview_dir(doc_path)
    stamp = out_dir / ".stamp"
    mtime = str(int(doc_path.stat().st_mtime))
    existing = sorted(out_dir.glob("page_*.png"))
    if existing and stamp.exists() and stamp.read_text(encoding="utf-8").strip() == mtime:
        return existing[:max_pages]

    for p in existing:
        try:
            p.unlink()
        except Exception:
            pass

    # PDF 直接渲染；否则先用 soffice 转 PDF
    if doc_path.suffix.lower() == ".pdf":
        pdf_path: Path | None = doc_path
    else:
        soffice = find_soffice()
        if not soffice:
            return []
        # 用独立 user profile，避免与用户正在开着的 LibreOffice 冲突
        with tempfile.TemporaryDirectory() as prof:
            prof_uri = Path(prof).as_uri()
            try:
                subprocess.run(
                    [
                        soffice, "--headless", "--norestore", "--nologo",
                        f"-env:UserInstallation={prof_uri}",
                        "--convert-to", "pdf", "--outdir", str(out_dir),
                        str(doc_path),
                    ],
                    check=True, timeout=timeout,
                    stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                )
            except Exception:
                return []
        cand = out_dir / (doc_path.stem + ".pdf")
        if cand.exists():
            pdf_path = cand
        else:
            pdfs = list(out_dir.glob("*.pdf"))
            pdf_path = pdfs[0] if pdfs else None
    if not pdf_path or not pdf_path.exists():
        return []

    try:
        import fitz
    except Exception:
        return []

    imgs: list[Path] = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    try:
        doc = fitz.open(str(pdf_path))
    except Exception:
        return []
    try:
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pix = page.get_pixmap(matrix=mat)
            out = out_dir / f"page_{i:02d}.png"
            pix.save(str(out))
            imgs.append(out)
    finally:
        doc.close()

    try:
        stamp.write_text(mtime, encoding="utf-8")
    except Exception:
        pass
    return imgs


def _count_pages_office_com(doc_path: Path) -> int | None:
    """Windows 兜底：soffice 缺失时用 Office COM 实测页数（docx→Word 页数，pptx→幻灯片数）。"""
    suffix = doc_path.suffix.lower()
    try:
        import pythoncom  # noqa: F401
        import win32com.client as win32
    except Exception:
        return None
    if suffix in (".docx", ".doc"):
        word = None
        try:
            word = win32.DispatchEx("Word.Application")
            word.Visible = False
            d = word.Documents.Open(str(doc_path), ReadOnly=True)
            n = int(d.ComputeStatistics(2))  # 2 = wdStatisticPages
            d.Close(False)
            return n
        except Exception:
            return None
        finally:
            try:
                if word is not None:
                    word.Quit()
            except Exception:
                pass
    if suffix in (".pptx", ".ppt"):
        ppt = None
        try:
            ppt = win32.DispatchEx("PowerPoint.Application")
            pres = ppt.Presentations.Open(str(doc_path), WithWindow=False)
            n = int(pres.Slides.Count)
            pres.Close()
            return n
        except Exception:
            return None
        finally:
            try:
                if ppt is not None:
                    ppt.Quit()
            except Exception:
                pass
    return None


def count_pages(doc_path: str | Path, *, timeout: int = 150) -> int | None:
    """实测文档页数：用 soffice 转 PDF 后用 fitz 统计页数。

    返回页数；soffice/fitz 不可用或转换失败时返回 None（调用方应据此回退）。
    soffice 缺失时（如本机仅装了 Office）回退用 Word/PowerPoint COM 实测页数。
    用独立临时目录，不污染预览缓存。
    """
    doc_path = Path(doc_path)
    if not doc_path.exists():
        return None
    try:
        import fitz
    except Exception:
        fitz = None  # type: ignore

    if doc_path.suffix.lower() == ".pdf":
        if fitz is None:
            return None
        pdf_path: Path | None = doc_path
        tmp_ctx = None
    else:
        soffice = find_soffice()
        if not soffice or fitz is None:
            # soffice/fitz 不可用 → Office COM 兜底（Windows）
            return _count_pages_office_com(doc_path)
        tmp_ctx = tempfile.TemporaryDirectory()
        out_dir = Path(tmp_ctx.name)
        with tempfile.TemporaryDirectory() as prof:
            prof_uri = Path(prof).as_uri()
            try:
                subprocess.run(
                    [
                        soffice, "--headless", "--norestore", "--nologo",
                        f"-env:UserInstallation={prof_uri}",
                        "--convert-to", "pdf", "--outdir", str(out_dir),
                        str(doc_path),
                    ],
                    check=True, timeout=timeout,
                    stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                )
            except Exception:
                tmp_ctx.cleanup()
                return None
        cand = out_dir / (doc_path.stem + ".pdf")
        if cand.exists():
            pdf_path = cand
        else:
            pdfs = list(out_dir.glob("*.pdf"))
            pdf_path = pdfs[0] if pdfs else None
        if not pdf_path or not pdf_path.exists():
            tmp_ctx.cleanup()
            return None

    try:
        d = fitz.open(str(pdf_path))
        n = d.page_count
        d.close()
        return n
    except Exception:
        return None
    finally:
        if tmp_ctx is not None:
            tmp_ctx.cleanup()


def extract_text(doc_path: str | Path, *, max_chars: int = 6000) -> str:
    """纯文字兜底：从 docx/pptx 抽取文本（无 soffice 时用）。"""
    doc_path = Path(doc_path)
    suffix = doc_path.suffix.lower()
    chunks: list[str] = []
    try:
        if suffix == ".docx":
            from docx import Document
            d = Document(str(doc_path))
            for para in d.paragraphs:
                t = (para.text or "").strip()
                if t:
                    chunks.append(t)
            for tbl in d.tables:
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        chunks.append(line)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(doc_path))
            for i, slide in enumerate(prs.slides, 1):
                chunks.append(f"—— Slide {i} ——")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = (shape.text_frame.text or "").strip()
                        if t:
                            chunks.append(t)
    except Exception as e:
        return f"(文本预览失败：{e})"
    text = "\n".join(chunks)
    return text[:max_chars] + ("\n…(已截断)" if len(text) > max_chars else "")
