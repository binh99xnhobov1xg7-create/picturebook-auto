"""Worksheet PPTX 生成器 v2.0（8 页固定模板，对齐真实 L5-1 样本）。

页面结构（强制 8 页）：
    Page 1  Vocabulary  - Match the words to their definitions   (5 对连线)
    Page 2  Vocabulary  - Use the words / phrase to fill blanks  (5 题填空 + 词库条)
    Page 3  Sentence    - Choose the correct sentence            (4 题二选一 + 绘本图)
    Page 4  Sentence    - Sentence practice 2                    (句型巩固)
    Page 5  Reading     - Choose the correct answer              (阅读理解)
    Page 6  Reading     - Reading extension                      (阅读理解)
    Page 7  Graphic Organizer - Skill-based chart                 (思维导图/图表)
    Page 8  Writing     - Write about [theme]                    (写作区)

字体/字号（v1.6 真实样本）：
    大标题 Poppins Bold 20pt #333333  / 副标题 Poppins Regular 12pt #666666
    题号  Poppins Bold 16pt 圆形粉底白字
    题干  Poppins Regular 16pt 黑色  /  Reading 长文 12pt 黑色

品牌外框（统一 8 页）：
    粉色外背景 (BRAND_COLORS[level]) + 内白圆角
    左上 VIPKID Dino Reading Club logo（Dino 头像 + 白色文字）
    右上 Name 五角形角标 (粉底白字)
    右下 footer "Level X - <Title>" 白字
"""
from __future__ import annotations

import contextlib
import hashlib
import json
import random
import re
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches as _RawInches, Pt, Emu

from config import BRAND_DIR, brand_color_rgb
from parser import BookOutline
from text_format import (
    _to_us_spelling,
    capitalize_names,
    format_word_answer,
    format_sentence_answer,
    smart_format_answer,
    is_sentence_like,
)
from worksheet_activity_policy import (
    activity_required_fields as _ws_activity_required_fields,
    activity_min_items as _ws_activity_min_items,
    allowed_activity_codes as _ws_allowed_activity_codes,
    instruction as _ws_activity_instruction,
)


# ---------- 几何尺寸 ----------
# v2.2：直接以官方模板画布为准 = PowerPoint「A4 纸张」横向 = 27.517 x 19.05 cm
#        = 10.833 x 7.5 in。不再二次放大（_WS_SCALE=1.0），坐标即原生英寸。
DESIGN_W = 10.833
DESIGN_H = 7.5
_WS_SCALE = 1.0

SLIDE_W = DESIGN_W
SLIDE_H = DESIGN_H


def Inches(value):
    """worksheet 坐标 = 原生英寸（官方 A4 横向画布 10.833x7.5）。"""
    return _RawInches(value * _WS_SCALE)

# 内容白底（对齐官方模板的内白圆角区 x=0.46 y=0.73 w=9.9 h=6.43）
CONTENT_X = 0.46
CONTENT_Y = 0.73
CONTENT_W = 9.90
CONTENT_H = 6.43
CONTENT_ROUND = 0.06  # 圆角调整比例（python-pptx adjustments[0]）

# 顶部 logo 区（露在粉色背景上）
LOGO_X = 0.40
LOGO_Y = 0.10
LOGO_ICON_W = 0.55
LOGO_ICON_H = 0.55

# Name 角标（v1.9：放成矩形标签，避免旋转 PENTAGON 跑位）
NAME_X = 8.85
NAME_Y = 0.18
NAME_W = 1.50
NAME_H = 0.42

# Footer（v1.9：往上挪到内容白底底边附近，避免被裁切）
FOOTER_X = 0.30
FOOTER_Y = 6.85
FOOTER_W = 10.23
FOOTER_H = 0.30


# ---------- 字号/颜色 ----------
FONT = "Poppins"
FONT_BOLD = "Poppins"
# underscore 字符在 Poppins 下被压扁，改用 Arial 才显示得清晰粗实
FONT_BLANK = "Arial"
ANSWER_BLANK = "_" * 20

# 验收口径：大标题 40pt 黑；副标题/command 控制为 14pt，避免指令抢占作答区
TITLE_PT = 40
SUBTITLE_PT = 14
BODY_PT = 18
READING_PT = 14
QNUM_PT = 18
HEADER_PT = 20  # mind map 表头
LOGO_TEXT_PT = 20
NAME_PT = 16
FOOTER_PT = 14.5

def _luminance(rgb: tuple) -> float:
    """相对亮度 0(黑)~1(白)。用于按底色亮度自动选黑/白前景字（Bug5：浅底白字看不清）。"""
    try:
        r, g, b = (int(rgb[0]) / 255.0, int(rgb[1]) / 255.0, int(rgb[2]) / 255.0)
    except Exception:
        return 0.0
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _readable_text_rgb(bg_rgb: tuple) -> "RGBColor":
    """根据底色亮度返回可读前景色：浅底→深炭灰，深底→白。阈值 0.6。"""
    return RGBColor(0x33, 0x33, 0x33) if _luminance(bg_rgb) >= 0.6 else RGBColor(0xFF, 0xFF, 0xFF)


TITLE_RGB = RGBColor(0x33, 0x33, 0x33)  # 深炭灰
SUB_RGB = RGBColor(0x66, 0x66, 0x66)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAME_FILL = RGBColor(0xF8, 0xC8, 0xDC)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# Mind Map 三列表头颜色（v1.6 紫粉 / 黄 / 绿）
MM_PURPLE = RGBColor(0xE6, 0xD8, 0xF2)
MM_YELLOW = RGBColor(0xFA, 0xEB, 0xC6)
MM_GREEN = RGBColor(0xCE, 0xE7, 0xCD)
MM_PURPLE_DARK = RGBColor(0xC9, 0xB3, 0xE2)
MM_YELLOW_DARK = RGBColor(0xF5, 0xDC, 0x95)
MM_GREEN_DARK = RGBColor(0xA8, 0xD3, 0xA6)

# Reading 红框
READ_BORDER = RGBColor(0xE9, 0x52, 0x83)
# P1 Activity Bank renderers.

def _visual_vocab_items(pairs: list[dict], images: list[Path], max_n: int = 4) -> list[dict]:
    items: list[dict] = []
    story_images = [p for p in (images or [])[2:] if p and Path(p).exists()]
    for i, pair in enumerate(pairs or []):
        word = str(pair.get("word", "")).strip()
        if not word:
            continue
        kind = _clean_vocab_clue_kind(word)
        img = story_images[i % len(story_images)] if story_images else None
        if not kind and img is None:
            continue
        row = dict(pair)
        row["visual_kind"] = kind
        row["img"] = img
        items.append(row)
        if len(items) >= max_n:
            break
    return items


def _meaning_mc_items(pairs: list[dict], max_n: int = 4) -> list[dict]:
    clean_pairs = [p for p in (pairs or []) if _valid_definition_pair(p)]
    defs = [str(p.get("def", "")).strip() for p in clean_pairs if str(p.get("def", "")).strip()]
    out: list[dict] = []
    for pair in clean_pairs[:max_n]:
        word = str(pair.get("word", "")).strip()
        correct = str(pair.get("def", "")).strip()
        if not word or not correct:
            continue
        wrong = [d for d in defs if d and d != correct]
        opts = [correct] + wrong[:2]
        rnd = random.Random(hash((word.lower(), "meaning")) & 0xFFFFFFFF)
        rnd.shuffle(opts)
        out.append({
            "kind": "mc",
            "q": f"What does {word} mean?",
            "options": opts,
            "correct": opts.index(correct),
            "answer": correct,
        })
    return out


def _draw_vocab_visual_box(slide, brand_rgb: tuple, word: str, img: Optional[Path],
                           x: float, y: float, w: float, h: float) -> None:
    placed = _draw_clean_vocab_clue(slide, _clean_vocab_clue_kind(word), x, y, w, h, brand_rgb)
    if placed:
        return
    if img and Path(img).exists():
        try:
            from PIL import Image as _PILImg
            with _PILImg.open(str(img)) as _pim:
                iw, ih = _pim.size
            aspect = iw / ih if ih else 1.35
            fit_h = h
            fit_w = fit_h * aspect
            if fit_w > w:
                fit_w = w
                fit_h = fit_w / aspect
            off_x = x + (w - fit_w) / 2
            off_y = y + (h - fit_h) / 2
            slide.shapes.add_picture(str(img), Inches(off_x), Inches(off_y),
                                     width=Inches(fit_w), height=Inches(fit_h))
            return
        except Exception:
            pass
    _draw_image_placeholder(slide, x, y, w, h)


def _build_p1_word_picture_match(slide, brand_rgb: tuple, items: list[dict]) -> None:
    _add_title(slide, "Vocabulary", _ws_activity_instruction("vocab_word_picture_matching"))
    n = min(len(items), 4)
    if n == 0:
        return
    order = _derange_order(n)
    area_top = CONTENT_Y + 1.38
    area_h = CONTENT_H - 1.78
    gap = 0.22
    row_h = (area_h - gap * (n - 1)) / n
    word_x = CONTENT_X + 0.70
    word_w = 3.00
    pic_x = CONTENT_X + 5.05
    pic_w = CONTENT_W - 5.75

    for i, item in enumerate(items[:n]):
        y = area_top + i * (row_h + gap)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(word_x), Inches(y + 0.10),
                                       Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*brand_rgb)
        badge.line.fill.background()
        badge.shadow.inherit = False
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = str(i + 1)
        br.font.name = FONT
        br.font.size = Pt(12)
        br.font.bold = True
        br.font.color.rgb = WHITE

        wc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(word_x + 0.55), Inches(y),
                                    Inches(word_w), Inches(row_h))
        wc.adjustments[0] = 0.25
        wc.fill.solid()
        wc.fill.fore_color.rgb = RGBColor(*brand_rgb)
        wc.line.fill.background()
        wc.shadow.inherit = False
        tf = wc.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = _clean_text(item.get("word", ""))
        r.font.name = FONT
        r.font.size = Pt(18 if len(r.text) <= 14 else 14)
        r.font.color.rgb = WHITE

        visual_item = items[order[i]] if order[i] < len(items) else item
        lb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pic_x), Inches(y + 0.10),
                                    Inches(0.42), Inches(0.42))
        lb.fill.solid()
        lb.fill.fore_color.rgb = RGBColor(*brand_rgb)
        lb.line.fill.background()
        lb.shadow.inherit = False
        lp = lb.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = chr(65 + i)
        lr.font.name = FONT
        lr.font.size = Pt(12)
        lr.font.bold = True
        lr.font.color.rgb = WHITE
        _draw_vocab_visual_box(slide, brand_rgb, str(visual_item.get("word", "")),
                               visual_item.get("img"), pic_x + 0.55, y, pic_w, row_h)


def _build_p1_choose_picture(slide, brand_rgb: tuple, items: list[dict]) -> None:
    _add_title(slide, "Vocabulary", _ws_activity_instruction("vocab_choose_picture"))
    choice_rows = items if items and items[0].get("_option_items") else _picture_choice_items(items)
    n = min(len(choice_rows), 4)
    if n == 0:
        return
    top = CONTENT_Y + 1.32
    gap_x, gap_y = 0.52, 0.28
    cell_w = (CONTENT_W - 0.90 - gap_x) / 2
    cell_h = (CONTENT_H - 1.72 - gap_y) / 2
    x0 = CONTENT_X + 0.45
    for i, item in enumerate(choice_rows[:n]):
        rr, cc = divmod(i, 2)
        x = x0 + cc * (cell_w + gap_x)
        y = top + rr * (cell_h + gap_y)
        word = str(item.get("word", "")).strip()
        opts = item.get("_option_items") or []

        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cell_w), Inches(0.38))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{i + 1}. {word}   (      )"
        r.font.name = FONT
        r.font.size = Pt(17)
        r.font.color.rgb = BLACK

        opt_top = y + 0.48
        opt_gap = 0.12
        opt_w = (cell_w - opt_gap * 2) / 3
        opt_h = cell_h - 0.62
        for j, opt in enumerate(opts[:3]):
            ox = x + j * (opt_w + opt_gap)
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ox), Inches(opt_top),
                                           Inches(0.32), Inches(0.32))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(*brand_rgb)
            badge.line.fill.background()
            badge.shadow.inherit = False
            bp = badge.text_frame.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = chr(65 + j)
            br.font.name = FONT
            br.font.size = Pt(9.5)
            br.font.bold = True
            br.font.color.rgb = WHITE
            _draw_vocab_visual_box(slide, brand_rgb, str(opt.get("word", "")),
                                   opt.get("img"), ox, opt_top + 0.38, opt_w, opt_h - 0.38)


def _picture_choice_items(items: list[dict]) -> list[dict]:
    all_items = list(items or [])[:4]
    out: list[dict] = []
    for item in all_items:
        word = str(item.get("word", "")).strip()
        if not word:
            continue
        wrong = [it for it in all_items if str(it.get("word", "")).lower() != word.lower()]
        opts = [item] + wrong[:2]
        rnd = random.Random(hash((word.lower(), "picture_options")) & 0xFFFFFFFF)
        rnd.shuffle(opts)
        row = dict(item)
        row["_option_items"] = opts
        row["options"] = [str(o.get("word", "")).strip() for o in opts]
        row["correct"] = next((idx for idx, opt in enumerate(opts)
                               if str(opt.get("word", "")).lower() == word.lower()), 0)
        row["answer"] = word
        out.append(row)
    return out


def _build_l34_vocab1(new_page, brand_rgb: tuple, data: dict, images: list[Path],
                      seed: int, outline: BookOutline) -> None:
    pairs = list(data.get("match_pairs") or [])
    valid_defs = [p for p in pairs if _valid_definition_pair(p)]
    visual_items = _visual_vocab_items(pairs, images, max_n=4)

    for code in _l34_activity_order(outline, 1, seed=seed, include_optional=False):
        if code == "vocab_word_picture_matching" and len(visual_items) >= _l34_min_items(code, 4):
            order = _derange_order(min(len(visual_items), 4))
            answer_letters = {}
            for row_idx, item_idx in enumerate(order):
                if item_idx < len(visual_items):
                    answer_letters[str(visual_items[item_idx].get("word", ""))] = chr(65 + row_idx)
            _record_l34_activity(outline, 1, code)
            _record_worksheet_items(outline, 1, [
                {
                    "word": it.get("word", ""),
                    "visual_kind": it.get("visual_kind", ""),
                    "answer_picture": answer_letters.get(str(it.get("word", "")), ""),
                }
                for it in visual_items[:4]
            ])
            _build_p1_word_picture_match(new_page(), brand_rgb, visual_items)
            return
        if code == "vocab_choose_picture" and len(visual_items) >= _l34_min_items(code, 4):
            choice_items = _picture_choice_items(visual_items)
            _record_l34_activity(outline, 1, code)
            _record_worksheet_items(outline, 1, [
                {
                    "word": it.get("word", ""),
                    "options": it.get("options", []),
                    "correct": it.get("correct", 0),
                    "answer": it.get("answer", ""),
                    "visual_kind": it.get("visual_kind", ""),
                }
                for it in choice_items
            ])
            _build_p1_choose_picture(new_page(), brand_rgb, choice_items)
            return
        if code == "vocab_choose_meaning" and len(valid_defs) >= _l34_min_items(code, 4):
            mc_items = _meaning_mc_items(pairs, max_n=4)
            if len(mc_items) >= _l34_min_items(code, 4):
                _record_l34_activity(outline, 1, code)
                _record_worksheet_items(outline, 1, mc_items)
                _build_mcq_page(new_page(), brand_rgb, "Vocabulary",
                                _ws_activity_instruction(code), mc_items)
                return
        if code == "vocab_word_definition_matching" and len(valid_defs) >= _l34_min_items(code, 4):
            _record_l34_activity(outline, 1, code)
            _record_worksheet_items(outline, 1, pairs)
            _build_p1_match(new_page(), brand_rgb, pairs, images)
            return

    if len(visual_items) >= 4:
        order = _derange_order(min(len(visual_items), 4))
        answer_letters = {}
        for row_idx, item_idx in enumerate(order):
            if item_idx < len(visual_items):
                answer_letters[str(visual_items[item_idx].get("word", ""))] = chr(65 + row_idx)
        _record_l34_activity(outline, 1, "vocab_word_picture_matching")
        _record_worksheet_items(outline, 1, [
            {
                "word": it.get("word", ""),
                "visual_kind": it.get("visual_kind", ""),
                "answer_picture": answer_letters.get(str(it.get("word", "")), ""),
            }
            for it in visual_items[:4]
        ])
        _build_p1_word_picture_match(new_page(), brand_rgb, visual_items)
    else:
        _record_l34_activity(outline, 1, "vocab_word_definition_matching")
        _record_worksheet_items(outline, 1, pairs)
        _build_p1_match(new_page(), brand_rgb, pairs, images)


# ============================================================
#  对外入口
# ============================================================

def _level_num(level: str) -> int:
    """取 level 数字（smart→0），用于分级页面逻辑。"""
    s = str(level or "").lower()
    if "smart" in s:
        return 0
    digits = "".join(ch for ch in s if ch.isdigit())
    try:
        return int(digits) if digits else 5
    except ValueError:
        return 5


def _record_l34_activity(outline: Optional[BookOutline], page: int, code: str) -> None:
    """Record the selected L3/L4 worksheet activity for QA/TG-only workflows."""
    if outline is None:
        return
    try:
        activity_map = getattr(outline, "_worksheet_activity_types", None)
        if not isinstance(activity_map, dict):
            activity_map = {}
            setattr(outline, "_worksheet_activity_types", activity_map)
        activity_map[page] = code
    except Exception:
        return


def _record_worksheet_items(outline: Optional[BookOutline], page: int, items: list[dict]) -> None:
    """Record structured page items for the worksheet manifest and later TG sync."""
    if outline is None:
        return
    try:
        page_items = getattr(outline, "_worksheet_manifest_items", None)
        if not isinstance(page_items, dict):
            page_items = {}
            setattr(outline, "_worksheet_manifest_items", page_items)
        page_items[page] = list(items or [])
    except Exception:
        return


def _l34_activity_order(outline: BookOutline, page: int, *, seed: int,
                        include_optional: bool = False) -> list[str]:
    lvl = _level_num(getattr(outline, "level", "") or "")
    text_type = getattr(outline, "fiction_type", "") or ""
    codes = _ws_allowed_activity_codes(
        page,
        level=lvl,
        text_type=text_type,
        include_optional=include_optional,
    )
    if not codes:
        return []
    offset = seed % len(codes)
    return codes[offset:] + codes[:offset]


def _l34_min_items(code: str, default: int = 4) -> int:
    return _ws_activity_min_items(code, default)


def _l34_go_activity_code(outline: BookOutline, go_mode: str) -> str:
    """Map the rendered organizer mode to the worksheet activity bank."""
    is_nonfic = "non" in (getattr(outline, "fiction_type", "") or "").lower()
    explicit_go = " ".join(
        str(getattr(outline, name, "") or "")
        for name in (
            "graphic_organizer",
            "graphic_organizer_desc",
            "reading_skill",
            "reading_strategy",
        )
    ).strip().lower()
    if explicit_go:
        if any(k in explicit_go for k in (
            "compare", "contrast", "venn", "same", "different",
        )):
            return "go_compare_chart"
        if any(k in explicit_go for k in (
            "timeline", "sequence", "sequencing", "plan chart", "journey",
            "steps", "process", "order", "procedure",
        )):
            return "go_sequence_chart"
        if any(k in explicit_go for k in (
            "problem", "solution", "action", "result", "story map",
            "story element", "swbst", "somebody", "wanted", "beginning",
            "middle", "end",
        )):
            return "go_problem_plan_result"
        if any(k in explicit_go for k in (
            "classification", "classify", "category", "categories",
            "fact", "fact web", "bubble", "main idea", "details",
            "topic", "habitat", "labeled diagram",
        )):
            return "go_fact_web"
    mode = (go_mode or "").strip().lower()
    if mode in {"l3bubble", "bubble", "factweb", "fact_web"}:
        return "go_fact_web"
    if is_nonfic:
        return "go_fact_web"
    if mode in {"timeline", "sequence", "planchart"}:
        return "go_sequence_chart"
    return "go_problem_plan_result"


def _worksheet_manifest_path(out_path: Path) -> Path:
    return out_path.with_name(out_path.stem + "_manifest.json")


def _worksheet_lesson_id(outline: BookOutline) -> str:
    lvl = _level_num(getattr(outline, "level", "") or "")
    book_raw = str(getattr(outline, "book_number", "") or "").strip()
    digits = "".join(ch for ch in book_raw if ch.isdigit())
    return f"L{lvl}-B{digits or book_raw or 'unknown'}"


def _worksheet_story_sentences(text: str) -> list[dict]:
    rows: list[dict] = []
    for idx, sent in enumerate(re.split(r"(?<=[.!?])\s+", _clean_text(text or "")), 1):
        sent = sent.strip()
        if sent:
            rows.append({"id": f"S{idx}", "text": sent})
    return rows


def _worksheet_content_hash(text: str) -> str:
    normalized = re.sub(r"\s+", " ", _clean_text(text or "")).strip().lower()
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:16] if normalized else ""


def _worksheet_image_checks(images: list[Path]) -> list[dict]:
    checks: list[dict] = []
    for idx, raw in enumerate(images or []):
        path = Path(raw)
        row = {
            "index": idx,
            "path": str(path),
            "exists": path.exists(),
            "bytes": path.stat().st_size if path.exists() else 0,
            "width": None,
            "height": None,
            "blank_like": None,
        }
        if path.exists():
            try:
                from PIL import Image, ImageStat  # type: ignore
                with Image.open(str(path)) as im:
                    row["width"], row["height"] = im.size
                    stat = ImageStat.Stat(im.convert("RGB").resize((32, 32)))
                    row["blank_like"] = max(stat.stddev or [0]) < 2.0
            except Exception:
                row["blank_like"] = None
        checks.append(row)
    return checks


def _worksheet_go_type(outline: BookOutline, go_mode: str = "") -> str:
    explicit = (
        getattr(outline, "graphic_organizer", "")
        or getattr(outline, "graphic_organizer_desc", "")
        or ""
    ).strip()
    if explicit:
        return explicit
    code = _l34_go_activity_code(outline, go_mode)
    return {
        "go_sequence_chart": "Sequence Chart",
        "go_problem_plan_result": "Problem-Plan-Result Chart",
        "go_fact_web": "Fact Web",
        "go_compare_chart": "Comparison Chart",
    }.get(code, code)


def _worksheet_source_preflight(outline: BookOutline, data: dict, lvl_n: int) -> dict:
    blockers: list[dict] = []
    warnings: list[dict] = []
    title = (getattr(outline, "title", "") or "").strip()
    reading_text = _clean_text(data.get("reading_text", "") or getattr(outline, "story_text", ""))
    pairs = list(data.get("match_pairs") or [])
    core_vocab = [str(p.get("word", "")).strip() for p in pairs if str(p.get("word", "")).strip()]
    valid_defs = [p for p in pairs if _valid_definition_pair(p)]

    def _issue(bucket: list[dict], code: str, message: str, severity: str) -> None:
        bucket.append({"code": code, "severity": severity, "message": message})

    if not title:
        _issue(blockers, "missing_title", "Worksheet title is empty.", "BLOCKER")
    if not reading_text:
        _issue(blockers, "missing_passage", "Book passage text is empty.", "BLOCKER")
    if lvl_n in (3, 4):
        if len(core_vocab) < 4:
            _issue(blockers, "core_vocab_too_few", "L3/L4 Worksheet needs at least 4 core vocabulary items.", "BLOCKER")
        if len(valid_defs) < min(4, len(core_vocab)):
            _issue(warnings, "p1_missing_definitions", "P1 vocabulary definitions are missing or placeholder-like.", "WARNING")
    low_text = f" {reading_text.lower()} "
    missing_vocab = [w for w in core_vocab if w.lower() not in low_text]
    if missing_vocab:
        _issue(
            warnings,
            "core_vocab_not_in_passage",
            "Core vocabulary not found verbatim in passage: " + ", ".join(missing_vocab),
            "WARNING",
        )
    return {"blockers": blockers, "warnings": warnings}


def _relax_nonblocking_worksheet_issues(source_validation: dict) -> dict:
    """Keep review-only issues from blocking Worksheet generation."""
    if not isinstance(source_validation, dict):
        return {"blockers": [], "warnings": []}
    blockers = list(source_validation.get("blockers") or [])
    warnings = list(source_validation.get("warnings") or [])
    kept_blockers: list[dict] = []
    for issue in blockers:
        code = str(issue.get("code", "") if isinstance(issue, dict) else "")
        message = str(issue.get("message", "") if isinstance(issue, dict) else issue)
        if code == "p1_missing_definitions" or "P1 vocabulary definitions are missing" in message:
            if isinstance(issue, dict):
                relaxed = dict(issue)
                relaxed["severity"] = "WARNING"
                warnings.append(relaxed)
            else:
                warnings.append({
                    "code": "p1_missing_definitions",
                    "severity": "WARNING",
                    "message": message,
                })
            continue
        kept_blockers.append(issue)
    return {"blockers": kept_blockers, "warnings": warnings}


def _worksheet_activity_validation(
    activity_map: dict,
    manifest_items: dict,
) -> dict:
    blockers: list[dict] = []
    warnings: list[dict] = []

    def _issue(bucket: list[dict], code: str, message: str, severity: str) -> None:
        bucket.append({"code": code, "severity": severity, "message": message})

    for page in range(1, 9):
        code = str(activity_map.get(page, "") or "").strip()
        if not code:
            _issue(warnings, f"p{page}_missing_activity", f"Page {page} has no recorded activity code.", "WARNING")
            continue
        required = _ws_activity_required_fields(code)
        if page <= 6:
            items = manifest_items.get(page, [])
            if not isinstance(items, list) or not items:
                _issue(warnings, f"p{page}_missing_items", f"Page {page} has no structured worksheet items.", "WARNING")
                continue
            missing_counts = 0
            for item in items:
                if not isinstance(item, dict):
                    missing_counts += 1
                    continue
                for field in required:
                    if field in {"answer_map", "word_bank"}:
                        continue
                    if field == "answer" and (
                        item.get("answer") or item.get("correct") is not None or item.get("answer_order") is not None
                    ):
                        continue
                    if field == "definition" and item.get("def"):
                        continue
                    if field == "clue" and item.get("q"):
                        continue
                    if field == "sentence" and (item.get("sentence") or item.get("q") or item.get("prompt")):
                        continue
                    if field == "sentence_frame" and (
                        item.get("sentence") or item.get("prompt") or item.get("q") or item.get("options")
                    ):
                        continue
                    if field == "question" and item.get("q"):
                        continue
                    if field == "statement" and item.get("q"):
                        continue
                    if field == "summary_sentence" and item.get("q"):
                        continue
                    if field == "events" and item.get("event"):
                        continue
                    if field == "wrong_sentence" and item.get("prompt", "").lower().startswith("correct"):
                        continue
                    if field == "scrambled_words" and item.get("scrambled"):
                        continue
                    if not item.get(field):
                        missing_counts += 1
                        break
            if missing_counts:
                _issue(
                    warnings,
                    f"p{page}_missing_required_fields",
                    f"Page {page} activity {code} has {missing_counts} item(s) missing required fields: {', '.join(required)}.",
                    "WARNING",
                )
    return {"blockers": blockers, "warnings": warnings}


def _write_worksheet_manifest(
    outline: BookOutline,
    out_path: Path,
    data: dict,
    images: list[Path],
    go_mode: str,
    source_validation: dict,
) -> Path:
    reading_text = _clean_text(data.get("reading_text", "") or getattr(outline, "story_text", ""))
    activity_map = getattr(outline, "_worksheet_activity_types", {}) or {}
    manifest_items = getattr(outline, "_worksheet_manifest_items", {}) or {}
    pairs = [
        {"word": str(p.get("word", "")).strip(), "definition": str(p.get("def", "")).strip()}
        for p in (data.get("match_pairs") or [])
        if str(p.get("word", "")).strip()
    ]
    activity_validation = _worksheet_activity_validation(activity_map, manifest_items)
    validation_results = {
        "source_preflight": source_validation,
        "activity_validation": activity_validation,
        "blockers": list(source_validation.get("blockers", [])) + list(activity_validation.get("blockers", [])),
        "warnings": list(source_validation.get("warnings", [])) + list(activity_validation.get("warnings", [])),
    }
    manifest = {
        "schema_version": "worksheet_manifest_v1",
        "lesson_spec": {
            "lesson_id": _worksheet_lesson_id(outline),
            "level": _level_num(getattr(outline, "level", "") or ""),
            "book_number": str(getattr(outline, "book_number", "") or ""),
            "title": str(getattr(outline, "title", "") or ""),
            "source_version": "S&S",
            "content_hash": _worksheet_content_hash(reading_text),
            "core_vocabulary": [p["word"] for p in pairs],
            "target_pattern": _sentence_frame_text(outline),
            "sample_sentence": _display_sentence_frame(outline),
            "reading_skill": str(getattr(outline, "reading_skill", "") or ""),
            "reading_strategy": str(getattr(outline, "reading_strategy", "") or ""),
            "graphic_organizer_type": _worksheet_go_type(outline, go_mode),
            "asset_lesson_id": _worksheet_lesson_id(outline),
        },
        "passage_sentences": _worksheet_story_sentences(reading_text),
        "page_plan": [
            {
                "page": page,
                "activity_code": activity_map.get(page, ""),
                "required_fields": list(_ws_activity_required_fields(activity_map.get(page, ""))),
            }
            for page in range(1, 9)
        ],
        "items": {
            "page_1": manifest_items.get(1, pairs),
            "page_2": manifest_items.get(2, []),
            "page_3": manifest_items.get(3, []),
            "page_4": manifest_items.get(4, []),
            "page_5": manifest_items.get(5, []),
            "page_6": manifest_items.get(6, []),
            "page_7": manifest_items.get(7, {
                "graphic_organizer_type": _worksheet_go_type(outline, go_mode),
                "activity_code": activity_map.get(7, ""),
            }),
            "page_8": manifest_items.get(8, {
                "activity_code": activity_map.get(8, ""),
                "must_use_go": True,
                "target_pattern": _sentence_frame_text(outline),
            }),
        },
        "image_checks": _worksheet_image_checks(images),
        "validation_results": validation_results,
        "release_status": "blocked" if validation_results.get("blockers") else "review_required",
    }
    path = _worksheet_manifest_path(out_path)
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


# ============================================================
#  句型考点引擎：故事时态（一般过去时）—— 学什么考什么
# ============================================================
# 过去式 → 原形（不规则 + 常见 magic-e / -ied 规则动词，避免启发式还原出错）
_PAST_TO_BASE: dict[str, str] = {
    # be / 高频不规则
    "was": "be", "were": "be", "felt": "feel", "saw": "see", "shook": "shake",
    "ran": "run", "said": "say", "took": "take", "gave": "give", "told": "tell",
    "went": "go", "had": "have", "made": "make", "came": "come", "brought": "bring",
    "thought": "think", "found": "find", "got": "get", "knew": "know", "drew": "draw",
    "ate": "eat", "sat": "sit", "stood": "stand", "heard": "hear", "held": "hold",
    "kept": "keep", "left": "leave", "met": "meet", "won": "win", "wrote": "write",
    "began": "begin", "drank": "drink", "swam": "swim", "sang": "sing", "flew": "fly",
    "grew": "grow", "threw": "throw", "blew": "blow", "broke": "break", "spoke": "speak",
    "woke": "wake", "chose": "choose", "rode": "ride", "drove": "drive", "rose": "rise",
    "fell": "fall", "bought": "buy", "caught": "catch", "taught": "teach", "fought": "fight",
    "slept": "sleep", "spent": "spend", "sent": "send", "built": "build", "paid": "pay",
    "burst": "burst", "read": "read", "put": "put", "cut": "cut", "hit": "hit",
    "let": "let", "set": "set", "shut": "shut", "hurt": "hurt", "cost": "cost",
    # magic-e 规则动词（启发式会还原错，单列）
    "shared": "share", "liked": "like", "baked": "bake", "smiled": "smile",
    "hoped": "hope", "moved": "move", "lived": "live", "used": "use", "closed": "close",
    "danced": "dance", "loved": "love", "saved": "save", "named": "name", "placed": "place",
    "arrived": "arrive", "decided": "decide", "invited": "invite", "noticed": "notice",
    "promised": "promise", "surprised": "surprise", "cared": "care", "waved": "wave",
    "smiled ": "smile", "raced": "race",
    # -ied
    "hurried": "hurry", "carried": "carry", "tried": "try", "cried": "cry",
    "studied": "study", "replied": "reply", "worried": "worry",
}


def _regular_base(past: str) -> Optional[str]:
    """规则动词过去式 → 原形（启发式，用于 _PAST_TO_BASE 未覆盖的 -ed 词）。"""
    if not past.endswith("ed") or len(past) < 4:
        return None
    if past.endswith("ied"):
        return past[:-3] + "y"          # tried→try
    stem = past[:-2]
    # 双写辅音：grabbed→grab, stopped→stop, hugged→hug
    if len(stem) >= 2 and stem[-1] == stem[-2] and stem[-1] not in "aeiou":
        return stem[:-1]
    return stem                          # listened→listen, helped→help, looked→look


def _past_to_base(word: str) -> Optional[str]:
    low = word.lower()
    if low in _PAST_TO_BASE:
        return _PAST_TO_BASE[low]
    return _regular_base(low)


def _present_3rd(base: str) -> str:
    """原形 → 一般现在时第三人称单数（用于造『现在时』错误项）。"""
    if base == "be":
        return "is"
    if base == "have":
        return "has"
    if base.endswith(("s", "x", "z", "ch", "sh", "o")):
        return base + "es"
    if base.endswith("y") and len(base) >= 2 and base[-2] not in "aeiou":
        return base[:-1] + "ies"
    return base + "s"


# 不是动词的常见 -ed/过去式同形词（避免误判）
_NOT_VERB = {"red", "bed", "bird", "good", "food", "wood", "need", "feed", "seed", "ahead"}


def _find_past_verb(sentence: str) -> Optional[tuple[str, str]]:
    """在句子里找第一个『过去式动词』token，返回 (原 token, 原形 base)。"""
    for tok in sentence.split():
        clean = "".join(ch for ch in tok if ch.isalpha())
        if not clean or clean.lower() in _NOT_VERB:
            continue
        base = _past_to_base(clean)
        if base and base != clean.lower():  # 必须确实是过去式（与原形不同）
            return tok, base
        # was/were/read/put 等过去=原形同形的，靠白名单确认
        if clean.lower() in ("was", "were"):
            return tok, "be"
    return None


_PAGE_PREFIX_RE = re.compile(r"^\s*(?:page|p|第)\s*\d+\s*[页:：.、)\-]?\s*", re.IGNORECASE)


def _story_sentences_for_grammar(outline: BookOutline) -> list[str]:
    out: list[str] = []
    for p in outline.pages:
        if p.page_type == "story" and (p.text or "").strip():
            # 去掉可能残留的分页标记前缀（"Page 1" / "P2" / "第3页"），避免混进句型题。
            txt = _PAGE_PREFIX_RE.sub("", p.text.strip())
            for s in re.split(r"(?<=[.!?])\s+", txt.strip()):
                s = _PAGE_PREFIX_RE.sub("", s.strip()).strip()
                if s and len(s.split()) >= 4:
                    out.append(capitalize_names(s))
    return out


def _sentence_pattern_fallback_items(
    outline: BookOutline, word_bank: Optional[list] = None, max_n: int = 4
) -> list[dict]:
    """句型页兜底题（Bug1：句型信息缺失导致整页空白）。

    优先用官方 S&S 的 sentence_pattern / example_sentence；缺失则用 grammar_focus + 词表，
    最后兜底"就故事写一句"。保证句型页永不空白。返回 {"prompt": ...} 列表（配整宽作答横线）。
    """
    entry = getattr(outline, "syllabus", None)
    pattern = (getattr(entry, "sentence_pattern", "") or "").strip() if entry else ""
    example = (getattr(entry, "example_sentence", "") or "").strip() if entry else ""
    items: list[dict] = []
    if example:
        items.append({"prompt": f"Read and copy the sentence:  {example}"})
    if pattern and pattern.lower() != example.lower():
        items.append({"prompt": f"Use this pattern to write a sentence:  {pattern}"})
    # 用核心词各造一句
    words: list[str] = []
    for w in (word_bank or []):
        word = (w.get("word") if isinstance(w, dict) else str(w or "")).strip()
        if word and word.isascii():
            words.append(word)
    for word in words:
        if len(items) >= max_n:
            break
        items.append({"prompt": f"Write a sentence using the word:  {word}"})
    if not items:
        focus = (outline.grammar_focus or "").strip()
        base = f" ({focus})" if focus else ""
        items.append({"prompt": f"Write one sentence about the story{base}."})
        items.append({"prompt": "Then write one more sentence of your own."})
    return items[:max_n]


def _sentence_frame_text(outline: BookOutline) -> str:
    entry = getattr(outline, "syllabus", None)
    frame = (getattr(entry, "sentence_pattern", "") or "").strip() if entry else ""
    if not frame:
        frame = (getattr(entry, "example_sentence", "") or "").strip() if entry else ""
    if not frame:
        frame = (outline.grammar_focus or "").strip()
    if re.search(r"[\u3400-\u9fff]", frame) or re.fullmatch(r"\?{3,}", frame):
        return "She will do homework first."
    if "[" in frame and "]" in frame:
        return "She will do homework first."
    return frame or "She will do homework first."


def _syllabus_sentence_source(outline: BookOutline) -> tuple[str, str]:
    """Return the official sentence pattern and example when the syllabus has them."""
    entry = getattr(outline, "syllabus", None)
    pattern = (getattr(entry, "sentence_pattern", "") or "").strip() if entry else ""
    example = (getattr(entry, "example_sentence", "") or "").strip() if entry else ""
    if not pattern:
        pattern = (getattr(entry, "sentence_frames", "") or "").strip() if entry else ""
    if not pattern:
        pattern = (getattr(entry, "syntax_focus", "") or "").strip() if entry else ""
    if not pattern:
        pattern = (outline.grammar_focus or "").strip()
    pattern = _clean_text(pattern)
    example = _clean_text(example)
    if re.search(r"[\u3400-\u9fff]", pattern) or re.fullmatch(r"\?{3,}", pattern or ""):
        pattern = ""
    if re.search(r"[\u3400-\u9fff]", example) or re.fullmatch(r"\?{3,}", example or ""):
        example = ""
    return pattern, example


def _display_sentence_frame(outline: BookOutline) -> str:
    """Student-facing frame text. Prefer a real example over bracket notation."""
    pattern, example = _syllabus_sentence_source(outline)
    if example:
        return example
    if pattern:
        return re.sub(r"\[[^\]]+\]", ANSWER_BLANK, pattern)
    return _sentence_frame_text(outline)


def _frame_signature(text: str) -> str:
    low = (text or "").lower()
    if " will " in f" {low} ":
        return "will"
    if re.search(r"\bthere\s+(is|are|was|were)\b", low):
        return "there"
    if re.search(r"\b(could|can)\b", low):
        return "modal"
    if " because " in f" {low} ":
        return "because"
    if " when " in f" {low} ":
        return "when"
    if " to " in f" {low} ":
        return "to"
    return ""


def _blank_for_frame(sentence: str, signature: str) -> tuple[str, str] | None:
    """Blank one teachable chunk while preserving the official sentence frame."""
    s = _clean_text(sentence).strip()
    if not s:
        return None
    if signature == "will":
        m = re.search(r"\bwill\s+(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(1) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(1)] + ANSWER_BLANK + (m.group(2) or "."), ans
    if signature == "there":
        m = re.search(r"\b(there\s+(?:is|are|was|were)\s+)(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(2) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(2)] + ANSWER_BLANK + (m.group(3) or "."), ans
    if signature == "modal":
        m = re.search(r"\b((?:can|could|could not|cannot|can't)\s+)(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(2) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(2)] + ANSWER_BLANK + (m.group(3) or "."), ans
    if signature == "because":
        m = re.search(r"\bbecause\s+(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(1) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(1)] + ANSWER_BLANK + (m.group(2) or "."), ans
    if signature == "when":
        m = re.search(r"\bwhen\s+(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(1) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(1)] + ANSWER_BLANK + (m.group(2) or "."), ans
    if signature == "to":
        m = re.search(r"\bto\s+(.+?)([.!?])?$", s, flags=re.I)
        if m:
            ans = (m.group(1) or "").strip().rstrip(".!?")
            if ans:
                return s[:m.start(1)] + ANSWER_BLANK + (m.group(2) or "."), ans

    words = re.findall(r"[A-Za-z][A-Za-z'-]*", s)
    if len(words) >= 5:
        ans = words[-2] if len(words[-1]) <= 3 else words[-1]
        return re.sub(rf"\b{re.escape(ans)}\b", ANSWER_BLANK, s, count=1), ans
    return None


def _official_sentence_frame_fill_items(outline: BookOutline, max_n: int = 4) -> tuple[list[dict], list[str]]:
    """Build sentence items from the official syllabus frame across L3-L6."""
    pattern, example = _syllabus_sentence_source(outline)
    if not (pattern or example):
        return [], []
    signature = _frame_signature(pattern or example)
    candidates: list[tuple[str, str]] = []
    sources = []
    if example:
        sources.append(example)
    sources.extend(_story_sentences_for_grammar(outline))
    seen: set[str] = set()
    for sent in sources:
        key = sent.strip().lower()
        if not key or key in seen:
            continue
        seen.add(key)
        low = f" {sent.lower()} "
        if signature == "will" and " will " not in low:
            continue
        if signature == "there" and not re.search(r"\bthere\s+(is|are)\b", low):
            continue
        if signature == "modal" and not re.search(r"\b(can|could|cannot|can't)\b", low):
            continue
        if signature == "because" and " because " not in low:
            continue
        if signature == "when" and " when " not in low:
            continue
        blanked = _blank_for_frame(sent, signature)
        if blanked:
            candidates.append(blanked)
        if len(candidates) >= max_n:
            break

    if not candidates and example:
        blanked = _blank_for_frame(example, signature)
        if blanked:
            candidates.append(blanked)
    if not candidates and pattern:
        prompt = re.sub(r"\[[^\]]+\]", ANSWER_BLANK, pattern)
        candidates.append((prompt, "own answer"))

    fills = [{"sentence": s, "answer": a} for s, a in candidates[:max_n]]
    bank: list[str] = []
    for _, ans in candidates:
        if ans and ans != "own answer" and ans not in bank:
            bank.append(ans)
    return fills, bank[:max_n]


def _official_sentence_frame_mc_items(
    outline: BookOutline, fills: list[dict], bank: list[str], max_n: int = 4
) -> list[dict]:
    """Sentence recognition page: choose the sentence that matches the syllabus frame."""
    out: list[dict] = []
    seen: set[str] = set()

    def _verb_3rd(base: str) -> str:
        base = (base or "").strip()
        if not base:
            return base
        if base.endswith(("s", "x", "z", "ch", "sh", "o")):
            return base + "es"
        if base.endswith("y") and len(base) > 1 and base[-2].lower() not in "aeiou":
            return base[:-1] + "ies"
        return base + "s"

    def _wrong_sentence_variant(sentence: str) -> str:
        """Make one clear grammar/pattern error without creating nonsense phrases."""
        s = sentence.strip().rstrip(".!?")
        replacements = [
            (r"\bthere are\b", "there is"),
            (r"\bthere is\b", "there are"),
            (r"\bcan ([a-z]+)\b", lambda m: "can " + _verb_3rd(m.group(1))),
            (r"\bwill ([a-z]+)\b", lambda m: "will " + _verb_3rd(m.group(1))),
            (r"\bI wear\b", "I wears"),
            (r"\bI use\b", "I uses"),
            (r"\bPeople wear\b", "People wears"),
            (r"\bWorkers wear\b", "Workers wears"),
            (r"\bFamilies can\b", "Families can"),
            (r"\bSome families are\b", "Some families is"),
            (r"\bSome families (cook|play|go|live|make|help|share|give)\b",
             lambda m: "Some families " + _verb_3rd(m.group(1))),
            (r"\bSome ocean animals live\b", "Some ocean animals lives"),
            (r"\bOcean animals live\b", "Ocean animals lives"),
            (r"\bCoral reefs are\b", "Coral reefs is"),
            (r"\bKelp looks\b", "Kelp look"),
            (r"\bA ([a-z]+) wears\b", lambda m: f"A {m.group(1)} wear"),
            (r"\bThe ([a-z]+) wears\b", lambda m: f"The {m.group(1)} wear"),
            (r"\bThey keep\b", "They keeps"),
            (r"\bThey make\b", "They makes"),
        ]
        for pat, repl in replacements:
            if re.search(pat, s, flags=re.I):
                wrong = re.sub(pat, repl, s, count=1, flags=re.I)
                return wrong.rstrip(".!?") + "."
        return ""

    for item in fills or []:
        prompt = _clean_text(item.get("sentence", ""))
        ans = _clean_text(item.get("answer", ""))
        if not prompt or not ans or "___" not in prompt:
            continue
        correct = re.sub(r"_{3,}", ans, prompt, count=1)
        wrong = _wrong_sentence_variant(correct)
        if not wrong:
            continue
        if not wrong or wrong.lower() == correct.lower():
            continue
        key = correct.lower()
        if key in seen:
            continue
        seen.add(key)
        options = [correct, wrong]
        if (hash(correct) & 1) == 1:
            options = [wrong, correct]
        out.append({"options": options, "correct": options.index(correct)})
        if len(out) >= max_n:
            break
    return out


def _official_sentence_frame_practice_items(outline: BookOutline, max_n: int = 4) -> list[dict]:
    """Second sentence page: guided production using the same official frame."""
    pattern, example = _syllabus_sentence_source(outline)
    if not (pattern or example):
        return []
    display = re.sub(r"\[[^\]]+\]", ANSWER_BLANK, pattern) if pattern else ""
    prompts: list[dict] = []
    if display and ANSWER_BLANK in display and "[" not in pattern:
        prompts.append({"prompt": display})
    elif example:
        blanked = _blank_for_frame(example, _frame_signature(pattern or example))
        prompts.append({"prompt": blanked[0] if blanked else example})

    sig = _frame_signature(pattern or example)
    if sig == "will":
        prompts.extend([
            {"prompt": f"I will {ANSWER_BLANK} first."},
            {"prompt": f"I will {ANSWER_BLANK} every day."},
            {"prompt": f"I will {ANSWER_BLANK} on {ANSWER_BLANK}."},
        ])
    elif sig == "there":
        prompts.extend([
            {"prompt": f"There is {ANSWER_BLANK}."},
            {"prompt": f"There are {ANSWER_BLANK}."},
            {"prompt": f"In the story, there are {ANSWER_BLANK}."},
            {"prompt": f"I can see {ANSWER_BLANK}."},
        ])
    elif sig == "modal":
        prompts.extend([
            {"prompt": f"I can {ANSWER_BLANK}."},
            {"prompt": f"We can {ANSWER_BLANK}."},
            {"prompt": f"The character could {ANSWER_BLANK}."},
        ])
    elif sig == "because":
        prompts.extend([
            {"prompt": f"It is {ANSWER_BLANK} because {ANSWER_BLANK}."},
            {"prompt": f"The character feels {ANSWER_BLANK} because {ANSWER_BLANK}."},
        ])
    elif sig == "when":
        prompts.extend([
            {"prompt": f"I wear {ANSWER_BLANK} when {ANSWER_BLANK}."},
            {"prompt": f"The character {ANSWER_BLANK} when {ANSWER_BLANK}."},
        ])
    else:
        for word in _verbatim_vocab(outline, 3):
            prompts.append({"prompt": f"Write a sentence with {word}: {ANSWER_BLANK}"})

    out: list[dict] = []
    seen: set[str] = set()
    for item in prompts:
        p = item.get("prompt", "").strip()
        if p and p.lower() not in seen:
            seen.add(p.lower())
            out.append({"prompt": p})
        if len(out) >= max_n:
            break
    return out


def _official_sentence_reorder_items(outline: BookOutline, max_n: int = 4) -> list[dict]:
    """Controlled P4 activity: reorder words from real target-pattern sentences."""
    pattern, example = _syllabus_sentence_source(outline)
    signature = _frame_signature(pattern or example)
    sources: list[str] = []
    if example:
        sources.append(example)
    sources.extend(_story_sentences_for_grammar(outline))

    out: list[dict] = []
    seen: set[str] = set()
    for sent in sources:
        clean = _clean_text(sent).strip()
        if not clean or clean.lower() in seen:
            continue
        low = f" {clean.lower()} "
        if signature == "will" and " will " not in low:
            continue
        if signature == "there" and not re.search(r"\bthere\s+(is|are|was|were)\b", low):
            continue
        if signature == "modal" and not re.search(r"\b(can|could|cannot|can't)\b", low):
            continue
        if signature == "because" and " because " not in low:
            continue
        if signature == "when" and " when " not in low:
            continue
        words = re.findall(r"[A-Za-z][A-Za-z'-]*|[.,!?]", clean)
        word_only = [w for w in words if re.search(r"[A-Za-z]", w)]
        if not (4 <= len(word_only) <= 12):
            continue
        order = list(range(len(word_only)))
        rnd = random.Random(hash(clean) & 0xFFFFFFFF)
        for _ in range(10):
            rnd.shuffle(order)
            if order != list(range(len(word_only))):
                break
        scrambled = " / ".join(word_only[i] for i in order)
        answer = clean if clean.endswith((".", "!", "?")) else clean + "."
        out.append({
            "scrambled": scrambled,
            "scrambled_words": word_only,
            "answer": answer,
            "source": "syllabus_example" if clean == example else "book_text",
        })
        seen.add(clean.lower())
        if len(out) >= max_n:
            break
    return out


def _sentence_frame_fill_items(outline: BookOutline, max_n: int = 4) -> tuple[list[dict], list[str]]:
    """L3 A1 sentence practice: assess the syllabus sentence frame with support."""
    sents = _story_sentences_for_grammar(outline)
    candidates: list[tuple[str, str]] = []
    for sent in sents:
        low = sent.lower()
        if " will " not in low:
            continue
        # Keep easy, visible chunks from the actual story.
        replacements = [
            ("do homework", "do homework", ANSWER_BLANK),
            ("clean her room", "clean her room", ANSWER_BLANK),
            ("practice the piano", "practice the piano", ANSWER_BLANK),
            ("play for one hour a day", "play", f"{ANSWER_BLANK} for one hour a day"),
        ]
        for phrase, answer, blanked in replacements:
            if phrase in low:
                pat = re.compile(re.escape(phrase), re.I)
                candidates.append((pat.sub(blanked, sent, count=1), answer))
                break
    if not candidates:
        words = [str(w).strip() for w in (outline.vocabulary_for_display or []) if str(w).strip()]
        for w in words[:max_n]:
            candidates.append((f"I will {ANSWER_BLANK} {w}.", "use"))
    fills = [{"sentence": s, "answer": a} for s, a in candidates[:max_n]]
    bank = []
    for _, ans in candidates:
        if ans not in bank:
            bank.append(ans)
    return fills, bank[:max_n]


def _sentence_frame_copy_items(outline: BookOutline, max_n: int = 4) -> list[dict]:
    fills, _ = _sentence_frame_fill_items(outline, max_n=max_n)
    prompts: list[dict] = []
    for f in fills[:max_n]:
        sent = f.get("sentence", "")
        answer = (f.get("answer") or "").strip().lower()
        if answer in {"play", "play for one hour a day"}:
            sent = f"I will {ANSWER_BLANK} for one hour a day."
            prompts.append({"prompt": sent})
            continue
        sent = re.sub(r"^\s*(She|He|They|Mia)\s+will\s+", "I will ", sent, flags=re.I)
        sent = re.sub(r"_{3,}", ANSWER_BLANK, sent)
        prompts.append({"prompt": sent})
    if not prompts:
        prompts = [
            {"prompt": f"I will {ANSWER_BLANK} first."},
            {"prompt": f"I will {ANSWER_BLANK} on {ANSWER_BLANK}."},
            {"prompt": f"I will {ANSWER_BLANK} every day."},
            {"prompt": f"I will {ANSWER_BLANK} for one hour a day."},
        ]
    return prompts[:max_n]


def _render_sentence_fallback(slide, brand_rgb: tuple, items: list[dict]) -> None:
    """在已建好（含标题）的句型页 slide 上渲染兜底书写题：每题 = 提示句 + 整宽作答横线。

    与 _build_prompt_line_page 同款排版，但独立绘制，供各 builder 在 n==0 时调用，杜绝空页。
    """
    n = min(len(items or []), 5)
    if n == 0:
        return
    area_top = CONTENT_Y + 1.34
    area_bottom = CONTENT_Y + CONTENT_H - 0.40
    avail = area_bottom - area_top
    x = CONTENT_X + 0.62
    box_w = CONTENT_W - 1.24
    slot_h = avail / n
    pt = 16.0
    lh = pt / 72.0 * 1.16
    for i, it in enumerate(items[:n]):
        y = area_top + i * slot_h
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(slot_h - 0.10))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{i + 1}.  {capitalize_names(_clean_text(it.get('prompt', '')))}"
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.color.rgb = BLACK
        # 作答横线（落在题槽下部，留出书写空间）
        line_y = y + slot_h - 0.30
        _draw_writing_line(slide, x, line_y, box_w, LIGHT_GRAY, 1.2)


def _tense_fill_items(outline: BookOutline, max_n: int = 4) -> tuple[list[dict], list[str]]:
    """第 2 句型页（考点=故事时态）：给出原形提示，挖空让学生写出正确过去式。

    例：『Anna ________ (feel) nervous on her first day.』 答案 felt。
    返回 (fills, word_bank=正确过去式列表)。
    """
    fills: list[dict] = []
    bank: list[str] = []
    for sent in _story_sentences_for_grammar(outline):
        if len(fills) >= max_n:
            break
        found = _find_past_verb(sent)
        if not found:
            continue
        tok, base = found
        past_word = "".join(ch for ch in tok if ch.isalpha())
        # 用『________ (base)』替换该动词，保留句中其余部分
        blanked = sent.replace(tok, f"________ ({base})", 1)
        fills.append({"sentence": blanked, "answer": past_word})
        bank.append(past_word)
    return fills, bank


_PLURAL_SUBJECTS = {"they", "we", "you", "i", "children", "kids", "friends", "students"}


def _present_form(base: str, plural: bool) -> str:
    """原形 → 一般现在时（按主语单复数）。"""
    if plural:
        return "are" if base == "be" else base
    return _present_3rd(base)


def _sentence_to_present(sent: str) -> Optional[str]:
    """把整句里所有过去式动词换成一般现在时，得到（错误的）现在时句子；无可换则 None。"""
    first = "".join(ch for ch in sent.split()[0] if ch.isalpha()).lower() if sent.split() else ""
    sent_plural = first in _PLURAL_SUBJECTS
    tokens = sent.split()
    changed = False
    prev_clean = ""  # 动词前一个实词 → 判断主语单复数（修「At recess they plays」）
    for i, tok in enumerate(tokens):
        clean = "".join(ch for ch in tok if ch.isalpha())
        if not clean:
            continue
        if clean.lower() in _NOT_VERB:
            prev_clean = clean.lower()
            continue
        base = _past_to_base(clean)
        if not base or base == clean.lower():
            prev_clean = clean.lower()
            continue  # 非过去式 / 过去=原形同形（put/read），跳过不制造对比
        # 主语单复数：优先看动词紧邻的前一个词（they/we/you/I → 复数），否则用句首判断
        plural = (prev_clean in _PLURAL_SUBJECTS) or (not prev_clean and sent_plural)
        present = _present_form(base, plural)
        if clean[:1].isupper():
            present = present[:1].upper() + present[1:]
        tokens[i] = tok.replace(clean, present, 1)
        prev_clean = clean.lower()
        changed = True
    return " ".join(tokens) if changed else None


def _tense_mc_items(outline: BookOutline, max_n: int = 4) -> list[dict]:
    """第 1 句型页（考点=故事时态）：两句二选一，选『正确过去式』那句。

    正确项 = 故事原句（过去式）；干扰项 = 整句动词换成一般现在时（错误）。
    """
    mcs: list[dict] = []
    for sent in _story_sentences_for_grammar(outline):
        if len(mcs) >= max_n:
            break
        wrong = _sentence_to_present(sent)
        if not wrong or wrong == sent:
            continue
        mcs.append({"options": [sent, wrong], "correct": 0})
    return mcs


# 兼容旧调用名
def _sentence_fill_items(outline: BookOutline, max_n: int = 4) -> tuple[list[dict], list[str]]:
    return _tense_fill_items(outline, max_n=max_n)


# ============================================================
#  句型考点引擎（现在时）：主谓一致 / 现在时填空 —— 适配非虚构等现在时文本
# ============================================================
# 现在时常见动作动词原形（用于在文中可靠识别"现在时动词"并做主谓一致变形）。
# 只对白名单内的词动手，避免把名词误判成动词（如 water/areas）。
_PRESENT_ACTION_BASES: set[str] = {
    "cover", "carry", "flow", "meet", "create", "help", "make", "live", "grow",
    "feed", "give", "take", "bring", "keep", "hold", "move", "show", "need",
    "mean", "use", "form", "join", "reach", "protect", "call", "find", "turn",
    "run", "fall", "rise", "mix", "fill", "feel", "look", "play", "work", "want",
    "like", "love", "eat", "drink", "swim", "jump", "see", "hear", "know",
    "think", "say", "tell", "read", "write", "sing", "walk", "talk", "open",
    "close", "start", "stop", "build", "plant", "catch", "throw", "fly", "drive",
    "wash", "shine", "melt", "freeze", "burn", "blow", "sail", "float", "sink",
    "include", "contain", "support", "produce", "provide", "connect", "spread",
    "begin", "happen", "appear", "remain", "become", "belong", "depend",
}


def _base_of_present(low: str) -> Optional[tuple[str, str]]:
    """识别一个现在时动词 token（小写、纯字母）→ 返回 (原形 base, 数 'sing'/'plur')。

    'sing' = 第三人称单数形态（covers/carries/is/has），'plur' = 原形/复数形态（cover/are/have）。
    仅命中白名单或 be/have 才返回，杜绝把名词误判为动词。"""
    if low in ("is", "has"):
        return ("be" if low == "is" else "have", "sing")
    if low in ("are", "have"):
        return ("be" if low == "are" else "have", "plur")
    if low in _PRESENT_ACTION_BASES:
        return (low, "plur")                       # 原形 = 复数/第一二人称
    # 第三人称单数 -s/-es/-ies → 还原原形再核对白名单
    if low.endswith("ies") and len(low) > 3:
        cand = low[:-3] + "y"
        if cand in _PRESENT_ACTION_BASES:
            return (cand, "sing")
    if low.endswith("es") and len(low) > 2:
        cand = low[:-2]
        if cand in _PRESENT_ACTION_BASES:
            return (cand, "sing")
    if low.endswith("s") and len(low) > 1:
        cand = low[:-1]
        if cand in _PRESENT_ACTION_BASES:
            return (cand, "sing")
    return None


def _wrong_agreement_form(base: str, number: str) -> str:
    """给出"错误的"主谓一致形态：单数↔复数互换。"""
    if base == "be":
        return "are" if number == "sing" else "is"
    if base == "have":
        return "have" if number == "sing" else "has"
    return base if number == "sing" else _present_3rd(base)


def _swap_token(sentence: str, old_tok: str, new_word: str) -> str:
    """把句中第一处 old_tok 换成 new_word，保留首字母大小写。"""
    clean = "".join(ch for ch in old_tok if ch.isalpha())
    if clean[:1].isupper():
        new_word = new_word[:1].upper() + new_word[1:]
    replaced = old_tok.replace(clean, new_word, 1)
    return sentence.replace(old_tok, replaced, 1)


def _find_present_verb(sentence: str) -> Optional[tuple[str, str, str]]:
    """找句中第一个可做主谓一致变形的现在时动词。返回 (原 token, base, number)。"""
    for tok in sentence.split():
        clean = "".join(ch for ch in tok if ch.isalpha())
        low = clean.lower()
        if not clean or low in _NOT_VERB:
            continue
        info = _base_of_present(low)
        if info:
            base, number = info
            return tok, base, number
    return None


def _present_agreement_mc_items(outline: BookOutline, max_n: int = 4) -> list[dict]:
    """现在时主谓一致二选一：正确句 vs 主谓不一致（错误）句。

    正确项 = 故事原句；干扰项 = 把主要现在时动词换成错误的单复数形态。
    """
    mcs: list[dict] = []
    seen: set[str] = set()
    for sent in _story_sentences_for_grammar(outline):
        if len(mcs) >= max_n:
            break
        found = _find_present_verb(sent)
        if not found:
            continue
        tok, base, number = found
        wrong_form = _wrong_agreement_form(base, number)
        wrong = _swap_token(sent, tok, wrong_form)
        if wrong == sent or sent in seen:
            continue
        seen.add(sent)
        mcs.append({"options": [sent, wrong], "correct": 0})
    return mcs


def _present_fill_items(outline: BookOutline, max_n: int = 4) -> tuple[list[dict], list[str]]:
    """现在时填空：挖空主要现在时动词，括号给原形，学生写出正确形态。

    例：『Oceans ________ (cover) most of Earth.』 答案 cover。
    返回 (fills, word_bank=正确形态列表)。
    """
    fills: list[dict] = []
    bank: list[str] = []
    seen: set[str] = set()
    for sent in _story_sentences_for_grammar(outline):
        if len(fills) >= max_n:
            break
        found = _find_present_verb(sent)
        if not found:
            continue
        tok, base, _number = found
        ans = "".join(ch for ch in tok if ch.isalpha())
        prompt_base = "be" if base == "be" else base
        blanked = sent.replace(tok, f"________ ({prompt_base})", 1)
        if blanked == sent or sent in seen:
            continue
        seen.add(sent)
        fills.append({"sentence": blanked, "answer": ans})
        bank.append(ans)
    return fills, bank


def _dominant_tense(outline: BookOutline) -> str:
    """统计故事主时态：返回 'past' 或 'present'。

    逐句看首个可识别动词是过去式还是现在时；过去式句占比 ≥ 40% → 'past'，否则 'present'。
    （非虚构科普文几乎全是一般现在时 → 'present'，从而切换到现在时考点。）"""
    sents = _story_sentences_for_grammar(outline)
    if not sents:
        return "past"
    past = present = 0
    for s in sents:
        if _find_past_verb(s):
            past += 1
        elif _find_present_verb(s):
            present += 1
    total = past + present
    if total == 0:
        return "past"
    return "past" if (past / total) >= 0.40 else "present"


# ============================================================
#  分级出题：词汇/句型按 Level 选型（用户拍板 2026-06-04）
# ============================================================
def _ws_seed(outline: BookOutline) -> int:
    """由书名+级别派生稳定轮换种子（同书恒定、不同书各异）——用于在【版式恒定】的前提下
    跨书轮换题型，告别"每本雷同"。与 ai_extractor._pool_seed 同套路，保证确定性。"""
    parts = [str(getattr(outline, "title", "") or ""), str(getattr(outline, "level", "") or "")]
    s = "|".join(p for p in parts if p)
    if not s:
        return 0
    return int(hashlib.md5(s.encode("utf-8")).hexdigest()[:8], 16)


def _lvl_band(lvl_n: int) -> str:
    """L0-2→'l02'，L3-4→'l34'，L5-6→'l56'。"""
    if lvl_n <= 2:
        return "l02"
    if lvl_n <= 4:
        return "l34"
    return "l56"


def _clean_text(s) -> str:
    """规整题面文本的空白（用户拍板 2026-06-09，修"词距忽大忽小"）。

    成因：AI/RR 抽出的句子常混入【不间断空格 U+00A0 / 制表符 / 重复空格 / 窄空格】等
    不可见的异形空白，渲染时这些空白宽度与普通空格不同 → 同一行词距忽大忽小。
    本函数把所有连续空白（含 Unicode 空白）统一压成单个普通空格并去首尾，
    保证整行左对齐时词距均匀（不动可见字符，不影响标点/大小写）。"""
    text = (
        str(s if s is not None else "")
        .replace("\u201c", '"').replace("\u201d", '"')
        .replace("\u2018", "'").replace("\u2019", "'")
        .replace("\u02bc", "'").replace("\u2032", "'").replace("\u00b4", "'").replace("`", "'")
        .replace("\u2033", '"')
        .replace("\uff02", '"').replace("\uff07", "'")
        .replace("\ufffe", "-").replace("\ufffd", "-").replace("\u00ad", "-")
        .replace("\u2010", "-").replace("\u2011", "-")
        .replace("\u2012", "-").replace("\u2013", "-").replace("\u2014", "-")
    )
    text = re.sub(r"(?<=[A-Za-z])\?(?=s\b)", "'", text)
    text = re.sub(r"\?([A-Za-z]{1,12})\?", r'"\1"', text)
    return re.sub(r"\s+", " ", text).strip()


def _first_letter_mask(word: str) -> str:
    """nervous → 'n _ _ _ _ _ _'（首字母提示 + 其余下划线）。"""
    w = (word or "").strip()
    if len(w) <= 1:
        return w
    return w[0] + " " + " ".join("_" for _ in w[1:])


def _suffix_mask(word: str) -> str:
    """nervous → 'nerv _ _ _'（保留词根、空出词尾；一个下划线=一个字母）。"""
    w = (word or "").strip()
    n = len(w)
    cut = max(2, n - 3)
    blanks = " ".join("_" for _ in range(n - cut))
    return w[:cut] + " " + blanks if blanks else w


def _word_fill_meaning_items(pairs: list[dict], max_n: int = 5) -> list[dict]:
    """L0-2 看词义/首字母写词：clue=释义，hint=首字母掩码。"""
    out: list[dict] = []
    for p in pairs[:max_n]:
        w = str(p.get("word", "")).strip()
        d = str(p.get("def", "")).strip()
        if not w or " " in w or _bad_placeholder_text(d):
            continue
        out.append({"clue": d, "answer": w, "hint": _first_letter_mask(w)})
    return out


def _bad_placeholder_text(text: str) -> bool:
    low = str(text or "").strip().lower()
    if not low:
        return True
    bad = (
        "placeholder", "todo", "word from the story", "meaning of ",
        "look back at the book", "text does not say this", "it is not true that",
        "copied exactly", "ai generated",
    )
    return any(x in low for x in bad)


def _valid_definition_pair(pair: dict) -> bool:
    word = _clean_text(pair.get("word", ""))
    definition = _clean_text(pair.get("def", ""))
    if not word or not definition or _bad_placeholder_text(definition):
        return False
    return word.lower() not in {x.lower() for x in re.findall(r"[A-Za-z][A-Za-z'-]*", definition)}


def _valid_mc_item(item: dict, *, min_options: int = 2, max_options: int = 4) -> bool:
    q = _clean_text(item.get("q", ""))
    opts = [_clean_text(o) for o in (item.get("options") or [])[:max_options]]
    opts = [o for o in opts if o]
    if not q or _bad_placeholder_text(q) or len(opts) < min_options:
        return False
    if len({o.lower() for o in opts}) != len(opts):
        return False
    try:
        correct = int(item.get("correct", 0))
    except Exception:
        return False
    return 0 <= correct < len(opts)


def _valid_short_item(item: dict) -> bool:
    q = _clean_text(item.get("q", ""))
    if not q or _bad_placeholder_text(q):
        return False
    if q.lower().startswith(("what ", "who ", "where ", "when ", "why ", "how ")):
        return "?" in q
    return True


def _valid_tf_item(item: dict) -> bool:
    q = _clean_text(item.get("q", ""))
    return bool(q and not _bad_placeholder_text(q) and q.endswith("."))


def _sanitize_reading_items(items: list[dict], *, preferred_kind: str | None = None,
                            max_n: int = 5) -> list[dict]:
    out: list[dict] = []
    seen: set[str] = set()
    for item in items or []:
        kind = item.get("kind", preferred_kind or "")
        ok = (
            _valid_mc_item(item, min_options=2) if kind == "mc"
            else _valid_tf_item(item) if kind == "tf"
            else _valid_short_item(item)
        )
        if not ok:
            continue
        key = _reading_question_key(item.get("q", ""))
        if kind == "mc":
            opt_key = "|".join(_clean_text(o).lower() for o in (item.get("options") or []))
            key = f"{key}|{opt_key}"
        if not key or key in seen:
            continue
        seen.add(key)
        out.append(item)
        if len(out) >= max_n:
            break
    return out


def _word_fill_morph_items(pairs: list[dict], max_n: int = 5) -> list[dict]:
    """L5-6 构词/拼写补全：clue=释义，hint=词根+空词尾。"""
    out: list[dict] = []
    for p in pairs[:max_n]:
        w = str(p.get("word", "")).strip()
        if len(w) < 4:
            continue
        d = str(p.get("def", "")).strip()
        out.append({"clue": d or "Complete the word.", "answer": w, "hint": _suffix_mask(w)})
    return out


def _unscramble_items(sentences: list[str], max_n: int = 5) -> list[dict]:
    """L0-2 连词成句：把 3-9 词的句子打乱词序。"""
    import random as _rnd
    out: list[dict] = []
    for s in sentences:
        words = s.rstrip(".!?").split()
        if not (3 <= len(words) <= 9):
            continue
        order = list(range(len(words)))
        for _ in range(8):
            _rnd.shuffle(order)
            if order != list(range(len(words))):
                break
        scrambled = "  /  ".join(words[i] for i in order)
        out.append({"scrambled": scrambled, "answer": s})
        if len(out) >= max_n:
            break
    return out


def _rewrite_items(outline: BookOutline, max_n: int = 5) -> list[dict]:
    """L5-6 句子改写：给现在时句子，改写成故事原句（过去时）。"""
    out: list[dict] = []
    for s in _story_sentences_for_grammar(outline):
        pres = _sentence_to_present(s)
        if not pres or pres == s:
            continue
        out.append({"prompt": pres, "answer": s})
        if len(out) >= max_n:
            break
    return out


# ---------- L3 题型（用户拍板 2026-06-09，对齐官方样板 L3-8 / L3-23） ----------

def _missing_letters_mask(word: str) -> str:
    """nest → 'n _ _ t'：保留首尾字母、中间逐字母挖空（一个下划线=一个字母）。
    用于 L3 词汇页"看图填缺失字母"。"""
    w = (word or "").strip()
    if len(w) <= 2:
        return w
    chars = list(w)
    out = [chars[0]]
    for c in chars[1:-1]:
        out.append("_" if c.isalpha() else c)
    out.append(chars[-1])
    return " ".join(out)


def _word_fill_pic_items(pairs: list[dict], images: list[Path], max_n: int = 4) -> list[dict]:
    """L3 看图填缺字母：每题 = 配图 + 缺字母掩码 + 小释义提示。
    配图复用绘本内页图（page_02 起，跳过封面/P1），作为视觉支持（L3 学生程度不高需要图）。"""
    out: list[dict] = []
    n_img = len(images or [])
    used = 0
    for p in pairs:
        w = str(p.get("word", "")).strip()
        if not w or " " in w or len(w) < 3:
            continue
        d = str(p.get("def", "")).strip()
        if d.lower().startswith("word from the story:"):
            d = ""
        img = None
        if n_img > 2:
            idx = 2 + (used % (n_img - 2))
            cand = images[idx]
            img = cand if (cand and Path(cand).exists()) else None
        out.append({"clue": d, "answer": w, "hint": _missing_letters_mask(w), "img": img})
        used += 1
        if len(out) >= max_n:
            break
    return out


def _clean_vocab_clue_kind(word: str) -> str:
    w = (word or "").strip().lower()
    if w in {"week", "day", "days"}:
        return "calendar"
    if w in {"homework", "study", "write", "worksheet"}:
        return "desk"
    if w in {"plan", "schedule"}:
        return "notebook"
    if w in {"practice", "piano", "music"}:
        return "piano"
    return ""


def _riddle_mc_items(pairs: list[dict], max_n: int = 4) -> list[dict]:
    """L3 词汇②：四选一谜语猜词。stem=释义，options=正确词 + 同表干扰词。"""
    import random as _rnd
    clean_pairs = [p for p in (pairs or []) if _valid_definition_pair(p)]
    words = [str(p.get("word", "")).strip() for p in clean_pairs if str(p.get("word", "")).strip()]
    out: list[dict] = []
    for p in clean_pairs:
        w = str(p.get("word", "")).strip()
        d = str(p.get("def", "")).strip()
        distractors = [x for x in words if x.lower() != w.lower()]
        rnd = _rnd.Random(hash(w) & 0xFFFFFFFF)
        rnd.shuffle(distractors)
        opts = [w] + distractors[:2]
        if len(opts) < 2:
            continue
        rnd.shuffle(opts)
        item = {
            "kind": "mc",
            "q": f"Which word means: {d}?",
            "options": opts,
            "correct": opts.index(w),
        }
        if _valid_mc_item(item, min_options=2):
            out.append(item)
        if len(out) >= max_n:
            break
    return out


def _scramble_letters(word: str, *, seed: int | None = None) -> str:
    """把一个单词的字母【真正打乱】，返回空格分隔的乱序字母（如 'sing' → 'g n i s'）。

    用确定性随机（默认按词本身做种）保证可复现；尽量保证：① 结果 != 原词；
    ② 首字母也被换掉（让乱序更明显，杜绝 's i n g' 这种"没打乱"的假题）。
    无法打乱（如单字母/全同字母）时原样返回。"""
    import random as _rnd
    w = (word or "").strip()
    letters = [c for c in w if c.isalpha()]
    if len(letters) <= 1:
        return " ".join(letters)
    rnd = _rnd.Random(seed if seed is not None else (hash(w.lower()) & 0xFFFFFFFF))
    order = letters[:]
    best = None
    for _ in range(16):
        rnd.shuffle(order)
        if "".join(order).lower() == w.lower():
            continue
        best = order[:]
        if len(letters) <= 2 or order[0].lower() != letters[0].lower():
            break
    if best is None:                       # 全同字母等极端情况：保底反转
        best = list(reversed(letters))
    return " ".join(best)


def _word_unscramble_items(unscramble_data: list[dict] | None, pairs: list[dict],
                           max_n: int = 4) -> list[dict]:
    """L3 词汇②：把单词字母打乱，让学生拼回原词（作答方式【单一】=填空写词，不带选项）。

    优先用 AI 标记为 unscramble 的词；不足时用词汇表(match_pairs)的词补齐。
    返回 _build_word_fill_page 可直接渲染的 {clue, hint, answer} 列表：
      clue = 真打乱后的字母（+ 可选释义提示，给 A1 学生作答支撑）；
      hint = 箭头 + 与字母数等长的作答下划线（每个下划线=一个字母）。"""
    seen: set[str] = set()
    out: list[dict] = []
    # 词→释义 速查（给 A1 学生作答支撑：拼词时附一句简短释义）
    meaning_map = {
        str(p.get("word", "")).strip().lower(): str(p.get("def", "")).strip()
        for p in (pairs or []) if str(p.get("word", "")).strip()
    }

    def _push(word: str, meaning: str = "") -> None:
        w = (word or "").strip()
        low = w.lower()
        if not w or " " in w or not w.isalpha() or len(w) < 3 or low in seen:
            return
        scr = _scramble_letters(w)
        if scr.replace(" ", "").lower() == low:   # 没打乱成功 → 跳过，杜绝假乱序
            return
        seen.add(low)
        clue = scr
        m = (meaning or "").strip() or meaning_map.get(low, "")
        if m and not m.lower().startswith("meaning of"):
            clue = f"{scr}    ({m})"
        out.append({
            "clue": clue,
            "hint": "\u2192   " + " ".join("_" for _ in w),
            "answer": w,
        })

    for it in (unscramble_data or []):
        if len(out) >= max_n:
            break
        _push(str(it.get("answer", "")), str(it.get("clue", "")))
    for p in (pairs or []):
        if len(out) >= max_n:
            break
        _push(str(p.get("word", "")), str(p.get("def", "")))
    return out


def _cloze_mc_items(fills: list[dict], word_bank: list[str], max_n: int = 4) -> list[dict]:
    """L3 句型②：选词补全句子（finish the sentence）。
    stem=挖空句，options=正确词 + 词库干扰词。

    防御（用户拍板 2026-06-09）：本函数专供【句型/Sentences 区】，只接受【真正的句子级】
    挖空题。拒绝被误塞进 fill_blanks 的单词级拼词/改写/排序条目，以及挖空后实词不足 3 个的
    伪句子——杜绝"Sentences"区再出现 'Unscramble: s i n g → __' 这类单词题。"""
    import random as _rnd
    bank = [str(w).strip() for w in (word_bank or []) if str(w).strip()]
    out: list[dict] = []
    for f in (fills or []):
        sent = (f.get("sentence") or "").strip()
        ans = (f.get("answer") or "").strip()
        if not sent or not ans or "____" not in sent:
            continue
        low = sent.lower()
        if low.startswith("unscramble") or low.startswith("rewrite") \
                or sent.lstrip().startswith("___"):
            continue
        real_words = [w for w in re.sub(r"_+", " ", sent).split()
                      if any(c.isalpha() for c in w)]
        if len(real_words) < 3:
            continue
        distractors = [x for x in bank if x.lower() != ans.lower()]
        rnd = _rnd.Random(hash(sent) & 0xFFFFFFFF)
        rnd.shuffle(distractors)
        opts = [ans] + distractors[:3]
        if len(opts) < 3:
            continue
        rnd.shuffle(opts)
        stem = sent.replace("____", "______")
        out.append({
            "kind": "mc",
            "q": stem,
            "options": opts,
            "correct": opts.index(ans),
        })
        if len(out) >= max_n:
            break
    return out


def _build_l34_vocab2(new_page, brand_rgb: tuple, data: dict,
                      images: Optional[list[Path]], seed: int,
                      lvl_n: int = 3,
                      outline: Optional[BookOutline] = None) -> None:
    """L3/L4 vocabulary application page with deterministic activity rotation."""
    pairs = data.get("match_pairs") or []
    fills = data.get("fill_blanks") or []
    bank = data.get("word_bank") or []
    activity_codes = _l34_activity_order(outline, 2, seed=seed, include_optional=False)
    core_words_for_type = [
        str(p.get("word", "")).strip()
        for p in (pairs or [])
        if str(p.get("word", "")).strip()
    ] or [str(w).strip() for w in (bank or []) if str(w).strip()]
    has_phrase_vocab = any(" " in w for w in core_words_for_type)
    if has_phrase_vocab:
        phrase_first = ["vocab_multi_word_cloze", "vocab_contextual_word_bank_cloze"]
        activity_codes = phrase_first + [c for c in activity_codes if c not in phrase_first]
    else:
        activity_codes = [c for c in activity_codes if c != "vocab_multi_word_cloze"]

    def _do_context_fill() -> bool:
        code = "vocab_multi_word_cloze" if any(" " in str(w).strip() for w in bank) else "vocab_contextual_word_bank_cloze"
        if len(fills) >= _l34_min_items(code, 3):
            _record_l34_activity(outline, 2, code)
            _record_worksheet_items(outline, 2, fills)
            _build_p2_fill(new_page(), brand_rgb, fills, bank, images)
            return True
        return False

    def _do_clue_choice() -> bool:
        riddles = _riddle_mc_items(pairs, max_n=4)
        if len(riddles) >= _l34_min_items("vocab_contextual_choice", 3):
            _record_l34_activity(outline, 2, "vocab_contextual_choice")
            _record_worksheet_items(outline, 2, riddles)
            _build_mcq_page(new_page(), brand_rgb, "Vocabulary",
                            _ws_activity_instruction("vocab_contextual_choice"),
                            riddles, answer_paren=False)
            return True
        return False

    def _do_write_word() -> bool:
        items = _word_fill_meaning_items(pairs, max_n=4)
        if len(items) >= _l34_min_items("vocab_contextual_clue_write", 3):
            _record_l34_activity(outline, 2, "vocab_contextual_clue_write")
            _record_worksheet_items(outline, 2, items)
            _build_word_fill_page(new_page(), brand_rgb, items, "Vocabulary",
                                  _ws_activity_instruction("vocab_contextual_clue_write"))
            return True
        return False

    def _do_definition_match() -> bool:
        if len([p for p in pairs if _valid_definition_pair(p)]) >= _l34_min_items("vocab_word_definition_matching", 4):
            _record_l34_activity(outline, 2, "vocab_word_definition_matching")
            _record_worksheet_items(outline, 2, pairs)
            _build_p1_match(new_page(), brand_rgb, pairs, images or [])
            return True
        return False

    choices = {
        "vocab_contextual_word_bank_cloze": _do_context_fill,
        "vocab_multi_word_cloze": _do_context_fill,
        "vocab_contextual_choice": _do_clue_choice,
        "vocab_contextual_clue_write": _do_write_word,
        "vocab_word_definition_matching": _do_definition_match,
    }
    for code in activity_codes:
        fn = choices.get(code)
        if fn and fn():
            return
    _record_l34_activity(outline, 2, "vocab_contextual_word_bank_cloze")
    _record_worksheet_items(outline, 2, fills)
    _build_p2_fill(new_page(), brand_rgb, fills, bank, images)


def _build_l34_sentence2(new_page, brand_rgb: tuple, outline: BookOutline,
                         frame_fills: list[dict], frame_bank: list[str],
                         fallback: list[dict], seed: int,
                         lvl_n: int = 3) -> None:
    """Second L3/L4 sentence page with stable variety while preserving syllabus frame."""
    frame_copy = _official_sentence_frame_practice_items(outline, max_n=4)
    cloze_mc = _cloze_mc_items(frame_fills, frame_bank, max_n=4)
    reorder_items = _official_sentence_reorder_items(outline, max_n=4)
    activity_codes = _l34_activity_order(outline, 4, seed=seed // 5, include_optional=False)

    def _do_guided_write() -> bool:
        if len(frame_copy) >= _l34_min_items("sentence_guided_writing", 3):
            _record_l34_activity(outline, 4, "sentence_guided_writing")
            _record_worksheet_items(outline, 4, frame_copy)
            _build_prompt_line_page(
                new_page(), brand_rgb, frame_copy, "Sentences",
                f"{_ws_activity_instruction('sentence_guided_writing')} Follow the example: {_display_sentence_frame(outline)}",
                prompt_key="prompt",
                fallback=fallback,
            )
            return True
        return False

    def _do_choose_word() -> bool:
        if len(cloze_mc) >= _l34_min_items("sentence_grammar_word_cloze", 3):
            _record_l34_activity(outline, 4, "sentence_grammar_word_cloze")
            _record_worksheet_items(outline, 4, cloze_mc)
            _build_mcq_page(new_page(), brand_rgb, "Sentences",
                            _ws_activity_instruction("sentence_grammar_word_cloze"),
                            cloze_mc)
            return True
        return False

    def _do_correct_sentence() -> bool:
        if len(frame_fills) >= _l34_min_items("sentence_grammar_correction", 3):
            items = []
            for item in frame_fills[:4]:
                sent = _clean_text(item.get("sentence", ""))
                ans = _clean_text(item.get("answer", ""))
                if not sent or not ans or "___" not in sent:
                    continue
                correct = re.sub(r"_{3,}", ans, sent, count=1)
                wrong = _false_tf_variant(correct).rstrip(".!?") + "."
                if wrong and wrong.lower() != correct.lower():
                    items.append({
                        "prompt": f"Correct the sentence: {wrong}",
                        "wrong_sentence": wrong,
                        "correct_sentence": correct,
                    })
            if len(items) >= _l34_min_items("sentence_grammar_correction", 3):
                _record_l34_activity(outline, 4, "sentence_grammar_correction")
                _record_worksheet_items(outline, 4, items)
                _build_prompt_line_page(
                    new_page(), brand_rgb, items, "Sentences",
                    _ws_activity_instruction("sentence_grammar_correction"),
                    prompt_key="prompt",
                    fallback=fallback,
                )
                return True
        return False

    def _do_reorder_words() -> bool:
        if len(reorder_items) >= _l34_min_items("sentence_reorder_words", 3):
            _record_l34_activity(outline, 4, "sentence_reorder_words")
            _record_worksheet_items(outline, 4, reorder_items)
            _build_prompt_line_page(
                new_page(), brand_rgb, reorder_items, "Sentences",
                _ws_activity_instruction("sentence_reorder_words"),
                prompt_key="scrambled",
                fallback=fallback,
            )
            return True
        return False

    choices = {
        "sentence_complete_frame": _do_guided_write,
        "sentence_guided_writing": _do_guided_write,
        "sentence_grammar_word_cloze": _do_choose_word,
        "sentence_grammar_correction": _do_correct_sentence,
        "sentence_reorder_words": _do_reorder_words,
    }
    for code in activity_codes:
        fn = choices.get(code)
        if fn and fn():
            return
    _do_guided_write()


# ---------- 官方 A4 模板（底版/母版）----------
# 7 个 slide：index = 级别数字（Smart=0, L1=1 … L6=6），各带原生 Logo 图 / Name 角标 / 配色 / 页脚。
_WS_TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "templates" / "Worksheet_A4_L0-L6.pptx"


def _template_slide_index(level: str) -> int:
    return max(0, min(6, _level_num(level)))


def _clone_template_slide(prs: Presentation, src_slide):
    """把模板某个级别的 slide 整体克隆成新 slide（保留原生 Logo/Name 角标/配色/页脚）。
    深拷贝形状 XML 并重映射图片关系 rId，确保图片正常显示。"""
    import copy
    from pptx.oxml.ns import qn

    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # 清掉版式带进来的占位形状
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # 重建图片关系 rId 映射（老 rId → 新 rId）
    rid_map: dict[str, str] = {}
    for rid, rel in src_slide.part.rels.items():
        if "image" in rel.reltype:
            rid_map[rid] = new_slide.part.relate_to(rel._target, rel.reltype)
    # 拷贝形状并改写 embed/link 引用
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in ("embed", "link"):
                a = qn("r:" + attr)
                old = node.get(a)
                if old in rid_map:
                    node.set(a, rid_map[old])
        new_slide.shapes._spTree.append(el)
    return new_slide


def _set_footer(slide, footer_text: str) -> None:
    """把克隆 slide 右下页脚文本框的内容改成本书的 'Level X - Title'。
    找不到就在官方页脚位置新建一个。"""
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        left_in = sh.left / 914400 if sh.left is not None else 0
        top_in = sh.top / 914400 if sh.top is not None else 0
        if top_in > 6.8 and left_in > 6.0:  # 右下角区域
            target = sh
            break
    if target is None:
        target = slide.shapes.add_textbox(
            Inches(7.46), Inches(7.16), Inches(2.94), Inches(0.35))
    tf = target.text_frame
    tf.clear()
    # 不换行、不自动缩放 + 给足宽度 → 长标题不再被截断
    tf.word_wrap = False
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        target.left = Inches(3.0)
        target.width = Inches(SLIDE_W - 3.0 - 0.33)
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = footer_text
    r.font.name = FONT
    # 标题长则自动降字号，保证整行不溢出
    fp = FOOTER_PT
    if len(footer_text) > 42:
        fp = 11.0
    elif len(footer_text) > 32:
        fp = 12.5
    r.font.size = Pt(fp)
    r.font.color.rgb = WHITE


def _fix_name_badge(slide) -> None:
    """Only adjust the template's built-in Name badge; never add another one."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = (sh.text_frame.text or "").strip().lower()
        if txt in ("name", "name:", "姓名", "姓名:"):
            tf = sh.text_frame
            tf.word_wrap = False
            try:
                tf.auto_size = MSO_AUTO_SIZE.NONE
            except Exception:
                pass
            return


def _retell_writing_scaffold(outline: BookOutline) -> dict:
    """生成【基于本文故事】的引导式复述写作脚手架（v2.2）。

    用户拍板：写作题必须扣住本篇文章——让孩子复述/分析这个故事，而不是写空泛抽象题
    （如 "Write about friendship"，孩子无从下手）。这里把 5 步换成针对本故事的引导问题，
    读完原文就能照着一步步写出来。step_labels 是引导问句（孩子照着回答即可成文）。
    """
    title = (outline.title or "the story").strip()
    return {
        "theme": outline.theme or title,
        "title": f"Retell the story: {title}",
        "subtitle": "Read the story again, then retell it in your own words.",
        "steps": ["", "", "", "", ""],
        # 引导问句（通用于任何叙事文，且明确指向"本文"）：
        "step_labels": [
            "Who is the story about, and how did they feel at the start?",
            "What happened first?",
            "What happened next?",
            "What was the best, funniest or most surprising part?",
            "How did the story end? How did they feel?",
        ],
        "min_words": 50, "max_words": 80,
    }


def _resolve_second_reading_mode(mode: str, lvl_n: int) -> str:
    """解析第 2 张 Reading 页内容模式（用户拍板 2026-06-04，P6 = 分级 PBL）。

    auto 分级（用户拍板 2026-06-09：L3 阅读② 改为"涂色 + 自主造句"简单迷你项目，
    贴近样板 L3-23 的"画图 + 补全句子"，并结合昨天讨论的"尽量涂色让孩子自主造句"）：
      • L0–L3 → color_say  涂色线稿 + 完成一句简单句子（动手 + 自主造句 + 词库辅助）
      • L4    → mindmap    思维导图(SWBST)
      • L5–L6 → writing    写作脚手架（顶部 SWBST 规划框 + 写作区 = 思维导图+写作）
    其余取值原样返回：reading / mindmap / writing / pbl / color_say。
    """
    m = (mode or "auto").strip().lower()
    if m in ("reading", "mindmap", "writing", "writing_official", "pbl", "color_say",
             "l3summary", "l3bubble", "timeline", "planchart"):
        return m
    if lvl_n <= 2:
        return "color_say"
    if lvl_n == 3:
        # L3 第②阅读页按内容自选（在调用处用 outline 体裁细化）：默认思维导图，
        # 非虚构走写作/复述脚手架。占位返回 auto_l3，由调用处解析。
        return "auto_l3"
    if lvl_n == 4:
        return "mindmap"
    return "writing"


def _second_reading_signal(outline: BookOutline) -> str:
    """Collect stable syllabus signals used to pick the final Reading page."""
    parts = [
        getattr(outline, "title", ""),
        getattr(outline, "theme", ""),
        getattr(outline, "fiction_type", ""),
        getattr(outline, "reading_skill", ""),
        getattr(outline, "reading_strategy", ""),
        getattr(outline, "graphic_organizer", ""),
        getattr(outline, "graphic_organizer_desc", ""),
    ]
    syl = getattr(outline, "syllabus", None)
    if syl is not None:
        parts.extend([
            getattr(syl, "reading_skill", ""),
            getattr(syl, "reading_strategy", ""),
            getattr(syl, "graphic_organizer", ""),
            getattr(syl, "go_description", ""),
        ])
    return " ".join(str(p or "") for p in parts).lower()


def _auto_second_reading_mode(outline: BookOutline, lvl_n: int) -> str:
    """SOP 6-page default: Reading page 2 is a compact GO/writing task.

    The 8-page SOP is a candidate-page workflow for outsourcing. The app
    directly outputs the curated 6-page version, so the final Reading page
    adapts to the book instead of hardcoding Timeline for every title.
    """
    signal = _second_reading_signal(outline)
    is_nonfic = "non" in (getattr(outline, "fiction_type", "") or "").lower()

    if any(k in signal for k in (
        "plan", "schedule", "seven-day", "seven day",
    )):
        return "planchart"
    if any(k in signal for k in (
        "timeline", "sequence", "sequencing", "steps in a process",
        "process", "order", "first", "next", "then", "finally",
    )):
        return "timeline"
    if any(k in signal for k in (
        "main idea", "supporting details", "fact web", "bubble map",
        "classification", "classify", "category", "categories", "kwl",
    )) or is_nonfic:
        return "l3bubble"
    if any(k in signal for k in (
        "text-to-self", "text to self", "text-to-world", "text to world",
        "connection", "connections",
    )):
        return "writing" if lvl_n >= 5 else "l3summary"
    if any(k in signal for k in (
        "story elements", "character", "setting", "problem", "solution",
        "before", "event", "after",
    )):
        return "mindmap"
    if lvl_n <= 2:
        return "color_say"
    if lvl_n in (3, 4):
        return "l3summary"
    return "writing"


def _unify_reading_questions(rq: list[dict]) -> tuple[list[dict], str]:
    """同页题型统一（用户硬要求）：阅读页只能"全选择"或"全判断"，不混排 short。

    取数量更多的同质题型（mc / tf）；不足 3 题则退回原始混合，避免空页。
    返回 (questions, kind)，kind ∈ {"mc","tf","mixed"}。
    """
    def _difficulty_key(q: dict) -> tuple:
        # 易→难排序（不改版式，只调顺序，方便孩子顺手作答）：
        # 闭合题(mc/tf)先于开放题(short)；指向越靠前页的事实题越简单先排；题干越短越简单。
        kind_rank = {"tf": 0, "mc": 1, "short": 2}.get((q.get("kind") or "").lower(), 3)
        try:
            page = int(q.get("page") or 99)
        except (TypeError, ValueError):
            page = 99
        return (kind_rank, page, len((q.get("q") or "")))

    def _dedupe(items: list[dict]) -> list[dict]:
        out: list[dict] = []
        seen: set[str] = set()
        for q in items:
            text = q.get("q") or ""
            if _is_generic_reading_prompt(text):
                continue
            key = _reading_question_key(text)
            if not key or key in seen:
                continue
            seen.add(key)
            out.append(q)
        return out

    mc = _dedupe([q for q in rq if q.get("kind") == "mc" and q.get("q")])
    tf = _dedupe([q for q in rq if q.get("kind") == "tf" and q.get("q")])
    chosen, kind = (mc, "mc") if len(mc) >= len(tf) else (tf, "tf")
    if len(chosen) < 3:
        # 混排兜底：按易→难稳定排序
        return sorted(_dedupe([q for q in rq if q.get("q")]), key=_difficulty_key), "mixed"
    # 同质题内部也按易→难（靠前页事实题先、短题先）稳定排序
    return sorted(chosen, key=_difficulty_key), kind


def _reading_subtitle(kind: str) -> str:
    """副标题随题型自适应，提示学生本页统一的作答方式。"""
    if kind == "mc":
        return "Read the passage. Circle the correct answer (A / B / C) for each question."
    if kind == "tf":
        return "Read the passage. Write T (true) or F (false) for each statement."
    return "Read the passage and answer the questions."


def _reading_ext_items(outline: BookOutline, data: dict, exclude: set[str],
                       max_n: int = 4) -> list[dict]:
    """L4 第②阅读页【综合理解延伸】题（判断 / 简答 / 填空，紧扣原文，绝不与第①页重复）。

    取材优先级：① AI 已抽的 tf/short 阅读题 → ② RR 阅读表达题(作简答) → ③ 原文原句生成
    True/False（交替"原句=True"与"反义句=False"）→ ④ 本文完形填空兜底。各项去重、去与第①页
    重复，最多 max_n 条；不足 3 条时调用方回退 SWBST 复述。"""
    import re as _re
    items: list[dict] = []
    seen: set[str] = {_reading_question_key(x) for x in (exclude or set()) if x}

    def _add(it: dict) -> None:
        text = (it.get("q") or "").strip()
        if _is_generic_reading_prompt(text):
            return
        key = _reading_question_key(text)
        if key and key not in seen and len(items) < max_n:
            seen.add(key)
            items.append(it)

    # ① AI 已抽 tf/short
    for q in (getattr(outline, "_reading_questions", []) or []):
        if q.get("kind") in ("tf", "short") and (q.get("q") or "").strip():
            _add({"kind": q["kind"], "q": q["q"].strip(), "page": q.get("page")})
    # ② RR 表达题 → 简答
    for q in (getattr(outline, "_rr_questions", []) or []):
        text = (q.get("q") or q.get("question") or "").strip()
        if text:
            _add({"kind": "short", "q": text, "page": q.get("page")})
    # ③ 原文原句 → True/False（交替真/反义）
    story = capitalize_names(_to_us_spelling(str(data.get("reading_text") or "")))
    sents = [s.strip() for s in _re.split(r"(?<=[.!?])\s+", story)
             if 20 <= len(s.strip()) <= 100]
    tf_count = 0
    for s in sents:
        if len(items) >= max_n:
            break
        core = s.rstrip(".!?").strip()
        if not core:
            continue
        if tf_count % 2 == 0:
            _add({"kind": "tf", "q": core + "."})
        else:
            neg = core[0].lower() + core[1:] if core else core
            _add({"kind": "tf", "q": f"It is not true that {neg}."})
        tf_count += 1
    # ④ 完形填空兜底（已规整的 fill_blanks）
    for f in (data.get("fill_blanks") or []):
        s = (f.get("sentence") or "").strip()
        if "____" in s:
            _add({"kind": "short", "q": s})
    return items[:max_n]


def _reading_question_key(text: str) -> str:
    text = _clean_text(text or "").lower()
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[\"'“”‘’.,!?;:()]+", "", text)
    return text.strip()


def _is_generic_reading_prompt(text: str) -> bool:
    low = (text or "").strip().lower()
    return (
        low.startswith("what happens here:")
        or low.startswith("what does the passage say in sentence")
        or low in {"question 1?", "question 2?", "question 3?", "question 4?"}
    )


def _fill_same_kind_reading_questions(
    questions: list[dict], kind: str, reading_text: str, *, max_n: int = 4
) -> list[dict]:
    """Top up the first Reading page so it does not look unfinished.

    Keep the page one-type-only (MC or T/F). If AI only provides 3 same-kind
    questions, derive one more simple item from the passage instead of leaving
    a sparse page.
    """
    import re as _re
    out: list[dict] = []
    seen: set[str] = set()
    for q in (questions or []):
        text = q.get("q") or ""
        if _is_generic_reading_prompt(text):
            continue
        key = _reading_question_key(text)
        if not key or key in seen:
            continue
        seen.add(key)
        out.append(q)
        if len(out) >= max_n:
            break
    if kind == "tf":
        mixed = _mixed_tf_items(reading_text, max_n=max_n)
        if len(mixed) >= max_n and not _has_false_tf_statement(out):
            return mixed[:max_n]
    sents = [
        s.strip().rstrip(".!?")
        for s in _re.split(r"(?<=[.!?])\s+", reading_text or "")
        if 18 <= len(s.strip()) <= 100
    ]
    for idx, sent in enumerate(sents):
        if len(out) >= max_n:
            break
        if not sent or sent.lower() in seen:
            continue
        if kind == "tf":
            q = sent + "."
            if idx % 2 == 1:
                q = _false_tf_variant(sent)
                if not q:
                    continue
            out.append({"kind": "tf", "q": q})
            seen.add(q.lower())
        elif kind == "mc":
            wrong = _false_tf_variant(sent).rstrip(".!?")
            if not wrong or wrong.lower() == sent.lower():
                continue
            out.append({
                "kind": "mc",
                "q": "Which sentence is in the story?",
                "options": [sent + ".", wrong + "."],
                "correct": 0,
            })
            seen.add("which sentence is in the story?")
    return out[:max_n]


def _has_false_tf_statement(questions: list[dict]) -> bool:
    false_markers = (" not ", " no ", " never ", " clean.", "one-day", "monday", "sad")
    for q in questions or []:
        text = f" {(q.get('q') or '').strip().lower()} "
        if q.get("answer") is False or q.get("correct") is False:
            return True
        if any(m in text for m in false_markers):
            return True
    return False


def _false_tf_variant(sentence: str) -> str:
    if re.search(r"\bsaid,\s*[\"“]", sentence, flags=re.I):
        out = re.sub(r"\bsaid,\s*", "did not say, ", sentence, count=1, flags=re.I).strip()
        if re.search(r"[.!?][\"”]$", out):
            return out
        return out.rstrip(".!?") + "."
    replacements = [
        (r"\bdid not know many\b", "knew many"),
        (r"\bdoes not know many\b", "knows many"),
        (r"\bfour seasons\b", "two seasons"),
        (r"\bflowers and new leaves\b", "snow and ice"),
        (r"\bcome out and play\b", "hide and sleep"),
        (r"\bhot days\b", "cold days"),
        (r"\bcold days\b", "hot days"),
        (r"\byellow leaves and fruits\b", "green leaves and flowers"),
        (r"\bwhite snow\b", "hot sun"),
        (r"\bwhite coat\b", "red coat"),
        (r"\bhospital\b", "school"),
        (r"\bkitchen\b", "park"),
        (r"\bhat and an apron\b", "helmet and boots"),
        (r"\bhelmet and boots\b", "hat and an apron"),
        (r"\bstrong boots\b", "soft slippers"),
        (r"\bon the farm\b", "in the hospital"),
        (r"\bin the office\b", "on the farm"),
        (r"\bwork safely\b", "work slowly"),
        (r"\bhide and sleep\b", "come out and play"),
        (r"\bswim and have fun\b", "hide and sleep"),
        (r"\bspring\b", "winter"),
        (r"\bsummer\b", "winter"),
        (r"\bautumn\b", "spring"),
        (r"\bwinter\b", "summer"),
        (r"\bmessy\b", "clean"),
        (r"\bclean\b", "messy"),
        (r"\bSunday\b", "Monday"),
        (r"\bseven-day\b", "one-day"),
        (r"\bhappy\b", "sad"),
        (r"\bproud\b", "worried"),
        (r"\bmany\b", "a few"),
        (r"\bwill\b", "will not"),
        (r"\bwork\b", "play"),
        (r"\bschool\b", "home"),
        (r"\bhome\b", "school"),
        (r"\bwear\b", "find"),
        (r"\bwears\b", "finds"),
        (r"\bvery different\b", "the same"),
        (r"\bdifferent\b", "the same"),
        (r"\bclean\b", "dirty"),
        (r"\bsafe\b", "unsafe"),
        (r"\btidy\b", "messy"),
        (r"\bthree or more\b", "only one"),
        (r"\bone\b", "two"),
        (r"\btwo\b", "three"),
        (r"\bthree\b", "one"),
    ]
    for pat, repl in replacements:
        if re.search(pat, sentence, flags=re.I):
            return re.sub(pat, repl, sentence, count=1, flags=re.I).rstrip(".!?") + "."
    return ""


def _split_story_sentences(text: str) -> list[str]:
    marked = re.sub(r"([.!?][\"”]?)\s+", r"\1|||", text or "")
    return [s.strip() for s in marked.split("|||") if s.strip()]


def _mixed_tf_items(reading_text: str, *, max_n: int = 4) -> list[dict]:
    story = _to_us_spelling(reading_text or "")
    low = story.lower()
    if all(k in low for k in ("mia", "homework", "seven-day", "piano")):
        return [
            {"kind": "tf", "q": "Mia has many things to do this week.", "answer": True},
            {"kind": "tf", "q": "Her room is clean.", "answer": False},
            {"kind": "tf", "q": "The piano show is on Sunday.", "answer": True},
            {"kind": "tf", "q": "Mia makes a one-day plan.", "answer": False},
        ][:max_n]
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 100 and not s.strip().endswith("?")
    ]
    if len(sents) > max_n and max_n > 1:
        idxs = sorted({round(i * (len(sents) - 1) / (max_n - 1)) for i in range(max_n)})
        sents = [sents[i] for i in idxs]
    if len(sents) > max_n:
        idxs = sorted({round(i * (len(sents) - 1) / (max_n - 1)) for i in range(max_n)})
        sents = [sents[i] for i in idxs]
    out: list[dict] = []
    seen: set[str] = set()
    for idx, sent in enumerate(sents):
        if len(out) >= max_n:
            break
        if idx % 2 == 0:
            q = sent + "."
            ans = True
        else:
            q = _false_tf_variant(sent)
            if not q:
                continue
            ans = False
        key = q.lower()
        if q and key not in seen:
            seen.add(key)
            out.append({"kind": "tf", "q": q, "answer": ans})
    return out[:max_n]


def _reading_cloze_items(reading_text: str, exclude: set[str] | None = None,
                         *, max_n: int = 5) -> list[dict]:
    """Create a second L3/L4 Reading page that differs from T/F.

    The item is a short written cloze based on actual story sentences, so the
    page stays text-evidence based without becoming another true/false page.
    """
    story = _to_us_spelling(reading_text or "")
    exclude_keys = {_reading_question_key(x) for x in (exclude or set()) if x}
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 100 and not s.strip().endswith("?")
    ]
    if len(sents) > max_n and max_n > 1:
        idxs = sorted({round(i * (len(sents) - 1) / (max_n - 1)) for i in range(max_n)})
        sents = [sents[i] for i in idxs]
    out: list[dict] = []
    seen: set[str] = set()
    for sent in sents:
        if len(out) >= max_n:
            break
        words = re.findall(r"[A-Za-z][A-Za-z'-]*", sent)
        if len(words) < 4:
            continue
        # Blank a short final chunk, but keep the sentence readable for A1/A1+.
        chunk_n = 2 if len(words[-1]) <= 4 and len(words) >= 5 else 1
        answer_words = words[-chunk_n:]
        answer = " ".join(answer_words)
        if len(answer) < 3:
            continue
        pattern = r"\b" + r"\s+".join(re.escape(w) for w in answer_words) + r"\b"
        prompt = re.sub(pattern, ANSWER_BLANK, sent, count=1)
        if prompt == sent:
            continue
        q = f"Complete the text summary: {prompt}"
        if not q.rstrip().endswith((".", "!", "?", "\"", "”")):
            q += "."
        key = _reading_question_key(q)
        if key in seen or key in exclude_keys:
            continue
        seen.add(key)
        out.append({"kind": "short", "q": q, "answer": answer})
    return out[:max_n]


def _reading_mc_sentence_items(reading_text: str, *, max_n: int = 5) -> list[dict]:
    """Create controlled MC reading items from actual story sentences."""
    story = _to_us_spelling(reading_text or "")
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 88 and not s.strip().endswith("?")
    ]
    out: list[dict] = []
    seen: set[str] = set()

    def _fact_choice_stem(sentence: str) -> str:
        """Make Page 5 stems specific enough that MC items do not look duplicated."""
        low = sentence.lower()
        subject_match = re.match(
            r"^((?:A|An|The|Some|Many|Different)\s+[A-Za-z]+|[A-Z][a-z]+(?:\s+and\s+[A-Z][a-z]+)?|[A-Za-z]+s)\b",
            sentence,
        )
        if subject_match:
            subject = subject_match.group(1).strip().lower()
            if subject not in {"this", "that", "these", "those", "their"}:
                return f"Choose the fact about {subject}."
        topic_rules = [
            (("wear", "uniform", "helmet", "boots", "suit", "clothes"), "clothes"),
            (("worker", "work", "job", "doctor", "chef", "farmer", "firefighter"), "work"),
            (("family", "families", "parents", "children"), "families"),
            (("school", "class", "student", "teacher"), "school"),
            (("home", "room", "house"), "home"),
            (("plan", "first", "next", "then", "finally", "every day"), "the plan"),
            (("animal", "animals", "live", "eat"), "animals"),
            (("place", "city", "country", "world"), "places"),
        ]
        for keys, topic in topic_rules:
            if any(k in low for k in keys):
                return f"Choose the fact about {topic}."
        m = re.match(r"^(A|An|The|Some|Many|Different)?\s*([A-Z][a-z]+|[a-z]+)", sentence)
        if m:
            raw = " ".join(x for x in m.groups() if x).strip().lower()
            if raw and raw not in {"a", "an", "the", "some", "many", "different"}:
                return f"Choose the fact about {raw}."
        return "Choose the correct fact."

    for sent in sents:
        if len(out) >= max_n:
            break
        false_one = _false_tf_variant(sent).rstrip(".!?")
        if not false_one or false_one.lower() == sent.lower():
            continue
        opts = [sent + ".", false_one + "."]
        key = sent.lower()
        if key in seen:
            continue
        seen.add(key)
        # Deterministic light shuffle based on the sentence.
        rnd = random.Random(hash(sent) & 0xFFFFFFFF)
        correct = opts[0]
        rnd.shuffle(opts)
        item = {
            "kind": "mc",
            "q": _fact_choice_stem(sent),
            "options": opts,
            "correct": opts.index(correct),
            "answer": correct,
            "evidence": sent + ".",
        }
        if _valid_mc_item(item, min_options=2):
            out.append(item)
    return out[:max_n]


def _reading_short_answer_items(reading_text: str, *, max_n: int = 5) -> list[dict]:
    """Simple text-evidence short-answer items for L3/L4 Reading page variety."""
    story = _to_us_spelling(reading_text or "")
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 92 and not s.strip().endswith("?")
    ]
    out: list[dict] = []
    seen: set[str] = set()

    def _push(q: str, answer: str) -> None:
        q = _clean_text(q)
        answer = _clean_text(answer)
        if not q or not answer:
            return
        if q.lower().startswith(("what ", "who ", "where ", "when ", "why ", "how ")) and "?" not in q:
            q += "?"
        key = _reading_question_key(q)
        if key in seen:
            return
        seen.add(key)
        out.append({"kind": "short", "q": q, "answer": answer, "evidence": answer})

    def _article_phrase(noun: str) -> str:
        noun = _clean_text(noun)
        if not noun:
            return noun
        return noun if re.match(r"^(a|an|the|some|many|different|clean)\b", noun, flags=re.I) else noun

    def _make_fact_question(sent: str) -> tuple[str, str] | None:
        text = _clean_text(sent).rstrip(".!?")
        if not text or text.endswith("?"):
            return None
        there_match = re.match(r"^there\s+(is|are)\s+(.+?)$", text, flags=re.I)
        if there_match:
            be, obj = there_match.groups()
            near_match = re.match(r"(.+?)\s+near\s+(.+?)(?:\s+with\s+.+)?$", obj, flags=re.I)
            if near_match:
                thing, place = near_match.groups()
                return f"What is near {place}?", _article_phrase(thing)
            return "What does the passage say there are?", _article_phrase(obj)
        subj_re = (
            r"(i|we|they|he|she|it|Mia and Tommy|Some ocean animals|Ocean animals|"
            r"Coral reefs|Deep-sea fish|[A-Z][a-z]+|(?:A|An|The|Some|Many|Different) [a-z]+(?: [a-z]+)?)"
        )
        plural_subjects = {
            "we", "they", "people", "children", "kids", "students", "families",
            "animals", "jobs", "workers", "some jobs", "ocean animals",
            "some ocean animals", "coral reefs", "deep-sea fish",
        }

        def _is_plural_subject(subj: str) -> bool:
            low = (subj or "").strip().lower()
            if low in plural_subjects:
                return True
            m = re.match(r"^(some|many|different)\s+([a-z]+)$", low)
            if m and (m.group(2).endswith("s") or m.group(2) in plural_subjects):
                return True
            return bool(low.endswith("s") and not low.endswith("ss"))

        def _student_subject(subj: str) -> str:
            low = (subj or "").strip().lower()
            if _is_plural_subject(subj) or re.match(r"^(a|an|the|some|many|different)\s+", subj or "", flags=re.I):
                return low
            return subj

        looks_like_match = re.match(rf"^{subj_re}\s+looks\s+like\s+(.+?)$", text, flags=re.I)
        if looks_like_match:
            subj, obj = looks_like_match.groups()
            return f"What does {_student_subject(subj)} look like?", obj
        wear_match = re.match(rf"^{subj_re}\s+(wear|wears)\s+(.+?)\s+for\s+(.+?)$", text, flags=re.I)
        if wear_match:
            subj, verb, obj, context = wear_match.groups()
            cue = obj.split(" and ")[0].split(",")[0].strip()
            aux = "do" if subj.lower() == "i" or _is_plural_subject(subj) else "does"
            subj_text = _student_subject(subj)
            return f"What {aux} {subj_text} wear for {context}? Clue: {cue}.", obj
        patterns = [
            (r"^(.+?)\s+(is|are|was|were)\s+(.+?)$", "What {be} {subj}?", "{obj}"),
            (rf"^{subj_re}\s+(wear|wears|use|uses|need|needs|have|has|make|makes|choose|chooses|keep|keeps|help|helps|bring|brings|live|lives|go|goes|see|sees|like|likes)\s+(.+?)$", "What does {subj} {verb_base}?", "{obj}"),
            (r"^((?:some|many|different)\s+[a-z]+|families|people|children|kids|students)\s+(.+?)$", "What do {subj} do?", "{rest}"),
            (r"^(people|children|kids|students|families|animals|jobs|some jobs)\s+(.+?)$", "What do {subj} do?", "{rest}"),
            (r"^(.+?)\s+can\s+(.+?)$", "What can {subj} do?", "{obj}"),
            (r"^(.+?)\s+will\s+(.+?)$", "What will {subj} do?", "{obj}"),
        ]
        for pat, qtpl, atpl in patterns:
            m = re.match(pat, text, flags=re.I)
            if not m:
                continue
            groups = m.groups()
            if len(groups) == 2:
                subj, obj = groups
                subj = _student_subject(subj)
                data = {"subj": subj, "obj": obj, "rest": obj, "be": ""}
            elif len(groups) == 3:
                subj, verb, obj = groups
                qtpl_use = qtpl
                atpl_use = atpl
                subj = _student_subject(subj)
                verb_base = {
                    "wears": "wear", "uses": "use", "needs": "need", "has": "have",
                    "makes": "make", "chooses": "choose", "keeps": "keep", "helps": "help",
                    "brings": "bring", "lives": "live", "goes": "go", "sees": "see",
                    "likes": "like",
                }.get(verb.lower(), verb.lower())
                if verb.lower() not in {"is", "are", "was", "were"} and _is_plural_subject(subj):
                    qtpl_use = "What do {subj} do?"
                    atpl_use = "{rest}"
                data = {"subj": subj, "verb": verb, "verb_base": verb_base, "obj": obj,
                        "rest": f"{verb} {obj}", "be": verb}
            else:
                continue
            q = (qtpl_use if len(groups) == 3 else qtpl).format(**data)
            a = (atpl_use if len(groups) == 3 else atpl).format(**data)
            if len(a.split()) <= 12:
                return q, _article_phrase(a)
        return None

    for sent in sents:
        if len(out) >= max_n:
            break
        made = _make_fact_question(sent)
        if made:
            _push(made[0], made[1])
            continue
        if sent.strip().endswith("?"):
            continue
        if len(out) < max_n:
            _push("Copy one fact from the passage.", sent)
    return out[:max_n]


def _reading_literal_mc_items(reading_text: str, *, max_n: int = 5) -> list[dict]:
    """Page 5 literal MC: direct fact questions with one answer option."""
    facts = _reading_short_answer_items(reading_text, max_n=max(6, max_n + 2))
    answers: list[str] = []
    for item in facts:
        ans = _clean_text(item.get("answer", ""))
        if ans and ans.lower() not in {a.lower() for a in answers} and len(ans.split()) <= 8:
            answers.append(ans)
    out: list[dict] = []
    for item in facts:
        q = _clean_text(item.get("q", ""))
        answer = _clean_text(item.get("answer", ""))
        if not q or not answer or answer not in answers:
            continue
        distractors = [a for a in answers if a.lower() != answer.lower()]
        if len(distractors) < 2:
            continue
        rnd = random.Random(hash(q + "|" + answer) & 0xFFFFFFFF)
        rnd.shuffle(distractors)
        options = [answer] + distractors[:2]
        rnd.shuffle(options)
        mc = {
            "kind": "mc",
            "q": q,
            "options": options,
            "correct": options.index(answer),
            "answer": answer,
            "evidence": answer,
        }
        if _valid_mc_item(mc, min_options=3):
            out.append(mc)
        if len(out) >= max_n:
            break
    return out


def _reading_text_correction_items(reading_text: str, *, max_n: int = 4) -> list[dict]:
    """L4 A2 reading: correct one wrong detail from the passage."""
    story = _to_us_spelling(reading_text or "")
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 88
    ]
    out: list[dict] = []
    seen: set[str] = set()
    for sent in sents:
        wrong = _false_tf_variant(sent).rstrip(".!?")
        if not wrong or wrong.lower() == sent.lower():
            continue
        key = wrong.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append({
            "kind": "short",
            "q": f"Correct the sentence: {wrong}.",
            "answer": sent,
            "wrong_sentence": wrong + ".",
            "correct_sentence": sent + ".",
            "evidence": sent + ".",
        })
        if len(out) >= max_n:
            break
    return out


def _reading_cause_effect_items(reading_text: str, *, max_n: int = 4) -> list[dict]:
    """L4 A2 reading: simple cause/effect or detail connection prompts."""
    story = _to_us_spelling(reading_text or "")
    sents = [
        capitalize_names(s.strip().rstrip(".!?"))
        for s in _split_story_sentences(story)
        if 18 <= len(s.strip()) <= 92
    ]
    out: list[dict] = []
    for i in range(min(len(sents) - 1, max_n)):
        out.append({
            "kind": "short",
            "q": f"What happens after this? {sents[i]}.",
            "answer": sents[i + 1],
            "cause": sents[i] + ".",
            "effect": sents[i + 1] + ".",
        })
    return out


def _reading_sequence_events(reading_text: str, *, max_n: int = 5) -> list[str]:
    """Short story events for a written sequence activity."""
    story = _to_us_spelling(reading_text or "")
    events: list[str] = []
    seen: set[str] = set()
    for sent in _split_story_sentences(story):
        text = capitalize_names(sent.strip().rstrip(".!?"))
        if not text or len(text) < 18 or len(text) > 88:
            continue
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        events.append(text + ".")
        if len(events) >= max_n:
            break
    return events


def _build_l34_sequence_reading_page(slide, brand_rgb: tuple, reading_text: str,
                                     events: list[str], *, start_no: int = 1) -> None:
    """Reading page variant: number story events in order."""
    _add_title(slide, "Reading", "Number the events in order.")
    text = _clean_text(reading_text)
    read_pt = 11.5 if len(text) > 520 else 12.5
    text_top = CONTENT_Y + 1.35
    text_h = 1.55 if len(text) > 520 else 1.85
    text_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.40), Inches(text_top),
        Inches(CONTENT_W - 0.80), Inches(text_h),
    )
    text_box.adjustments[0] = 0.03
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = WHITE
    text_box.line.color.rgb = RGBColor(*brand_rgb)
    text_box.line.width = Pt(1.5)
    text_box.shadow.inherit = False
    tf = text_box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(read_pt)
    r.font.color.rgb = BLACK

    items = list(events[:5])
    rnd = random.Random(hash("|".join(items)) & 0xFFFFFFFF)
    scrambled = items[:]
    for _ in range(8):
        rnd.shuffle(scrambled)
        if scrambled != items:
            break
    body = [f"{start_no + i}.  ____  {event}" for i, event in enumerate(scrambled)]
    _flow_lines(
        slide, body, text_top + text_h + 0.22,
        CONTENT_Y + CONTENT_H - 0.30 - (text_top + text_h + 0.22),
        CONTENT_X + 0.78, CONTENT_W - 1.56, 15,
        gap_min=0.14, gap_max=0.48,
    )


def _build_l4_vocab2(new_page, brand_rgb: tuple, data: dict,
                     images: Optional[list[Path]], seed: int) -> None:
    """L4 词汇②页：在 4 种词汇题型间【按书轮换】（每种版式自身恒定、内容紧扣本书词汇），
    解决"几乎每本都用同一道原文挖空/占位填空"的单调问题：
      0 → 原文/词义挖空      (fill in the blank, _build_p2_fill)
      1 → 看义选词           (choose the correct word, 四选一)
      2 → 看义写词           (write the word, 首字母提示)
      3 → 选词补全句         (use the word in context, 句子四选一)
    每种先校验可用题量阈值，不足则顺延到必定 ≥2 题的挖空页，杜绝单题/空页（Problem 1 兜底）。"""
    pairs = data.get("match_pairs") or []
    fills = data.get("fill_blanks") or []
    bank = data.get("word_bank") or []

    riddles = _riddle_mc_items(pairs, max_n=4)            # 看义选词
    meaning = _word_fill_meaning_items(pairs, max_n=6)    # 看义写词
    cloze_mc = _cloze_mc_items(fills, bank, max_n=4)      # 选词补全句

    def _do_riddle() -> bool:
        if len(riddles) >= 3:
            _build_mcq_page(new_page(), brand_rgb, "Vocabulary",
                            "Read each clue and circle the correct word.", riddles,
                            answer_paren=False)
            return True
        return False

    def _do_meaning() -> bool:
        if len(meaning) >= 2:
            _build_word_fill_page(new_page(), brand_rgb, meaning, "Vocabulary",
                                  "Read the meaning and write the word. The first letter is given.")
            return True
        return False

    def _do_cloze_mc() -> bool:
        if len(cloze_mc) >= 3:
            _build_mcq_page(new_page(), brand_rgb, "Vocabulary",
                            "Choose the correct word to complete each sentence.", cloze_mc)
            return True
        return False

    choice = seed % 4
    if choice == 1 and _do_riddle():
        return
    if choice == 2 and _do_meaning():
        return
    if choice == 3 and _do_cloze_mc():
        return
    # choice == 0，或所选题型可用题量不足 → 回退链（仍尽量换一种花样），最后落到必 ≥2 的挖空页
    if choice != 0 and (_do_meaning() or _do_riddle()):
        return
    _build_p2_fill(new_page(), brand_rgb, fills, bank, images)


def build_worksheet(
    outline: BookOutline,
    out_path: Path,
    *,
    image_paths: Optional[Iterable[Path]] = None,
    sentence_image_mode: str = "reuse",  # reuse=复用绘本图 / none=不配图
    second_reading_mode: str = "auto",   # auto/reading/mindmap/writing/writing_official/pbl/color_say
    coloring_image: Optional[Path] = None,  # L0-2 涂色线稿图（batch 生成；缺省则用留白画框）
) -> Path:
    """生成 worksheet（所有级别统一 8 页）：

    2 词汇(Vocabulary) + 2 句型(Sentence) + 2 阅读(Reading) + 1 思维导图(Graphic Organizer) + 1 写作(Writing)。
      阅读① = 原文 + 理解题（垂直列表）。
      阅读② = 阅读理解延伸题。
      Graphic Organizer 内容由 second_reading_mode 决定：
        auto    → 按级别/文本自动选择思维导图
        reading → 阅读理解延伸（题目分两页）
        mindmap → 思维导图(SWBST)
        pbl     → 读后 PBL 迷你项目

    阅读理解题优先取 outline._reading_questions（AI 专门抽的 mc/tf/short），
    无则兜底旧 reading_mcs。Sentence 页复用 image_paths[2..5]（故事 P2..P5）。
    """
    data = _resolve_worksheet_data(outline)
    brand_rgb = brand_color_rgb(outline.level)
    level_label = _level_label(outline.level)
    footer_text = f"{level_label} - {outline.title}"
    # v2.1: 用专门的 Dino 头 icon（不是设定卡 dino_logo.png）
    logo_icon = BRAND_DIR / "dino_head_icon.png"
    if not logo_icon.exists():
        logo_icon = BRAND_DIR / "dino_logo.png"  # 兜底
    images = list(image_paths or [])
    lvl_n = _level_num(outline.level)
    source_validation = _worksheet_source_preflight(outline, data, lvl_n)
    source_validation = _relax_nonblocking_worksheet_issues(source_validation)
    if source_validation.get("blockers"):
        msg = "; ".join(i.get("message", "") for i in source_validation["blockers"])
        raise ValueError(f"Worksheet source preflight failed: {msg}")

    # v2.2：优先以官方 A4 模板为底版（按级别克隆 slide → 100% 还原边框/Logo/Name/配色/页脚）；
    # 模板缺失时退回旧的自绘外框。
    use_template = _WS_TEMPLATE.exists()
    if use_template:
        prs = Presentation(str(_WS_TEMPLATE))
        n_template = len(prs.slides._sldIdLst)
        tpl_src = prs.slides[_template_slide_index(outline.level)]
        blank = None
    else:
        prs = Presentation()
        prs.slide_width = _RawInches(SLIDE_W)
        prs.slide_height = _RawInches(SLIDE_H)
        n_template = 0
        tpl_src = None
        blank = prs.slide_layouts[6]

    def new_page():
        if use_template:
            s = _clone_template_slide(prs, tpl_src)
            _set_footer(s, footer_text)
            _fix_name_badge(s)
        else:
            s = prs.slides.add_slide(blank)
            _draw_brand_frame(s, brand_rgb, footer_text, logo_icon)
        return s

    band = _lvl_band(lvl_n)

    # ===== 2 词汇页（分级选型；学什么考什么）=====
    # 词汇页①：L3/L4 按 SOP 固定为理解页，优先图文/词义匹配，不再与 P2 连续做拼写题。
    _build_p1_match(new_page(), brand_rgb, data["match_pairs"], images)
    _record_worksheet_items(outline, 1, data["match_pairs"])
    if lvl_n in (3, 4):
        sld_lst = prs.slides._sldIdLst
        old_slide = sld_lst[-1]
        old_r_id = getattr(old_slide, "rId", None)
        sld_lst.remove(old_slide)
        if old_r_id:
            with contextlib.suppress(Exception):
                prs.part.drop_rel(old_r_id)
        _build_l34_vocab1(new_page, brand_rgb, data, images, _ws_seed(outline), outline)
    elif False:
        valid_defs = len([p for p in data["match_pairs"] if _valid_definition_pair(p)])
        img_count = len([p for p in (images or []) if p and Path(p).exists()])
        page1_code = (
            "vocab_word_definition_matching"
            if valid_defs >= _l34_min_items("vocab_word_definition_matching", 4)
            else "vocab_word_picture_matching"
            if img_count >= _l34_min_items("vocab_word_picture_matching", 4)
            else "vocab_word_definition_matching"
        )
        _record_l34_activity(outline, 1, page1_code)
    # 词汇页② 分级：L0-2 看义/首字母写词；L3 谜语四选一；L4 原文挖空；L5-6 构词补全
    if band == "l02":
        wf = _word_fill_meaning_items(data["match_pairs"], max_n=6)
        if wf:
            _build_word_fill_page(new_page(), brand_rgb, wf, "Vocabulary",
                                  "Read the meaning and write the word. The first letter is given.")
        else:
            _build_p2_fill(new_page(), brand_rgb, data["fill_blanks"], data["word_bank"], images)
    elif band == "l56":
        wf = _word_fill_morph_items(data["match_pairs"], max_n=6)
        if wf:
            _build_word_fill_page(new_page(), brand_rgb, wf, "Vocabulary",
                                  "Complete each word using the correct ending.")
        else:
            _build_p2_fill(new_page(), brand_rgb, data["fill_blanks"], data["word_bank"], images)
    else:  # L3/L4 词汇②：跨书轮换题型，避免每本都只做语境填空。
        _build_l34_vocab2(new_page, brand_rgb, data, images, _ws_seed(outline), lvl_n,
                          outline=outline)

    # ===== 2 句型页（考点 = 本课语法焦点；学什么考什么）=====
    # 时态自适应（用户拍板 2026-06-06）：现在时为主的文本（如非虚构科普）→ 现在时考点
    #   （主谓一致 / 现在时填空）；过去式为主 → 维持过去式考点。避免"现在时文出过去式题"
    #   既考点错、又因可转动词太少导致一页一题。每页保底 3-4 题填满版面。
    tense = _dominant_tense(outline)
    # Bug1 兜底：句型引擎可能因故事太短/无可转动词而 0 产出 → 用官方 S&S 句型/例句兜底，杜绝空页。
    sent_fb = _sentence_pattern_fallback_items(outline, data.get("word_bank"))

    # —— 句型页① ——
    # L3-L6: the syllabus Sentence Pattern / Example Sentence is authoritative.
    # Only fall back to tense heuristics when the syllabus has no usable frame.
    official_sentences_done = False
    if lvl_n >= 3:
        frame_fills, frame_bank = _official_sentence_frame_fill_items(outline, max_n=4)
        if frame_fills:
            frame_mcs = _official_sentence_frame_mc_items(outline, frame_fills, frame_bank, max_n=4)
            if frame_mcs:
                _record_l34_activity(outline, 3, "sentence_target_pattern_choice")
                _record_worksheet_items(outline, 3, frame_mcs)
                _build_p3_sentence(
                    new_page(), brand_rgb, frame_mcs, images,
                    show_images=(sentence_image_mode != "none"),
                    subtitle=f"{_ws_activity_instruction('sentence_target_pattern_choice')} Example: {_display_sentence_frame(outline)}",
                    fallback=sent_fb,
                )
            else:
                _record_l34_activity(outline, 3, "sentence_complete_frame")
                _record_worksheet_items(outline, 3, frame_fills)
                _build_p2_fill(
                    new_page(), brand_rgb, frame_fills, frame_bank, images,
                    title="Sentences",
                    subtitle=f"Use the example: {_display_sentence_frame(outline)}",
                    fallback=sent_fb,
                )
            official_sentences_done = True
    if not official_sentences_done:
        if tense == "present":
            sent_mcs = _present_agreement_mc_items(outline, max_n=4)
            if len(sent_mcs) < 3:
                sent_mcs = (sent_mcs + _tense_mc_items(outline, max_n=4))[:4]
            mc_sub = "Tick (\u2713) the sentence that is correct."
        else:
            sent_mcs = _tense_mc_items(outline, max_n=4)
            mc_sub = "Tick (\u2713) the sentence that uses the correct past tense."
        if len(sent_mcs) < 3 and data.get("sentence_mcs"):
            sent_mcs = (sent_mcs + data["sentence_mcs"])[:4]
            mc_sub = "Tick (\u2713) the correct sentence."
        if lvl_n in (3, 4):
            _record_l34_activity(outline, 3, "sentence_correct_sentence_choice")
        _build_p3_sentence(new_page(), brand_rgb, sent_mcs, images,
                           show_images=(sentence_image_mode != "none"), subtitle=mc_sub,
                           fallback=sent_fb)

    # —— 句型页② 分级 + 时态自适应（保底 ≥3 题）——
    def _present_fill_page():
        sf, _bk = _present_fill_items(outline, max_n=4)
        if len(sf) < 3:
            xf, _ = _tense_fill_items(outline, max_n=4)
            sf = (sf + xf)[:4]
        _build_p2_fill(new_page(), brand_rgb, sf, [], images, title="Sentences",
                       subtitle="Write the correct form of the verb to complete each sentence.",
                       fallback=sent_fb)

    if official_sentences_done:
        _build_l34_sentence2(new_page, brand_rgb, outline, frame_fills, frame_bank,
                             sent_fb, _ws_seed(outline), lvl_n)
    elif lvl_n == 3:
        frame_copy = _sentence_frame_copy_items(outline, max_n=4)
        _build_prompt_line_page(
            new_page(), brand_rgb, frame_copy, "Sentences",
            "Write about your own plan. Use your own ideas.",
            prompt_key="prompt",
            fallback=sent_fb,
        )
    elif band == "l02":
        unsc = _unscramble_items(_story_sentences_for_grammar(outline), max_n=6)
        if len(unsc) >= 3:
            _build_prompt_line_page(new_page(), brand_rgb, unsc, "Sentences",
                                    "Put the words in order to make a sentence.",
                                    prompt_key="scrambled", fallback=sent_fb)
        elif tense == "present":
            _present_fill_page()
        else:
            sf, _ = _tense_fill_items(outline, max_n=4)
            _build_p2_fill(new_page(), brand_rgb, sf, [], images, title="Sentences",
                           subtitle="Write the correct verb to complete each sentence.",
                           fallback=sent_fb)
    elif tense == "present":
        # 现在时为主：l34/l56 第 2 句型页统一用现在时填空（学什么考什么）
        _present_fill_page()
    elif band == "l56":
        rw = _rewrite_items(outline, max_n=6)
        if len(rw) >= 3:
            _build_prompt_line_page(
                new_page(), brand_rgb, rw, "Sentences",
                "Rewrite each sentence in the past tense.",
                prompt_key="prompt",
                example="I walk to school.  →  I walked to school.",
                fallback=sent_fb,
            )
        else:
            sf, _ = _tense_fill_items(outline, max_n=4)
            _build_p2_fill(new_page(), brand_rgb, sf, [], images, title="Sentences",
                           subtitle="Write the correct past tense of each verb in brackets.",
                           fallback=sent_fb)
    else:  # l34（过去式为主）→ 写动词过去式
        sent_fills, sent_bank = _tense_fill_items(outline, max_n=4)
        if not sent_fills:
            sent_fills, sent_bank = _sentence_fill_items(outline, max_n=4)
        _build_p2_fill(
            new_page(), brand_rgb, sent_fills, [], images,
            title="Sentences",
            subtitle="Write the correct past tense of each verb in brackets.",
            fallback=sent_fb,
        )

    # ===== 阅读 / 思维导图 / 写作 =====
    # 阅读理解题优先用专用 _reading_questions（mc/tf/short），否则兜底旧 reading_mcs。
    reading_text = capitalize_names(data["reading_text"])
    rq = list(getattr(outline, "_reading_questions", []) or [])
    if not rq:
        rq = [
            {"kind": "mc", "q": m.get("q", ""), "options": m.get("options", []),
             "correct": m.get("correct", 0)}
            for m in (data.get("reading_mcs") or []) if m.get("q")
        ]
    # 保底填满（用户拍板 2026-06-06）：worksheet 阅读题不足 3 题时，用 RR 阅读表达题
    # （AI 已生成、内容正确、紧扣原文）补齐到 4 题，杜绝阅读页"只有两题、下半页空白"。
    _valid = [q for q in rq if q.get("q")]
    if len(_valid) < 4:
        seen_q = {(q.get("q") or "").strip().lower() for q in _valid}
        for rrq in (getattr(outline, "_rr_questions", []) or []):
            if len(_valid) >= 4:
                break
            text = (rrq.get("q") or rrq.get("question") or "").strip()
            if not text or text.lower() in seen_q:
                continue
            seen_q.add(text.lower())
            item = {"kind": "short", "q": text,
                    "answer": rrq.get("answer", ""), "page": rrq.get("page")}
            rq.append(item)
            _valid.append(item)

    # 所有级别统一 8 页：2 词汇 + 2 句型 + 2 阅读 + 1 Graphic Organizer + 1 Writing。
    # 第 7 页 Graphic Organizer 按 second_reading_mode（auto=按级别/文本）选择：
    #   mindmap 思维导图(SWBST) / timeline 顺序图 / bubble 信息图 / pbl 迷你项目
    # 同页题型统一：整套阅读题先归一为"全选择"或"全判断"，两页同质、不混排。
    rq_uni, rq_kind = _unify_reading_questions(rq)
    sub_uni = _reading_subtitle(rq_kind)

    # 用户拍板（2026-06-09 修订）：【所有 L3-6 reading 页都放绘本原文】（彩框原文在上 + 题目在下），
    # 去掉 (P#) 回看提示（原文已在页面上）。L0-2 仍不配原文（靠常规排）。
    show_passage = lvl_n >= 3

    mode = _resolve_second_reading_mode(second_reading_mode, lvl_n)
    if mode == "auto_l3":
        # L3 第②阅读页（用户拍板 2026-06-09）：A1+ 不做长篇自由写作（太难）。按体裁分流增加多样性：
        #   · 非虚构(科普) → 代码原生【气泡思维导图】(中心主题 + 4 个关键点气泡)，更贴合"我学到了什么"；
        #   · 虚构(故事)   → 【看图回想 + 补全句子】引导式小结(fill-in + 词库)。
        mode = _auto_second_reading_mode(outline, lvl_n)

    # L4 第②阅读页跨书轮换（Problem 3：减少"每本都 SWBST 复述"）：仅当调用方未显式指定模式
    # （auto，batch 默认）时生效，按书在 4 种读后形式间轮换：
    #   reading2  综合理解题（判断 / 简答 / 填空，紧扣原文）
    #   mindmap   SWBST 五步复述（保留 ~1/4）
    #   l3bubble  关键信息气泡图（"我学到了什么"，非复述）
    #   l3summary 看故事补全句（fill-in 式小结）
    if (second_reading_mode or "auto").strip().lower() == "auto" and lvl_n == 4:
        mode = _auto_second_reading_mode(outline, lvl_n)
    go_mode = mode
    if go_mode in {"reading", "reading2", "writing", "writing_official"}:
        go_mode = _auto_second_reading_mode(outline, lvl_n)
    if go_mode in {"reading", "reading2", "writing", "writing_official"}:
        go_mode = "mindmap"
    setattr(outline, "_worksheet_second_reading_mode", go_mode)
    mode = "reading"

    if mode == "reading":
        if lvl_n in (3, 4):
            page1_q = []
            page1_code = ""
            try:
                book_no_seed = int("".join(ch for ch in str(getattr(outline, "book_number", "") or "") if ch.isdigit()) or "0")
            except Exception:
                book_no_seed = 0
            for code in _l34_activity_order(outline, 5, seed=(_ws_seed(outline) // 11) + book_no_seed):
                if code == "reading_supported_sentence_choice":
                    candidate = _reading_mc_sentence_items(reading_text, max_n=5)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="mc", max_n=5)
                elif code == "reading_true_false_literal":
                    candidate = _mixed_tf_items(reading_text, max_n=5)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="tf", max_n=5)
                    if len(candidate) < 4:
                        candidate = _fill_same_kind_reading_questions(rq_uni[:5], "tf", reading_text, max_n=5)
                        candidate = _sanitize_reading_items(candidate, preferred_kind="tf", max_n=5)
                elif code == "reading_literal_mc":
                    candidate = _reading_literal_mc_items(reading_text, max_n=5)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="mc", max_n=5)
                else:
                    continue
                if len(candidate) >= _l34_min_items(code, 4):
                    page1_q = candidate
                    page1_code = code
                    break
            if not page1_q:
                page1_q = _mixed_tf_items(reading_text, max_n=5)
                page1_q = _sanitize_reading_items(page1_q, preferred_kind="tf", max_n=5)
                page1_code = "reading_true_false_literal"
            _record_l34_activity(outline, 5, page1_code)
            _record_worksheet_items(outline, 5, page1_q)
            page1_subtitle = _ws_activity_instruction(page1_code)
            used_q = {_reading_question_key(q.get("q") or "") for q in page1_q}
            is_nonfic_reading = "non" in (getattr(outline, "fiction_type", "") or "").lower()
            sequence_events = _reading_sequence_events(reading_text, max_n=5)
            use_sequence_page = False
            page2_q = []
            page2_subtitle = ""
            page2_code = ""
            # Page 6 is integration, so use pedagogy-first bank priority rather
            # than pure rotation: sequence/structure before summary fallback.
            page6_signal = " ".join(
                str(getattr(outline, name, "") or "")
                for name in ("reading_skill", "reading_strategy", "graphic_organizer", "graphic_organizer_desc")
            ).lower()
            page6_seed = 0 if any(k in page6_signal for k in (
                "sequence", "sequencing", "process", "steps", "timeline", "order"
            )) else ((_ws_seed(outline) // 17) + book_no_seed)
            for code in _l34_activity_order(outline, 6, seed=page6_seed):
                if code == "reading_sequence_ordering":
                    if not is_nonfic_reading and len(sequence_events) >= _l34_min_items(code, 4):
                        use_sequence_page = True
                        page2_code = code
                        page2_q = [
                            {"event": event, "answer_order": idx + 1, "is_open": False}
                            for idx, event in enumerate(sequence_events)
                        ]
                        break
                    continue
                if code == "reading_text_based_correction":
                    candidate = _reading_text_correction_items(reading_text, max_n=4)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="short", max_n=4)
                elif code == "reading_cause_effect":
                    candidate = _reading_cause_effect_items(reading_text, max_n=4)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="short", max_n=4)
                elif code == "reading_short_answer":
                    candidate = _reading_short_answer_items(reading_text, max_n=5)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="short", max_n=5)
                elif code == "reading_summary_cloze":
                    candidate = _reading_cloze_items(reading_text, used_q, max_n=4)
                    candidate = _sanitize_reading_items(candidate, preferred_kind="short", max_n=4)
                else:
                    continue
                if len(candidate) >= _l34_min_items(code, 4):
                    page2_q = candidate
                    page2_code = code
                    break
            if not page2_code:
                page2_q = _reading_short_answer_items(reading_text, max_n=5)
                page2_q = _sanitize_reading_items(page2_q, preferred_kind="short", max_n=5)
                page2_code = "reading_short_answer" if len(page2_q) >= 4 else "reading_summary_cloze"
                if len(page2_q) < 4:
                    page2_q = _reading_cloze_items(reading_text, used_q, max_n=4)
                    page2_q = _sanitize_reading_items(page2_q, preferred_kind="short", max_n=4)
            _record_l34_activity(outline, 6, page2_code)
            _record_worksheet_items(outline, 6, page2_q)
            page2_subtitle = "" if use_sequence_page else _ws_activity_instruction(page2_code)
            _build_reading_page(
                new_page(), brand_rgb, reading_text, page1_q,
                subtitle=page1_subtitle,
                start_no=1, show_passage=show_passage,
                force_single_col=True,
            )
            if use_sequence_page:
                _build_l34_sequence_reading_page(
                    new_page(), brand_rgb, reading_text, sequence_events,
                    start_no=1,
                )
            else:
                _build_reading_page(
                    new_page(), brand_rgb, reading_text, page2_q,
                    subtitle=page2_subtitle,
                    start_no=1, show_passage=show_passage,
                    force_single_col=True,
                )
        else:
            ordered = rq_uni
            first_kind = rq_kind if rq_kind in ("mc", "tf") else "tf"
            page1_q = ordered[:4]
            if len(page1_q) < 4 and first_kind in ("mc", "tf"):
                page1_q = _fill_same_kind_reading_questions(page1_q, first_kind, reading_text, max_n=4)
                sub_uni = _reading_subtitle(first_kind)
            used_q = {_reading_question_key(q.get("q") or "") for q in page1_q}
            page2_q = [
                q for q in ordered[4:]
                if _reading_question_key(q.get("q") or "") not in used_q
            ][:4]
            _build_reading_page(new_page(), brand_rgb, reading_text, page1_q,
                                subtitle=sub_uni, start_no=1, show_passage=show_passage)
            if len(page2_q) >= 3:
                _build_reading_page(new_page(), brand_rgb, reading_text, page2_q,
                                    subtitle=sub_uni, start_no=1, show_passage=show_passage)
            else:
                ext = _reading_ext_items(outline, data, used_q, max_n=4)
                if len(ext) < 3:
                    ext = _mixed_tf_items(reading_text, max_n=4)
                _build_reading_fill_page(
                    new_page(), brand_rgb, "Reading",
                    "Read the story again and answer the questions.",
                    ext, start_no=1)
    else:
        # 横版阅读页（用户拍板 2026-06-06）：①同页题型必须统一（全选择 or 全判断，绝不混排）；
        # ②锁定 4 题（2×2 行对齐：左右各 2，横向两题对齐、纵向两题对齐）。先用归一后的同质题；
        # 不足 4 时只补【同种】，杜绝混入异类。若归一是 mixed（任何单一题型都 <3），则改取池中
        # 数量最多的单一题型，保证整页统一。
        q_cap = 4
        if rq_kind in ("mc", "tf"):
            first_q = list(rq_uni[:q_cap])
            if len(first_q) < 4:
                for q in rq:
                    if len(first_q) >= 4:
                        break
                    if q.get("q") and q.get("kind") == rq_kind and q not in first_q:
                        first_q.append(q)
            first_q = _fill_same_kind_reading_questions(first_q, rq_kind, reading_text, max_n=q_cap)
        else:
            by_kind: dict[str, list[dict]] = {}
            for q in rq:
                if q.get("q"):
                    by_kind.setdefault(q.get("kind", "mc"), []).append(q)
            best_k = max(("mc", "tf", "short"), key=lambda k: len(by_kind.get(k, [])))
            first_q = by_kind.get(best_k, [])[:q_cap]
            sub_uni = _reading_subtitle(best_k if best_k in ("mc", "tf") else "mixed")
        _build_reading_page(new_page(), brand_rgb, reading_text, first_q,
                            subtitle=sub_uni, start_no=1, show_passage=show_passage)
        if mode == "writing":
            _build_p5_writing(new_page(), brand_rgb, data["writing"], title="Reading")
        elif mode == "writing_official":
            # 官方 L4-13 风格：页面大标题 = Writing（写作脚手架 + 书写区）
            _build_p5_writing(new_page(), brand_rgb, data["writing"], title="Writing")
        elif mode == "pbl":
            _build_pbl_page(new_page(), brand_rgb, data, outline, title="Reading")
        elif mode == "color_say":
            _build_color_say_page(new_page(), brand_rgb, data, outline,
                                  coloring_image=coloring_image, title="Reading")
        elif mode == "l3summary":
            _build_l3_summary_page(new_page(), brand_rgb, data, outline, title="Reading")
        elif mode == "l3bubble":
            _build_l3_bubble_map(new_page(), brand_rgb, data, outline, title="Reading")
        elif mode == "timeline":
            _build_timeline_page(new_page(), brand_rgb, data, outline, title="Reading")
        elif mode == "planchart":
            _build_plan_chart_page(new_page(), brand_rgb, data, outline, title="Reading")
        elif mode == "reading2":
            # 综合理解延伸页（判断 / 简答 / 填空）；题量不足时回退 SWBST 复述，杜绝空页。
            _used = {(q.get("q") or "").strip().lower() for q in first_q}
            ext = _reading_ext_items(outline, data, _used, max_n=4)
            if len(ext) >= 3:
                _build_reading_fill_page(
                    new_page(), brand_rgb, "Reading",
                    "Read the story again and answer the questions.",
                    ext, start_no=1)
            else:
                _build_p6_mindmap(new_page(), data["mind_map"])
        else:  # mindmap
            _build_p6_mindmap(new_page(), data["mind_map"])

    # Page 7: Graphic Organizer. Page 8: Writing.
    if lvl_n in (3, 4):
        go_code = _l34_go_activity_code(outline, go_mode)
        _record_l34_activity(outline, 7, go_code)
        _record_l34_activity(outline, 8, "writing_organizer_to_writing")
        _record_worksheet_items(outline, 7, [{
            "graphic_organizer_type": _worksheet_go_type(outline, go_mode),
            "activity_code": go_code,
            "labels": [],
            "word_bank": [str(w).strip() for w in (outline.vocabulary_for_display or data.get("word_bank") or []) if str(w).strip()][:6],
            "fillable_fields": 3,
        }])
        _record_worksheet_items(outline, 8, [{
            "organizer_reference": _worksheet_go_type(outline, go_mode),
            "sentence_starters": True,
            "writing_lines": 4 if lvl_n <= 3 else 5,
            "target_pattern": _sentence_frame_text(outline),
        }])
        _build_l34_graphic_organizer_page(new_page(), brand_rgb, data, outline, go_mode)
        _build_l34_writing_page(new_page(), brand_rgb, outline)
    else:
        if go_mode == "pbl":
            _build_pbl_page(new_page(), brand_rgb, data, outline, title="Graphic Organizer")
        elif go_mode == "color_say":
            _build_color_say_page(new_page(), brand_rgb, data, outline,
                                  coloring_image=coloring_image, title="Graphic Organizer")
        elif go_mode == "l3summary":
            _build_l3_summary_page(new_page(), brand_rgb, data, outline, title="Graphic Organizer")
        elif go_mode == "l3bubble":
            _build_l3_bubble_map(new_page(), brand_rgb, data, outline, title="Graphic Organizer")
        elif go_mode == "timeline":
            _build_timeline_page(new_page(), brand_rgb, data, outline, title="Graphic Organizer")
        elif go_mode == "planchart":
            _build_plan_chart_page(new_page(), brand_rgb, data, outline, title="Graphic Organizer")
        else:
            _build_p6_mindmap(new_page(), data["mind_map"], title="Graphic Organizer")
        _build_p5_writing(new_page(), brand_rgb, data["writing"], title="Writing")

    # 删除模板自带的 7 个级别原始 slide，只留我们克隆出来的内容页
    if use_template and n_template:
        sld_lst = prs.slides._sldIdLst
        for sld in list(sld_lst)[:n_template]:
            r_id = getattr(sld, "rId", None)
            sld_lst.remove(sld)
            if r_id:
                with contextlib.suppress(Exception):
                    prs.part.drop_rel(r_id)

    # 全局收尾：锁死所有方格(自选图形)的 auto_size → 禁止随文字缩放，
    # 保证所有级别、所有页的方格严格等大、排版统一（用户硬要求 2026-06-04）。
    _lock_box_autosize(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    _write_worksheet_manifest(outline, out_path, data, images, go_mode, source_validation)
    return out_path


def _lock_box_autosize(prs) -> None:
    """把所有 slide 上【自选图形方格】的文本框 auto_size 设为 NONE。

    渲染器(PowerPoint/LibreOffice)默认会让自选图形“随文字自适应大小”，导致同一组
    方格因文字长短被撑成不同尺寸。锁成 NONE 后方格尺寸完全由代码控制 → 整齐统一。
    仅作用于 AUTO_SHAPE（方格/卡片），不影响流式文本框。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for slide in prs.slides:
        for sh in slide.shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.has_text_frame:
                    sh.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            except Exception:
                pass


def attach_worksheet_questions(
    outline: BookOutline, data, *, reading_q_count: int = 4
) -> None:
    """挂载 AI 抽取的 worksheet 内容到 outline。

    v1.9：兼容两种输入格式：
      - dict：直接用（match_pairs/word_bank/fill_blanks/sentence_mcs/reading_text/reading_mcs/writing/mind_map）
      - list[dict]（AI 抽出来的 6 道题 list，每条 {type, items, ...}）→ 自动跑 adapter

    v2.0: 加 reading_q_count 参数（4/6/8），控制 Reading MC 页题量。
    v6: 修「写死 4 与 L5-6 规格冲突」——L5/6 有两张 Reading 页（每页 4 题，共 8），
        故 L5/6 至少需要 8 题；这里按 level 自动抬高下限，避免第二张 Reading 页内容空/薄。
    """
    lvl_n = _level_num(getattr(outline, "level", "") or "")
    if lvl_n >= 5:
        reading_q_count = max(reading_q_count, 8)
    if isinstance(data, list):
        # 原始 6 题（含 title/instruction/answer_key）也存一份，供 Teacher Guide 同源取 Answer Key
        setattr(outline, "_worksheet_questions", data)
        data = _questions_list_to_template_data(data, outline)
    if isinstance(data, dict):
        data = dict(data)
        data["_reading_q_count"] = reading_q_count
    setattr(outline, "_worksheet_data", data)


def _verbatim_vocab(outline: BookOutline, n: int = 5) -> list[str]:
    """块6（用户拍板 2026-06-08）：worksheet 取词优先用大纲 verbatim（防被 AI 抽取覆盖）。"""
    syl = getattr(outline, "syllabus", None)
    if syl is not None:
        try:
            if outline.is_dual_vocab_level and getattr(syl, "vocab_mastery", None):
                w = list(syl.vocab_mastery)
            else:
                w = syl.vocab_words()
        except Exception:
            w = []
        w = [x.strip() for x in w if (x or "").strip()]
        if w:
            return w[:n]
    return (outline.vocabulary_mastery or outline.vocabulary_simple or [])[:n]


def _questions_list_to_template_data(qlist: list[dict], outline: BookOutline) -> dict:
    """把 AI 抽取的"按 level 池子的题目"映射到 worksheet PPTX 模板的固定字段。

    Worksheet 模板有 8 页固定结构：
      P1 Match (word ↔ definition)
      P2 Fill blanks (word_bank + 5 句 fill_blanks)
      P3 Sentence MC (4 题 2 选 1，可配图)
      P4 Sentence practice 2
      P5 Reading
      P6 Reading
      P7 Graphic Organizer
      P8 Writing

    AI 抽取的题型五花八门（match_definition / fill_blank / true_false / inference / unscramble / …）
    本函数把它们分发到对应模板字段，并用 outline 数据兜底空字段。
    """
    out: dict = {
        "match_pairs": [],
        "word_bank": [],
        "fill_blanks": [],
        "word_unscramble": [],   # 单词级"字母拼词"任务（归词汇区，绝不进 fill_blanks/句型区）
        "sentence_mcs": [],
        "reading_text": "",
        "reading_mcs": [],
        "writing": {},
        "mind_map": [],
    }

    for q in qlist or []:
        if not isinstance(q, dict):
            continue
        qtype = (q.get("type") or "").lower()
        # AI 抽取偶尔把 items 给成字符串列表（而非 dict 列表），统一规整成 dict，
        # 字符串元素塞进 _str，供 match/fill 等分支兜底使用，避免 .get 崩溃。
        items = [
            it if isinstance(it, dict) else {"_str": str(it)}
            for it in (q.get("items") or [])
            if it is not None
        ]
        extra = q.get("extra") or ""

        if qtype == "match_definition":
            out["match_pairs"] = [
                {"word": it.get("word") or it.get("_str", ""),
                 "def": it.get("def") or it.get("definition", "")}
                for it in items if (it.get("word") or it.get("_str"))
            ][:5]
            if not out["word_bank"]:
                out["word_bank"] = [
                    it.get("word") or it.get("_str", "") for it in items
                ][:5]

        elif qtype in ("fill_blank", "fill_blank_simple", "fill_blank_advanced", "emotion_fill"):
            out["fill_blanks"] = [
                {"sentence": it.get("sentence") or it.get("_str", ""),
                 "answer": it.get("answer", "")}
                for it in items if (it.get("sentence") or it.get("_str"))
            ][:5]
            if not out["word_bank"]:
                if extra:
                    bank = [w.strip() for w in extra.split(",") if w.strip()]
                else:
                    bank = [it.get("answer", "") for it in items]
                out["word_bank"] = bank[:5]

        elif qtype in ("true_false", "true_false_simple"):
            for it in items[:4]:
                stmt = it.get("statement", "")
                if not stmt:
                    continue
                # AI 给的 T/F 转成 sentence MC：正确句 vs 语法可读的反义句
                ans = (it.get("answer") or "T").upper()
                opt_true = stmt
                core = stmt.rstrip(".!?")
                opt_false = f"It is not true that {core}." if core else ""
                out["sentence_mcs"].append({
                    "options": [opt_true, opt_false],
                    "correct": 0 if ans == "T" else 1,
                })

        elif qtype in ("inference", "reading_mc"):
            for it in items[:4]:
                stem = it.get("q") or it.get("question", "")
                if not stem:
                    continue
                opts = it.get("options") or []
                if len(opts) < 2:
                    continue
                out["reading_mcs"].append({
                    "q": stem,
                    "options": list(opts)[:3],
                    "correct": int(it.get("correct", 0)),
                })

        elif qtype in ("plot_chart", "plot_chart_pbl"):
            # AI 给 {label: Setting/Problem/Solution/..., answer: 内容}
            buf: dict[str, str] = {}
            for it in items:
                buf[(it.get("label") or "").lower()] = it.get("answer", "")
            character = buf.get("characters") or buf.get("character") or "Main character"
            problem = buf.get("problem") or buf.get("conflict") or ""
            solution = buf.get("solution") or buf.get("resolution") or buf.get("ending") or ""
            out["mind_map"].append({
                "character": character,
                "problem": problem,
                "solution": solution,
            })

        elif qtype == "compare_contrast":
            for it in items[:3]:
                out["mind_map"].append({
                    "character": it.get("topic", "Comparison"),
                    "problem": it.get("side_a", ""),
                    "solution": it.get("side_b", ""),
                })

        elif qtype in ("essay_short", "personal_write", "personal_simple",
                       "draw_favorite", "open_ended_pbl", "research_pbl"):
            # v2.2：统一改成【基于本文】的引导式复述脚手架（不再用空泛 "Write about X" 题）
            out["writing"] = _retell_writing_scaffold(outline)

        elif qtype == "unscramble":
            # v7（用户拍板 2026-06-09）：这是【单词级】"字母拼词"任务，归【词汇区】，
            # 绝不再塞进 fill_blanks（句型区）——否则会被句型②的选词补全(cloze)误当作句子题，
            # 渲染出"Sentences / Choose the correct word… / Unscramble: s i n g → __ (A.hobby…)"
            # 这种题型↔区块↔标题↔指令全错、且又拼又选自相矛盾的页面。
            # 只收单词答案（多词=句子级，交给句型区自带的连词成句引擎，不在此处理）；
            # 字母乱序统一在渲染时用 _scramble_letters 真打乱（不信任 AI 给的 scrambled）。
            for it in items[:6]:
                ans = (it.get("answer") or "").strip()
                if ans and " " not in ans and ans.isalpha() and len(ans) >= 3:
                    out["word_unscramble"].append({
                        "answer": ans,
                        "clue": (it.get("clue") or it.get("def")
                                 or it.get("hint") or "").strip(),
                    })

        elif qtype in ("word_order", "word_order_simple", "story_sequence"):
            # v6：不再丢弃——落进 fill_blanks 槽位（版式不变），呈现为"按顺序编号"的排序题。
            # 兼容多种字段：{text/event/sentence, order} 或纯字符串。
            if not out["fill_blanks"]:
                seq = []
                for it in items:
                    txt = it.get("text") or it.get("event") or it.get("sentence") or it.get("_str", "")
                    if not txt:
                        continue
                    try:
                        order = int(it.get("order")) if it.get("order") is not None else None
                    except (TypeError, ValueError):
                        order = None
                    seq.append((order, txt))
                # 有 order 的按 order 给答案；缺 order 的按出现顺序兜底编号
                ordered = [t for t in seq if t[0] is not None]
                ordered.sort(key=lambda t: t[0])
                answer_map = {txt: i + 1 for i, (_, txt) in enumerate(ordered)} if ordered \
                    else {txt: i + 1 for i, (_, txt) in enumerate(seq)}
                for _, txt in seq[:5]:
                    out["fill_blanks"].append({
                        "sentence": f"___  {txt}",
                        "answer": str(answer_map.get(txt, "")),
                    })

        elif qtype in ("rewrite_tense", "rewrite_voice"):
            if not out["fill_blanks"]:
                for it in items[:5]:
                    prm = it.get("prompt", "")
                    ans = it.get("answer", "")
                    if prm and ans:
                        out["fill_blanks"].append({
                            "sentence": f"Rewrite: {prm} → ____",
                            "answer": ans,
                        })

        elif qtype in ("color_match", "circle_match", "word_to_pic"):
            # 简单 vocab cue，用作 word_bank 兜底
            if not out["word_bank"]:
                out["word_bank"] = [
                    it.get("word", "") for it in items if it.get("word")
                ][:5]

    # ----- 兜底字段：用 outline 数据补全 -----
    pages = outline.pages or []
    story_text = capitalize_names(" ".join(
        (p.text or "").strip() for p in pages if (p.text or "").strip()
    ).strip())

    if not out["reading_text"]:
        out["reading_text"] = story_text or "Story text goes here."

    # Match 兜底：用 vocab（块6：大纲 verbatim 优先）
    if not out["match_pairs"]:
        words = _verbatim_vocab(outline, 5)
        out["match_pairs"] = [
            {"word": w, "def": f"meaning of {w}"} for w in words
        ]
        if not out["word_bank"]:
            out["word_bank"] = words[:]

    # word_bank 兜底
    if not out["word_bank"]:
        out["word_bank"] = _verbatim_vocab(outline, 5)

    # Fill 兜底
    if not out["fill_blanks"] and out["word_bank"]:
        out["fill_blanks"] = [
            {"sentence": f"I feel ____ when I see this.", "answer": w}
            for w in out["word_bank"][:5]
        ]

    story_sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", story_text) if s.strip()]

    # Sentence MC 兜底：用故事原文 + 语法可读的反义句
    if not out["sentence_mcs"]:
        for s in story_sentences[:4]:
            core = s.rstrip(".!?")
            out["sentence_mcs"].append({
                "options": [s, f"It is not true that {core}." if core else ""],
                "correct": 0,
            })

    # Reading MC 兜底：从故事生成简单 5W（按上限取，不固定 4）
    if not out["reading_mcs"] and story_sentences:
        for i, s in enumerate(story_sentences[:8]):
            out["reading_mcs"].append({
                "q": f"What does the passage say in sentence {i + 1}?",
                "options": [
                    s[:60] + ("..." if len(s) > 60 else ""),
                    "Something opposite happens.",
                    "The story does not mention it.",
                ],
                "correct": 0,
            })

    # Writing 兜底（v2.2：基于本文的引导式复述）
    if not out["writing"]:
        out["writing"] = _retell_writing_scaffold(outline)

    # Mind Map 兜底
    if not out["mind_map"]:
        out["mind_map"] = [{
            "character": "Main character",
            "problem": "What is the problem in the story?",
            "solution": "How is it solved?",
        }]

    return out


# ============================================================
#  品牌外框
# ============================================================

def _draw_brand_frame(slide, brand_rgb: tuple, footer_text: str, logo_icon: Path) -> None:
    """所有 6 页统一的外框：粉背景 + 圆角白底 + 左上 logo + 右上 Name + 右下 footer。"""
    # 1. 外背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*brand_rgb)
    bg.line.fill.background()
    bg.shadow.inherit = False

    # 2. 内白底圆角
    content = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X), Inches(CONTENT_Y),
        Inches(CONTENT_W), Inches(CONTENT_H),
    )
    content.adjustments[0] = CONTENT_ROUND
    content.fill.solid()
    content.fill.fore_color.rgb = WHITE
    content.line.fill.background()
    content.shadow.inherit = False

    # 3. 左上 logo (icon + 文字)
    if logo_icon and logo_icon.exists():
        try:
            slide.shapes.add_picture(
                str(logo_icon),
                Inches(LOGO_X), Inches(LOGO_Y),
                width=Inches(LOGO_ICON_W), height=Inches(LOGO_ICON_H),
            )
        except Exception:
            pass
    text_x = LOGO_X + LOGO_ICON_W + 0.10
    tb = slide.shapes.add_textbox(
        Inches(text_x), Inches(LOGO_Y), Inches(4.5), Inches(LOGO_ICON_H),
    )
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run()
    r.text = "VIPKID Dino Reading Club"
    r.font.name = FONT
    r.font.size = Pt(LOGO_TEXT_PT)
    r.font.bold = True
    r.font.color.rgb = WHITE  # 块7：品牌色标题条统一白字（跨级别一致）

    # 4. 右上 Name 角标（v2.0 改回盾形：上方矩形 + 下方三角尖角，对齐官方模板）
    # Name 颜色用 brand_rgb 加深 30%（更接近模板的酒红/暗粉效果）
    name_dark = RGBColor(
        max(0, brand_rgb[0] - 50),
        max(0, brand_rgb[1] - 30),
        max(0, brand_rgb[2] - 30),
    )

    # 矩形部分（上半部）
    name_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(NAME_X), Inches(NAME_Y),
        Inches(NAME_W), Inches(NAME_H),
    )
    name_rect.fill.solid()
    name_rect.fill.fore_color.rgb = name_dark
    name_rect.line.fill.background()
    name_rect.shadow.inherit = False
    tf = name_rect.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Name"
    r.font.name = FONT
    r.font.size = Pt(NAME_PT)
    r.font.color.rgb = _readable_text_rgb(name_dark)  # Bug5：底色亮度自适应
    r.font.bold = True

    # 三角形尖角部分（朝下，宽度比矩形窄，居中下方）
    tri_w = NAME_W * 0.55
    tri_h = 0.22
    tri_x = NAME_X + (NAME_W - tri_w) / 2
    tri_y = NAME_Y + NAME_H
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(tri_x), Inches(tri_y),
        Inches(tri_w), Inches(tri_h),
    )
    tri.rotation = 180  # 翻转让尖朝下
    tri.fill.solid()
    tri.fill.fore_color.rgb = name_dark
    tri.line.fill.background()
    tri.shadow.inherit = False

    # 5. 底部 footer（v1.9：全宽 + 右对齐，往上挪到白底里面避免被裁切）
    fo = slide.shapes.add_textbox(
        Inches(FOOTER_X), Inches(FOOTER_Y),
        Inches(FOOTER_W), Inches(FOOTER_H),
    )
    fo.text_frame.margin_left = fo.text_frame.margin_right = Inches(0.10)
    p = fo.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = footer_text
    r.font.name = FONT
    r.font.size = Pt(FOOTER_PT)
    r.font.color.rgb = RGBColor(*brand_rgb)
    r.font.bold = True


def _get_subtitle(qtype_id: str, fallback: str) -> str:
    """v2.0 从题型库取标准英文 instruction；找不到就用兜底文本。"""
    try:
        from worksheet_question_types import get_type
        t = get_type(qtype_id)
        if t and t.en_instr:
            return t.en_instr
    except Exception:
        pass
    return fallback


def _add_title(slide, title: str, subtitle: str) -> None:
    """大标题 + 副标题（居中，位于内容白底顶部）。"""
    # Title
    tb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(CONTENT_Y + 0.14),
        Inches(CONTENT_W), Inches(0.60),
    )
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.bold = True
    # 块7：长标题自动缩字号保证【一行】，避免换行压到副标题/内容，或撑宽到右上 Name 角标
    title_pt = float(TITLE_PT)
    title_w = CONTENT_W - 0.30
    for cand in (40.0, 36.0, 32.0, 28.0, 24.0, 22.0):
        cpl = max(8, int(title_w / (cand / 72.0 * 0.60)))  # 粗体约 0.6em/字
        if len(title or "") <= cpl:
            title_pt = cand
            break
    else:
        title_pt = 22.0
    r.font.size = Pt(title_pt)
    r.font.color.rgb = TITLE_RGB

    # Subtitle —— 自动缩字号保证【一行】，避免长指令换行压到下方内容框
    sub_pt = float(SUBTITLE_PT)
    sub_w = CONTENT_W - 0.40
    for cand in (14.0, 13.0, 12.0):
        cpl = max(10, int(sub_w / (cand / 72.0 * 0.50)))
        if len(subtitle or "") <= cpl:
            sub_pt = cand
            break
    else:
        sub_pt = 12.0
    # 副标题（=command 指令）：与大标题拉开一段距离（用户反馈大标题/副标题贴太近），
    #   同时仍把下方空间留给首题。
    sb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(CONTENT_Y + 0.88),
        Inches(CONTENT_W), Inches(0.32),
    )
    sb.text_frame.margin_left = sb.text_frame.margin_right = 0
    sb.text_frame.word_wrap = False
    p2 = sb.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = FONT
    r2.font.size = Pt(sub_pt)
    r2.font.color.rgb = SUB_RGB


# ============================================================
#  Page 1 — Vocabulary 连线
# ============================================================

def _derange_order(n: int) -> list[int]:
    """返回 0..n-1 的一个错位排列（保证每个位置 i 的取值 != i）。

    用于连线题：右列定义要打乱顺序，且不能有任何一条定义恰好落在它对应单词的同一行
    （否则等于把正确答案直接画成水平对齐）。n<=1 时无法错位，原样返回。
    """
    if n <= 1:
        return list(range(n))
    idx = list(range(n))
    for _ in range(50):
        random.shuffle(idx)
        if all(idx[i] != i for i in range(n)):
            return idx
    # 兜底：整体循环移位（天然错位）
    return [(i + 1) % n for i in range(n)]


def _build_p1_match(slide, brand_rgb: tuple, pairs: list[dict], images: list[Path]) -> None:
    """5 对连线：左列小图（绘本图） + 中列粉色实心词卡 ↔ 右列白底粉边定义卡。

    v1.8 新增：每个 vocab 配一张绘本小图（用 page_02..page_06.png 循环）当 visual cue。
    v2.2 修复：右列定义【打乱顺序】渲染（错位排列），让连线题真正需要学生自己连线，
              而不是行对行直接把正确答案对齐给出。
    """
    _add_title(slide, "Vocabulary", _get_subtitle("vocab_match_definition", "Match the words to their definitions."))

    n = min(len(pairs), 6)
    if n == 0:
        return

    # 右列定义的错位渲染顺序：def_render_order[row] = 该行要显示 pairs 里第几个的定义
    def_render_order = _derange_order(n)

    # 区域
    area_top = CONTENT_Y + 1.30
    area_bottom = CONTENT_Y + CONTENT_H - 0.30
    area_h = area_bottom - area_top

    # 三列布局：放大绘本配图 + 缩小文字（用户反馈：图太小、字太大）
    img_x = CONTENT_X + 0.40
    img_w = 1.55
    word_x = img_x + img_w + 0.28
    word_w = 1.85
    def_x = word_x + word_w + 0.30
    def_w = CONTENT_W - (def_x - CONTENT_X) - 0.40
    row_gap = 0.18
    row_h = (area_h - row_gap * (n - 1)) / n
    word_pt = 15

    # 统一字号：让【最长定义】在固定行高内放得下 → 所有方格同尺寸、同字号、排版整齐
    _defs = [
        str((pairs[def_render_order[i]] if def_render_order[i] < len(pairs)
             else pairs[i]).get("def", ""))
        for i in range(n)
    ]
    def_pt = 10
    for cand in (13, 12, 11, 10):
        lh = cand / 72.0 * 1.16
        max_lines = max((_est_lines(d, def_w - 0.30, cand) for d in _defs), default=1)
        if max_lines * lh <= row_h - 0.14:
            def_pt = cand
            break

    for i, pair in enumerate(pairs[:n]):
        y = area_top + i * (row_h + row_gap)

        # 小图（左，绘本插画，正方形等比缩放）
        img_idx = (i % max(1, len(images) - 2)) + 2  # 从 page_02 起
        img_path = images[img_idx] if (img_idx < len(images) and images[img_idx]) else None
        if img_path and img_path.exists():
            try:
                from PIL import Image as _PILImg
                with _PILImg.open(str(img_path)) as _pim:
                    iw, ih = _pim.size
                aspect = iw / ih if ih else 1.0
                fit_w = img_w
                fit_h = fit_w / aspect
                if fit_h > row_h:
                    fit_h = row_h
                    fit_w = fit_h * aspect
                off_x = img_x + (img_w - fit_w) / 2
                off_y = y + (row_h - fit_h) / 2
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(off_x), Inches(off_y),
                    width=Inches(fit_w), height=Inches(fit_h),
                )
            except Exception:
                _draw_image_placeholder(slide, img_x, y, img_w, row_h)
        else:
            _draw_image_placeholder(slide, img_x, y, img_w, row_h)

        # 词卡（中，粉色实心，白字）
        wc = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(word_x), Inches(y), Inches(word_w), Inches(row_h),
        )
        wc.adjustments[0] = 0.28
        wc.fill.solid()
        wc.fill.fore_color.rgb = RGBColor(*brand_rgb)
        wc.line.fill.background()
        wc.shadow.inherit = False
        tf = wc.text_frame
        tf.margin_left = tf.margin_right = Inches(0.1)
        tf.margin_top = tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE   # 锁死方格尺寸，禁止随文字缩放（保证等大）
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(pair.get("word", "")).strip()
        r.font.name = FONT
        r.font.size = Pt(word_pt)
        r.font.color.rgb = WHITE
        r.font.bold = False

        # 定义卡（右，白底粉边，黑字）
        dc = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(def_x), Inches(y), Inches(def_w), Inches(row_h),
        )
        dc.adjustments[0] = 0.28
        dc.fill.solid()
        dc.fill.fore_color.rgb = WHITE
        dc.line.color.rgb = RGBColor(*brand_rgb)
        dc.line.width = Pt(1.8)
        dc.shadow.inherit = False
        tf = dc.text_frame
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.margin_top = tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE   # 锁死方格尺寸，禁止随文字缩放（保证等大）
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT  # v1.8: 定义文本左对齐更易读
        r = p.add_run()
        # v2.2: 定义按错位顺序取，避免和左侧单词行对行（即直接给出正确答案）
        def_pair = pairs[def_render_order[i]] if def_render_order[i] < len(pairs) else pair
        r.text = str(def_pair.get("def", "")).strip()
        r.font.name = FONT
        r.font.size = Pt(def_pt)
        r.font.color.rgb = BLACK


# ============================================================
#  Page 2 — Vocabulary 填空
# ============================================================

def _build_p2_fill(slide, brand_rgb: tuple, fills: list[dict], word_bank: list[str],
                   images: Optional[list[Path]] = None,
                   title: str = "Vocabulary", subtitle: str | None = None,
                   fallback: Optional[list[dict]] = None) -> None:  # noqa: ARG001
    """顶部粉色词库条 + 5 道填空题（用 ____ 表示空）。images 参数留作未来扩展。

    title/subtitle 可覆盖（第 2 句型页复用本渲染器时传 title='Sentence'）。
    fallback → fills 为空时渲染的兜底书写题（Bug1：句型页杜绝空页）。
    """
    if not fills and fallback:
        _add_title(slide, title, "Read, copy and write the sentences.")
        _render_sentence_fallback(slide, brand_rgb, fallback)
        return
    _add_title(slide, title,
               subtitle or _get_subtitle("vocab_fill_blank", "Use the words to fill each blank."))

    has_bank = bool(word_bank)

    # 词库条（粉色实心圆角，水平排列词）—— 仅在有词库时绘制
    bank_top = CONTENT_Y + 1.80
    phrase_bank = title.strip().lower().startswith("sentence") and any(" " in str(w).strip() for w in word_bank)
    bank_h = 0.92 if phrase_bank else 0.55
    bank_x = CONTENT_X + 1.50
    bank_w = CONTENT_W - 3.00
    if word_bank:
        if phrase_bank:
            card_gap = 0.18
            card_w = (bank_w - card_gap) / 2
            card_h = (bank_h - card_gap) / 2
            for i, word in enumerate(word_bank[:4]):
                rr, cc = divmod(i, 2)
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(bank_x + cc * (card_w + card_gap)),
                    Inches(bank_top + rr * (card_h + card_gap)),
                    Inches(card_w), Inches(card_h),
                )
                card.adjustments[0] = 0.25
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*brand_rgb)
                card.line.fill.background()
                card.shadow.inherit = False
                tf = card.text_frame
                tf.margin_left = tf.margin_right = Inches(0.08)
                tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = _clean_text(word)
                r.font.name = FONT
                r.font.size = Pt(12.5 if len(str(word)) > 18 else 14)
                r.font.color.rgb = WHITE
        else:
            bk = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bank_x), Inches(bank_top), Inches(bank_w), Inches(bank_h),
            )
            bk.adjustments[0] = 0.4
            bk.fill.solid()
            bk.fill.fore_color.rgb = RGBColor(*brand_rgb)
            bk.line.fill.background()
            bk.shadow.inherit = False
            tf = bk.text_frame
            tf.margin_left = tf.margin_right = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            # 水平排列：词之间用空格
            joined = "    ".join(_clean_text(w) for w in word_bank)
            r = p.add_run()
            r.text = joined
            r.font.name = FONT
            r.font.size = Pt(BODY_PT)
            r.font.color.rgb = WHITE  # 块7：品牌色标题条统一白字（跨级别一致）
            r.font.bold = False

    # 填空题（垂直排列）—— 最多 6 题（用户：4-6 题，保持整洁）
    n = min(len(fills), 6)
    if n == 0:
        return
    # 有词库 → 题目区从词库下方开始；无词库（如"写出动词过去式"页）→ 紧贴标题下方，
    #   避免顶部大片空白（用户反馈：中间一大块空、不好看）。
    qa_top = (bank_top + bank_h + 0.45) if has_bank else (CONTENT_Y + 1.30)
    qa_bottom = CONTENT_Y + CONTENT_H - 0.30
    qa_h = qa_bottom - qa_top
    box_w = CONTENT_W - 1.60
    # 无词库时字号略大，把版面填得均衡好看
    fill_pt = BODY_PT if has_bank else (BODY_PT + 2)

    items = [
        f"{i + 1}.  {capitalize_names(_ensure_blank(_clean_text(qa.get('sentence', ''))))}"
        for i, qa in enumerate(fills[:n])
    ]
    # B1：按单题实际换行数自适应每题高度；B4：题间统一间距、整体垂直居中。
    # gap_max 放宽到 1.1"：题少（3-4 题）时也把版面铺满，避免上下大片留白（用户拍板 2026-06-06）。
    _flow_lines(slide, items, qa_top, qa_h, CONTENT_X + 0.80, box_w, fill_pt,
                gap_min=0.18, gap_max=1.1)


def _est_lines(text: str, box_w_in: float, font_pt: float) -> int:
    """估算一段文本在给定宽度/字号下的换行行数（保守略多估，避免重叠）。"""
    cpl = max(8, int(box_w_in / (font_pt / 72.0 * 0.52)))
    n = 1
    for seg in str(text).split("\n"):
        n += max(0, -(-len(seg) // cpl) - 1) if seg else 0
    return max(1, n)


def _flow_lines(slide, items: list[str], area_top: float, area_h: float,
                x_in: float, box_w: float, font_pt: float,
                *, gap_min: float = 0.18, gap_max: float = 0.55,
                color: RGBColor = BLACK) -> None:
    """流式排版一组题干：每题高度随实际行数自适应，题间统一间距，整体垂直居中。

    解决"按行号均分 + 居中"导致的单/双行题间距忽大忽小、换行题被下一题挤压的问题。
    """
    if not items:
        return
    while font_pt > 11:
        line_h_try = font_pt / 72.0 * 1.20
        heights_try = [_est_lines(t, box_w, font_pt) * line_h_try + 0.06 for t in items]
        if sum(heights_try) + gap_min * max(0, len(items) - 1) <= area_h:
            break
        font_pt -= 0.5
    line_h = font_pt / 72.0 * 1.20
    pad = 0.06
    heights = [_est_lines(t, box_w, font_pt) * line_h + pad for t in items]
    total_content = sum(heights)
    n = len(items)
    # 统一题间距 = 剩余空间均分，夹在 [gap_min, gap_max]
    slack = area_h - total_content
    gap = slack / (n + 1) if n > 0 else 0
    gap = max(gap_min, min(gap_max, gap))
    block_h = total_content + gap * (n - 1)
    # 整体垂直居中（剩余空间不足时从顶部开始）
    y = area_top + max(0.0, (area_h - block_h) / 2.0)

    for text, h in zip(items, heights):
        tb = slide.shapes.add_textbox(Inches(x_in), Inches(y), Inches(box_w), Inches(h))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.0)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.12
        _emit_with_underscore_lock(p, text, font_pt, color)
        y += h + gap


def _ensure_blank(text: str) -> str:
    """若 sentence 里没有 ____ 占位，自动加 underscores。"""
    if "____" in text or "_____" in text or "___" in text:
        return text
    # 把第一个 [blank] / [BLANK] / [BLK] 替换为 ____
    for token in ("[blank]", "[BLANK]", "[BLK]"):
        if token in text:
            return text.replace(token, "________", 1)
    return text


def _emit_with_underscore_lock(paragraph, text: str, size_pt: int, color: RGBColor,
                               *, bold: bool = False) -> None:
    """把文本里连续的下划线段（____）切出来用 Arial 字体渲染（Poppins 下 _ 太扁、断成一截截），
    其余部分仍用 Poppins。bold 控制非下划线文字是否加粗（下划线段始终用 Arial 实线）。"""
    import re as _re
    parts = _re.split(r"(_{3,})", text)  # 至少 3 个 _ 才视作占位
    if len(parts) == 1:
        # 无下划线段，直接整段 Poppins
        r = paragraph.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size_pt)
        r.font.bold = bold
        r.font.color.rgb = color
        return
    for chunk in parts:
        if not chunk:
            continue
        r = paragraph.add_run()
        r.text = chunk
        if chunk.startswith("___"):
            r.font.name = FONT_BLANK  # Arial，下划线粗实连续
            r.font.bold = False
        else:
            r.font.name = FONT
            r.font.bold = bold
        r.font.size = Pt(size_pt)
        r.font.color.rgb = color


def _draw_writing_line(slide, x_in: float, y_in: float, w_in: float,
                       color: RGBColor = LIGHT_GRAY, weight_pt: float = 1.2) -> None:
    """画一条整宽作答横线（用直线连接符，比下划线字符更整齐、长度精准）。"""
    ln = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x_in), Inches(y_in), Inches(x_in + w_in), Inches(y_in),
    )
    ln.line.color.rgb = color
    ln.line.width = Pt(weight_pt)
    try:
        ln.shadow.inherit = False
    except Exception:
        pass


# ============================================================
#  Page 3 — Sentence MC (二选一 + 绘本配图)
# ============================================================

def _build_p3_sentence(slide, brand_rgb: tuple, mcs: list[dict], images: list[Path],
                       show_images: bool = True, *,
                       subtitle: str = "Tick (\u2713) the sentence that uses the correct past tense.",
                       fallback: Optional[list[dict]] = None) -> None:
    """4 题二选一，每题左侧绘本图 + 右侧 A/B 选项 + 行首圆圈题号。

    show_images=False → 不配图（选项满宽展开），供老师在「配图来源」选『不配图』时使用。
    subtitle → 按考点（过去式/主谓一致…）覆盖副标题。
    fallback → mcs 为空时渲染的兜底书写题（Bug1：杜绝空页）。
    """
    # 带图每行占位多 → 最多 4 题；不配图 → 选项满宽，可放 5 题（用户：句型要多一点但保持整洁）
    n = min(len(mcs), 4 if show_images else 5)
    if n == 0:
        # 句型题引擎无产出 → 用官方句型/例句兜底，保证不空页
        _add_title(slide, "Sentences", "Read, copy and write the sentences.")
        _render_sentence_fallback(slide, brand_rgb, fallback or [])
        return
    _add_title(slide, "Sentences", subtitle)

    area_top = CONTENT_Y + 1.40
    area_bottom = CONTENT_Y + CONTENT_H - 0.30
    area_h = area_bottom - area_top
    row_gap = 0.10
    row_h = (area_h - row_gap * (n - 1)) / n

    qnum_size = 0.45
    qnum_x = CONTENT_X + 0.30
    img_x = qnum_x + qnum_size + 0.20
    img_w = 1.80 if show_images else 0.0
    opt_x = (img_x + img_w + 0.30) if show_images else (qnum_x + qnum_size + 0.30)
    opt_w = CONTENT_W - (opt_x - CONTENT_X) - 0.40

    # 自适应字号：配图时选项区变窄，长句易折行→拥挤重叠。按"最长选项 + 可用文字宽度"
    # 反推一个尽量让每条选项落在 1 行的字号（标定：18pt 时约 0.133in/字符），夹在 [12, BODY_PT]。
    cb_size = 0.22
    text_w = max(2.0, opt_w - cb_size - 0.20)
    max_chars = 1
    for _mc in mcs[:n]:
        for _j, _opt in enumerate((_mc.get("options") or [])[:2]):
            max_chars = max(max_chars, len(f"{chr(ord('A') + _j)}. {_opt}"))
    char_w_at_18 = 0.133
    fit_pt = text_w / (max_chars * (char_w_at_18 / 18.0))
    opt_pt = max(12.0, min(float(BODY_PT), fit_pt))

    for i, mc in enumerate(mcs[:n]):
        y = area_top + i * (row_h + row_gap)
        cy = y + (row_h - qnum_size) / 2

        # 圆形题号
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(qnum_x), Inches(cy),
            Inches(qnum_size), Inches(qnum_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*brand_rgb)
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(QNUM_PT)
        r.font.color.rgb = WHITE  # 块7：品牌色标题条统一白字（跨级别一致）

        # 图（复用绘本 page_02..page_05），保持长宽比，居中放入 img_w x (row_h-0.10) 框
        img_idx = i + 2  # P3 题 1 用绘本 page_02
        img_path = images[img_idx] if (img_idx < len(images) and images[img_idx]) else None
        max_box_w = img_w
        max_box_h = row_h - 0.10
        if not show_images:
            pass  # 不配图：跳过图片/占位，选项满宽展开
        elif img_path and img_path.exists():
            try:
                from PIL import Image as _PILImg
                with _PILImg.open(str(img_path)) as _pim:
                    iw, ih = _pim.size
                aspect = iw / ih if ih else 1.0
                # 尝试按宽适配
                fit_w = max_box_w
                fit_h = fit_w / aspect
                if fit_h > max_box_h:
                    fit_h = max_box_h
                    fit_w = fit_h * aspect
                off_x = img_x + (max_box_w - fit_w) / 2
                off_y = y + 0.05 + (max_box_h - fit_h) / 2
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(off_x), Inches(off_y),
                    width=Inches(fit_w), height=Inches(fit_h),
                )
            except Exception:
                _draw_image_placeholder(slide, img_x, y + 0.05, max_box_w, max_box_h)
        else:
            _draw_image_placeholder(slide, img_x, y + 0.05, max_box_w, max_box_h)

        # 选项 A/B（垂直堆叠 + 复选框）
        options = mc.get("options") or []
        opt_h = (row_h - 0.05) / max(len(options), 1)
        for j, opt in enumerate(options[:2]):
            oy = y + j * opt_h
            # 复选框：对齐到选项【首行】（顶部），双行选项时不再悬在两行中间
            cb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(opt_x), Inches(oy + 0.09),
                Inches(cb_size), Inches(cb_size),
            )
            cb.fill.solid()
            cb.fill.fore_color.rgb = WHITE
            cb.line.color.rgb = BLACK
            cb.line.width = Pt(1.0)
            cb.shadow.inherit = False
            # 选项文字（顶部对齐 + 行距，便于和复选框首行齐平）
            tb = slide.shapes.add_textbox(
                Inches(opt_x + cb_size + 0.15), Inches(oy),
                Inches(opt_w - cb_size - 0.20), Inches(opt_h),
            )
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.12
            r = p.add_run()
            letter = chr(ord("A") + j)
            r.text = f"{letter}. {opt}"
            r.font.name = FONT
            r.font.size = Pt(opt_pt)
            r.font.color.rgb = BLACK


def _draw_image_placeholder(slide, x: float, y: float, w: float, h: float) -> None:
    ph = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    ph.adjustments[0] = 0.05
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    ph.line.color.rgb = LIGHT_GRAY
    ph.line.width = Pt(1.0)
    ph.shadow.inherit = False


# ============================================================
#  Page 4 — Reading 全文 + 8 道 3 选
# ============================================================

def _build_p4_reading(slide, brand_rgb: tuple, text: str, mcs: list[dict]) -> None:
    _add_title(slide, "Reading", _get_subtitle("read_mc_questions", "Choose the correct answer for each question."))

    text = (text or "").strip()

    # 顶部红框 = Reading 全文。短文字数差异很大（长文会撑破固定框），
    # 按字符数自适应字号 + 行距，确保始终落在框内不压字。
    tlen = len(text)
    if tlen > 780:
        read_pt, read_ls = 10.5, 1.12
    elif tlen > 560:
        read_pt, read_ls = 11.5, 1.15
    elif tlen > 380:
        read_pt, read_ls = 12.5, 1.20
    else:
        read_pt, read_ls = float(READING_PT), 1.25

    text_top = CONTENT_Y + 1.35
    text_h = 2.15
    text_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.40), Inches(text_top),
        Inches(CONTENT_W - 0.80), Inches(text_h),
    )
    text_box.adjustments[0] = 0.03
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = WHITE
    text_box.line.color.rgb = READ_BORDER
    text_box.line.width = Pt(1.5)
    text_box.shadow.inherit = False
    # python-pptx 不直接支持 dashed line per shape — 这里用实线代替

    tf = text_box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = _clean_text(text)
    r.font.name = FONT
    r.font.size = Pt(read_pt)
    r.font.color.rgb = BLACK
    p.line_spacing = read_ls

    # 下方 4 道 MC（2x2 排版）
    n = min(len(mcs), 4)
    if n == 0:
        return
    mc_top = text_top + text_h + 0.20
    mc_bottom = CONTENT_Y + CONTENT_H - 0.25
    mc_h_total = mc_bottom - mc_top
    rows_per_col = (n + 1) // 2  # 4 题 → 2/2
    row_h = mc_h_total / max(rows_per_col, 1)

    col_w = (CONTENT_W - 0.80) / 2
    col_x = [CONTENT_X + 0.40, CONTENT_X + 0.40 + col_w]

    # 题干+选项字号按"最长一题"的总字符数自适应，避免换行后相互重叠/压页脚
    def _q_chars(mc: dict) -> int:
        return len(str(mc.get("q", ""))) + sum(
            len(str(o)) for o in (mc.get("options") or [])[:3]
        )

    max_chars = max((_q_chars(m) for m in mcs[:n]), default=0)
    if max_chars > 150:
        q_pt = 10.5
    elif max_chars > 115:
        q_pt = 11.5
    elif max_chars > 85:
        q_pt = 12.5
    else:
        q_pt = 14.0

    for i, mc in enumerate(mcs[:n]):
        col_idx = 0 if i < rows_per_col else 1
        row_idx = i if col_idx == 0 else (i - rows_per_col)
        x = col_x[col_idx] + 0.10
        y = mc_top + row_idx * row_h
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(col_w - 0.20), Inches(row_h - 0.05),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        # 第一段：题号 + 题干
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = f"{i + 1}. {_clean_text(mc.get('q', ''))}"
        r.font.name = FONT
        r.font.size = Pt(q_pt)
        r.font.color.rgb = BLACK
        r.font.bold = False
        # 后续段：选项 A/B/C 每题占一段（避免选项过长被截断）
        opts = (mc.get("options") or [])[:3]
        for j, opt in enumerate(opts):
            po = tf.add_paragraph()
            po.alignment = PP_ALIGN.LEFT
            po.space_after = Pt(1)
            po.line_spacing = 1.05
            ro = po.add_run()
            ro.text = f"    {chr(ord('A') + j)}. {_clean_text(opt)}"
            ro.font.name = FONT
            ro.font.size = Pt(q_pt)
            ro.font.color.rgb = BLACK


def _build_word_fill_page(slide, brand_rgb: tuple, items: list[dict],
                          title: str, subtitle: str) -> None:
    """Word Fill 版式：每题 = 释义提示行 + 下方"掩码 + 作答空"。
    L0-2 用首字母掩码，L5-6 用词尾掩码（构词）。"""
    _add_title(slide, title, subtitle)
    n = min(len(items), 6)
    if n == 0:
        return
    area_top = CONTENT_Y + 1.30
    area_bottom = CONTENT_Y + CONTENT_H - 0.46   # 留足底部页脚间距
    avail = area_bottom - area_top
    box_w = CONTENT_W - 1.60
    clues = [f"{i + 1}.  {capitalize_names(_clean_text(it.get('clue', '')))}" for i, it in enumerate(items[:n])]

    # 自适应字号：题少时放大、间距加宽 → 把版面填满、排版均衡（不再底部大片留白）
    def _hs(pt: float):
        lh = pt / 72.0 * 1.18
        return [(_est_lines(c, box_w, pt) + 1) * lh + 0.12 for c in clues]

    body_pt = float(BODY_PT)
    for cand in (BODY_PT + 4, BODY_PT + 2, BODY_PT, BODY_PT - 1):
        if sum(_hs(float(cand))) + 0.20 * (n - 1) <= avail:
            body_pt = float(cand)
            break
    heights = _hs(body_pt)
    total = sum(heights)
    gap = max(0.20, min(1.10, (avail - total) / (n + 1)))
    y = area_top + max(0.0, (avail - (total + gap * (n - 1))) / 2.0)
    for it, clue, h in zip(items[:n], clues, heights):
        tb = slide.shapes.add_textbox(
            Inches(CONTENT_X + 0.80), Inches(y), Inches(box_w), Inches(h),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.14
        r = p.add_run()
        r.text = clue
        r.font.name = FONT
        r.font.size = Pt(body_pt)
        r.font.color.rgb = BLACK
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = 1.14
        _emit_with_underscore_lock(p2, f"      {it.get('hint', '')}", body_pt + 1, BLACK)
        y += h + gap


def _build_mcq_page(slide, brand_rgb: tuple, title: str, subtitle: str,
                    questions: list[dict], *, start_no: int = 1,
                    answer_paren: bool = True) -> None:
    """L3 选择题通用排版（单栏，最多 4 选项 A/B/C/D，字号+间距自适应填满整页）。

    用于词汇谜语四选一 / 句型选词补全。对齐用户标杆 L3-30：
    题量尽量 4 道（不足则 3 道），字号随题量自适应放大、整页铺满不显空，
    每题 = 常规字重题干 + 一行 A/B/C/D 选项。
    answer_paren：题干末尾是否带 "(   )" 作答括号——
      · 选词补全(finish the sentence) → True（写字母）；
      · 圈词(circle the word)         → False（直接圈选项，不需括号）。
    用户标杆 L3-30：题干为常规字重（非加粗）。"""
    import math as _m
    qs = [q for q in (questions or []) if q.get("q")]
    _add_title(slide, title, subtitle)
    if not qs:
        return
    n = len(qs)
    x = CONTENT_X + 0.60
    box_w = CONTENT_W - 1.20
    area_top = CONTENT_Y + 1.46
    area_bottom = CONTENT_Y + CONTENT_H - 0.32
    avail = area_bottom - area_top

    def _cpl(pt: float, w: float) -> int:
        return max(8, int(6.6 * w * (12.5 / pt)))

    def _opts_lines(q: dict) -> list[str]:
        opts = [_clean_text(o) for o in (q.get("options") or []) if str(o).strip()][:4]
        return [f"{chr(65 + j)}. {o}" for j, o in enumerate(opts)]

    def _stem_text(q: dict, i: int) -> str:
        s = f"{start_no + i}. {capitalize_names(_clean_text(q.get('q', '')))}"
        if answer_paren:
            s += "   (        )"
        return s

    def _slot_need(sp: float) -> list:
        op = max(13.0, sp - 2.0)
        lh_s = sp / 72.0 * 1.16
        lh_o = op / 72.0 * 1.16
        out = []
        for i, q in enumerate(qs):
            slines = max(1, _m.ceil(len(_stem_text(q, i)) / _cpl(sp, box_w)))
            olines = sum(
                max(1, _m.ceil(len(opt) / _cpl(op, box_w - 0.30)))
                for opt in _opts_lines(q)
            ) or 1
            out.append(slines * lh_s + olines * lh_o + 0.16)
        return out

    # 选最大的题干字号：内容能塞进"每题平分到的槽高"，整页铺满
    stem_pt = 16.0
    for cand in (26, 24, 22, 21, 20, 19, 18, 17, 16):
        per = _slot_need(float(cand))
        if max(per) <= (avail / n) - 0.10 and sum(per) <= avail:
            stem_pt = float(cand)
            break
    opt_pt = max(13.0, stem_pt - 2.0)
    slot_h = avail / n

    for i, q in enumerate(qs):
        y = area_top + i * slot_h
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(slot_h - 0.06))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.14
        _emit_with_underscore_lock(p, _stem_text(q, i), int(round(stem_pt)), BLACK, bold=False)
        for opt in _opts_lines(q):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = 1.10
            p2.space_before = Pt(3)
            _set_indent(p2, 0.30)
            ro = p2.add_run()
            ro.text = opt
            ro.font.name = FONT
            ro.font.size = Pt(opt_pt)
            ro.font.color.rgb = BLACK


def _build_reading_fill_page(slide, brand_rgb: tuple, title: str, subtitle: str,
                             questions: list[dict], *, start_no: int = 1) -> None:
    """L3 阅读页（无原文）：单栏、字号+间距自适应填满整页，杜绝"太空"。

    支持 tf / short / mc 三类题：
      • mc → A/B/C/D 选项行 + 作答括号；
      • tf → "Your answer ( T / F ): ____"；
      • short → 题干 + 两条作答横线。
    题量尽量 4 道（不足 3 道也铺满），避免 _build_reading_page 在短题时产生大片空隙。"""
    import math as _m
    qs = [q for q in (questions or []) if q.get("q")]
    _add_title(slide, title, subtitle)
    if not qs:
        return
    n = len(qs)
    x = CONTENT_X + 0.60
    box_w = CONTENT_W - 1.20
    area_top = CONTENT_Y + 1.46
    area_bottom = CONTENT_Y + CONTENT_H - 0.32
    avail = area_bottom - area_top
    slot_h = avail / n

    def _cpl(pt: float, w: float) -> int:
        return max(8, int(6.6 * w * (12.5 / pt)))

    def _stem_text(q: dict, i: int) -> str:
        s = f"{start_no + i}. {capitalize_names(_clean_text(q.get('q', '')))}"
        pg = q.get("page")
        if pg:
            s += f"   (P{pg})"
        if q.get("kind", "mc") == "mc":
            s += "   (        )"
        return s

    def _ans_lines(q: dict) -> int:
        k = q.get("kind", "mc")
        return 1 if k in ("mc", "tf") else 2  # short：两条作答横线

    def _stem_need(sp: float) -> list:
        lh = sp / 72.0 * 1.18
        out = []
        for i, q in enumerate(qs):
            slines = max(1, _m.ceil(len(_stem_text(q, i)) / _cpl(sp, box_w)))
            out.append(slines * lh + _ans_lines(q) * (max(13.0, sp - 2.0) / 72.0 * 1.7) + 0.16)
        return out

    stem_pt = 16.0
    for cand in (24, 22, 21, 20, 19, 18, 17, 16):
        per = _stem_need(float(cand))
        if max(per) <= (slot_h - 0.10) and sum(per) <= avail:
            stem_pt = float(cand)
            break
    ans_pt = max(13.0, stem_pt - 2.0)

    for i, q in enumerate(qs):
        y = area_top + i * slot_h
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(slot_h - 0.06))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.16
        rs = p.add_run()
        rs.text = _stem_text(q, i)
        rs.font.name = FONT
        rs.font.size = Pt(stem_pt)
        rs.font.bold = True
        rs.font.color.rgb = BLACK
        kind = q.get("kind", "mc")
        if kind == "mc":
            opts = [_clean_text(o) for o in (q.get("options") or []) if str(o).strip()][:4]
            line = "        ".join(f"{chr(65 + j)}. {o}" for j, o in enumerate(opts))
            p2 = tf.add_paragraph()
            p2.line_spacing = 1.16
            p2.space_before = Pt(4)
            _set_indent(p2, 0.30)
            ro = p2.add_run()
            ro.text = line
            ro.font.name = FONT
            ro.font.size = Pt(ans_pt)
            ro.font.color.rgb = BLACK
        elif kind == "tf":
            p2 = tf.add_paragraph()
            p2.line_spacing = 1.16
            p2.space_before = Pt(6)
            _set_indent(p2, 0.30)
            ro = p2.add_run()
            ro.text = "Your answer ( T / F ):  "
            ro.font.name = FONT
            ro.font.size = Pt(ans_pt)
            ro.font.color.rgb = BLACK
            ru = p2.add_run()
            ru.text = "______"
            ru.font.name = FONT_BLANK
            ru.font.size = Pt(ans_pt)
            ru.font.color.rgb = BLACK
        else:  # short：两条书写横线
            n_us = max(24, int((box_w - 0.30) / 0.105))
            for _ in range(2):
                p2 = tf.add_paragraph()
                p2.line_spacing = 1.7
                p2.space_before = Pt(4)
                _set_indent(p2, 0.30)
                ru = p2.add_run()
                ru.text = "_" * n_us
                ru.font.name = FONT_BLANK
                ru.font.size = Pt(ans_pt)
                ru.font.color.rgb = BLACK


def _build_word_fill_pic_page(slide, brand_rgb: tuple, items: list[dict]) -> None:
    """L3 词汇①：看图填缺失字母（2×2 网格）。
    每格 = 序号徽标 + 绘本配图 + 缺字母掩码 + 小释义。对齐官方样板 L3-8/L3-23。"""
    _add_title(slide, "Vocabulary", "Look at the picture and fill in the missing letters.")
    n = min(len(items), 4)
    if n == 0:
        return
    grid_top = CONTENT_Y + 1.34
    grid_bottom = CONTENT_Y + CONTENT_H - 0.34
    gx_gap, gy_gap = 0.70, 0.34
    cols = 2
    rows = (n + 1) // 2
    cell_w = (CONTENT_W - 0.80 - gx_gap) / 2.0
    cell_h = (grid_bottom - grid_top - gy_gap * (rows - 1)) / rows
    # 单行（2 题）时别把单元高度拉满整页，限高并整体竖向居中，配更大的图，杜绝下方大片空白
    if rows == 1:
        cell_h = min(cell_h, 3.6)
        grid_top = grid_top + (grid_bottom - grid_top - cell_h) / 2.0
    x0 = CONTENT_X + 0.40
    img_h = max(1.2, min(2.7 if rows == 1 else 1.95, cell_h - 1.05))

    for i, it in enumerate(items[:n]):
        rr, cc = divmod(i, cols)
        cx = x0 + cc * (cell_w + gx_gap)
        cy = grid_top + rr * (cell_h + gy_gap)

        # 序号徽标（圆，品牌色）
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.52), Inches(0.52))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*brand_rgb)
        badge.line.fill.background()
        badge.shadow.inherit = False
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{i + 1}"
        br.font.name = FONT
        br.font.bold = True
        br.font.size = Pt(15)
        br.font.color.rgb = WHITE

        # 配图（徽标右侧居中）
        img = it.get("img")
        img_w_box = cell_w - 0.62
        img_x = cx + 0.62
        img_y = cy
        placed = _draw_clean_vocab_clue(slide, _clean_vocab_clue_kind(it.get("answer", "")),
                                        img_x, img_y, img_w_box, img_h, brand_rgb)
        if not placed and img and Path(img).exists():
            try:
                from PIL import Image as _PILImg
                with _PILImg.open(str(img)) as _pim:
                    iw, ih = _pim.size
                aspect = iw / ih if ih else 1.4
                fit_h = img_h
                fit_w = fit_h * aspect
                if fit_w > img_w_box:
                    fit_w = img_w_box
                    fit_h = fit_w / aspect
                off_x = img_x + (img_w_box - fit_w) / 2
                slide.shapes.add_picture(
                    str(img), Inches(off_x), Inches(img_y),
                    width=Inches(fit_w), height=Inches(fit_h))
                placed = True
            except Exception:
                placed = False
        if not placed:
            _draw_image_placeholder(slide, img_x, img_y, img_w_box, img_h)

        # 缺字母掩码（大号、居中，图下方）
        mask_y = cy + img_h + 0.10
        mtb = slide.shapes.add_textbox(
            Inches(cx + 0.10), Inches(mask_y), Inches(cell_w - 0.20), Inches(0.46))
        mtf = mtb.text_frame
        mtf.margin_left = mtf.margin_right = 0
        mtf.margin_top = mtf.margin_bottom = 0
        mtf.word_wrap = True
        mp = mtf.paragraphs[0]
        mp.alignment = PP_ALIGN.CENTER
        _emit_with_underscore_lock(mp, _missing_letters_mask(it.get("answer", "")), 22, BLACK)
        for _r in mp.runs:
            _r.font.bold = True

        # 小释义提示（灰，掩码下方）
        clue = capitalize_names(_clean_text(it.get("clue", "")))
        if clue:
            ctb = slide.shapes.add_textbox(
                Inches(cx + 0.10), Inches(mask_y + 0.48), Inches(cell_w - 0.20), Inches(0.40))
            ctf = ctb.text_frame
            ctf.margin_left = ctf.margin_right = 0
            ctf.margin_top = ctf.margin_bottom = 0
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = clue
            cr.font.name = FONT
            cr.font.italic = True
            cr.font.size = Pt(11)
            cr.font.color.rgb = SUB_RGB


def _draw_clean_vocab_clue(slide, kind: str, x: float, y: float, w: float, h: float,
                           brand_rgb: tuple) -> bool:
    """Draw a clean worksheet clue instead of using full reader-page screenshots."""
    if not kind:
        return False
    bc = RGBColor(*brand_rgb)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bg.adjustments[0] = 0.08
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xF2)
    bg.line.color.rgb = bc
    bg.line.width = Pt(1.0)
    bg.shadow.inherit = False

    def box(px, py, pw, ph, fill=WHITE, line=bc):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + px), Inches(y + py),
            Inches(pw), Inches(ph)
        )
        sh.adjustments[0] = 0.08
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = line
        sh.line.width = Pt(1.0)
        sh.shadow.inherit = False
        return sh

    def label(text, px, py, pw, ph, pt=10.5):
        tb = slide.shapes.add_textbox(Inches(x + px), Inches(y + py), Inches(pw), Inches(ph))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = True
        r.font.color.rgb = bc

    if kind == "calendar":
        cal = box(w * 0.20, h * 0.14, w * 0.60, h * 0.68)
        # Draw a seven-block calendar without extra English cue words.
        for i in range(6):
            xx = w * 0.26 + i * w * 0.08
            dot = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x + xx), Inches(y + h * 0.62),
                Inches(w * 0.035), Inches(h * 0.035)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = bc
            dot.line.fill.background()
    elif kind == "desk":
        box(w * 0.18, h * 0.24, w * 0.64, h * 0.36)
        for i in range(3):
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + w * 0.30), Inches(y + h * (0.34 + i * 0.08)),
                Inches(x + w * 0.66), Inches(y + h * (0.34 + i * 0.08)),
            ).line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x + w * 0.26), Inches(y + h * 0.66),
            Inches(x + w * 0.74), Inches(y + h * 0.66),
        ).line.color.rgb = bc
    elif kind == "notebook":
        box(w * 0.24, h * 0.16, w * 0.52, h * 0.64)
        for i in range(4):
            cb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + w * 0.32), Inches(y + h * (0.28 + i * 0.10)),
                Inches(w * 0.05), Inches(h * 0.05),
            )
            cb.fill.solid()
            cb.fill.fore_color.rgb = WHITE
            cb.line.color.rgb = bc
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + w * 0.42), Inches(y + h * (0.305 + i * 0.10)),
                Inches(x + w * 0.66), Inches(y + h * (0.305 + i * 0.10)),
            ).line.color.rgb = bc
    elif kind == "piano":
        kb = box(w * 0.14, h * 0.32, w * 0.72, h * 0.32)
        for i in range(6):
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + w * (0.14 + i * 0.12)), Inches(y + h * 0.32),
                Inches(x + w * (0.14 + i * 0.12)), Inches(y + h * 0.64),
            ).line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        for i in (1, 2, 4):
            sh = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x + w * (0.14 + i * 0.12 - 0.025)),
                Inches(y + h * 0.32), Inches(w * 0.05), Inches(h * 0.18)
            )
            sh.fill.solid()
            sh.fill.fore_color.rgb = BLACK
            sh.line.fill.background()
    else:
        return False
    return True


def _build_prompt_line_page(slide, brand_rgb: tuple, items: list[dict],
                            title: str, subtitle: str, *,
                            prompt_key: str, example: Optional[str] = None,
                            fallback: Optional[list[dict]] = None) -> None:
    """连词成句 / 句子改写 版式：每题 = 提示句 + 一条【整宽】作答横线。

    用户拍板（2026-06-04）：
      • 作答横线必须和题目一样长（整条内容宽）、尽量多留书写空间；
      • 所有题（含第 6 题）必须落在白底内、绝不被页脚遮挡 → 字号/题距自适应。
    fallback → items 为空时渲染的兜底书写题（Bug1：杜绝空页）。
    """
    n = min(len(items), 6)
    if n == 0:
        _add_title(slide, title, "Read, copy and write the sentences.")
        _render_sentence_fallback(slide, brand_rgb, fallback or [])
        return
    _add_title(slide, title, subtitle)
    area_top = CONTENT_Y + 1.28
    x = CONTENT_X + 0.42
    box_w = CONTENT_W - 0.84            # 题目 + 横线统一宽度（更宽 → 书写空间更大）
    if example:
        eb = slide.shapes.add_textbox(
            Inches(x), Inches(area_top), Inches(box_w), Inches(0.34),
        )
        ep = eb.text_frame.paragraphs[0]
        ep.alignment = PP_ALIGN.LEFT
        er = ep.add_run()
        er.text = f"e.g.  {_clean_text(example)}"
        er.font.name = FONT
        er.font.italic = True
        er.font.size = Pt(13)
        er.font.color.rgb = SUB_RGB
        area_top += 0.40
    area_bottom = CONTENT_Y + CONTENT_H - 0.40   # 给页脚留足，杜绝遮挡
    avail = area_bottom - area_top
    prompts = [f"{i + 1}.  {capitalize_names(_clean_text(it.get(prompt_key, '')))}"
               for i, it in enumerate(items[:n])]

    # —— 对齐官方模板（L4-13 Sentences 页）：每题一个等高「题槽」slot —— 
    #   题干在 slot 顶部，作答横线落在 slot 下部 → 题干与横线之间留出充足书写区，
    #   学生就写在横线上方/横线上（用户拍板 2026-06-04：要给学生书写空间）。
    #   字号尽量大（官方 20pt），但保证每题最多 2 行且书写区不小于 0.34in。
    def _max_qlines(pt: float) -> int:
        return max(_est_lines(pr, box_w, pt) for pr in prompts)

    WRITE_MIN = 0.42       # 题干底 → 横线 的最小书写空间
    pt = 14.0
    for cand in (20.0, 18.0, 16.0, 15.0, 14.0):
        lh = cand / 72.0 * 1.16
        slot_need = _max_qlines(cand) * lh + WRITE_MIN + 0.14
        if slot_need * n <= avail:
            pt = cand
            break
    lh = pt / 72.0 * 1.16
    slot_h = avail / n
    for i, prompt in enumerate(prompts):
        y_q = area_top + i * slot_h
        plines = _est_lines(prompt, box_w, pt)
        q_h = plines * lh
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y_q), Inches(box_w), Inches(q_h),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.14
        _emit_with_underscore_lock(p, prompt, int(round(pt)), BLACK)
        # 作答横线：落在题槽下部，与题干之间留足书写区（≥ WRITE_MIN）
        q_bottom = y_q + q_h
        line_y = y_q + slot_h - 0.20
        line_y = max(line_y, q_bottom + WRITE_MIN)
        line_y = min(line_y, y_q + slot_h - 0.10)
        _draw_writing_line(slide, x, line_y, box_w)


def _set_indent(paragraph, inches: float) -> None:
    """给段落设置固定左缩进（marL，EMU）+ indent=0 → 整段悬挂缩进。

    用固定英寸值而非空格前缀，保证 A/B/C 选项首字符严格落在同一条垂直线上，
    且换行后的续行也对齐（杜绝参差不齐）。"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(inches * 914400)))
    pPr.set("indent", "0")


def _build_reading_page(
    slide, brand_rgb: tuple, text: str, questions: list[dict],
    *, subtitle: str = "Choose the correct answer for each question.",
    start_no: int = 1, show_passage: bool = True, title: str = "Reading",
    force_single_col: bool = False,
) -> None:
    """阅读页 v2.2：大标题（默认 Reading，可定制）+ 灰副标题 +（可选）完整原文框 + 题目（左右两列对称）。

    show_passage：
      • True（L5/6）：顶部印完整原文框，下方 4 题 2 左 2 右对称。
      • False（其他级别）：不印原文（常规排），题目用 (P#) 提示回看绘本，占满整页更舒展。
    title：页面大标题（默认 Reading；L3 选择题页可复用此版式作 Vocabulary/Sentences MCQ 页）。
    题目 kind 支持 mc / tf / short，禁止网格卡片。题号从 start_no 起。"""
    # 不印原文时，副标题里的 "Read the passage..." 不再适用 → 换成回看绘本的指引
    if not show_passage and subtitle.strip().lower().startswith("read the passage"):
        if "true" in subtitle.lower() or "(true)" in subtitle.lower():
            subtitle = "Look back at the book. Write T (true) or F (false) for each statement."
        else:
            subtitle = "Look back at the book. Choose the correct answer for each question."
    _add_title(slide, title, subtitle)

    text = (text or "").strip()

    if show_passage:
        tlen = len(text)
        if tlen > 780:
            read_pt, read_ls = 10.5, 1.12
        elif tlen > 560:
            read_pt, read_ls = 11.5, 1.15
        elif tlen > 380:
            read_pt, read_ls = 12.5, 1.20
        else:
            read_pt, read_ls = float(READING_PT), 1.25

        text_top = CONTENT_Y + 1.35
        # 自适应正文框高：短文压缩，给下方题目（含 P#/Hint）腾出空间，避免裁题。
        import math as _mh
        _cpl_passage = max(40, int(94 * (12.5 / read_pt)))
        _passage_lines = max(1, _mh.ceil(tlen / _cpl_passage))
        text_h = min(2.4, max(1.0, _passage_lines * (read_pt / 72.0 * read_ls) + 0.30))
        text_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(CONTENT_X + 0.40), Inches(text_top),
            Inches(CONTENT_W - 0.80), Inches(text_h),
        )
        text_box.adjustments[0] = 0.03
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = WHITE
        # 边框跟随级别品牌色（用户硬要求 2026-06-06）：L5 粉、L4 蓝绿…与模板统一，不再固定红框。
        text_box.line.color.rgb = RGBColor(*brand_rgb)
        text_box.line.width = Pt(1.5)
        text_box.shadow.inherit = False
        tf = text_box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.20)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.06)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = _clean_text(text)
        r.font.name = FONT
        r.font.size = Pt(read_pt)
        r.font.color.rgb = BLACK
        p.line_spacing = read_ls
        list_top = text_top + text_h + 0.18
    else:
        # 不印原文：题目从标题下方直接开始，占满整页
        list_top = CONTENT_Y + 1.35

    qs = [q for q in (questions or []) if q.get("q")]
    if not qs:
        return

    # ===== 判断题(T/F) 专用排版（用户标杆 L3-30）=====
    # 全为 tf 时：单栏、前置 "(   )" 作答括号 + 陈述句，整页等距铺满；
    # 不再用 "题干 + Your answer (T / F): ___" 两行式（更紧凑、更像官方判断题）。
    if all(q.get("kind", "mc") == "tf" for q in qs):
        import math as _mtf
        tf_bottom = CONTENT_Y + CONTENT_H - 0.30
        tf_avail = tf_bottom - list_top
        n_tf = len(qs)
        tf_x = CONTENT_X + 0.55
        tf_w = CONTENT_W - 1.10

        def _tf_lines(pt: float) -> list:
            cpl = max(12, int(6.6 * tf_w * (12.5 / pt)))
            return [max(1, _mtf.ceil((len(str(q.get("q", ""))) + 8) / cpl)) for q in qs]

        tf_pt = 16.0
        for cand in (18.0, 17.0, 16.0, 15.0, 14.0):
            lh = cand / 72.0 * 1.18
            if sum(l * lh for l in _tf_lines(cand)) + (n_tf + 1) * 0.16 <= tf_avail:
                tf_pt = cand
                break
        lh = tf_pt / 72.0 * 1.18
        heights = [l * lh for l in _tf_lines(tf_pt)]
        slot_h = tf_avail / max(1, n_tf)
        for i, q in enumerate(qs):
            y = list_top + i * slot_h + max(0, (slot_h - heights[i]) / 2)
            tb = slide.shapes.add_textbox(
                Inches(tf_x), Inches(y), Inches(tf_w), Inches(heights[i] + 0.06))
            tfb = tb.text_frame
            tfb.margin_left = tfb.margin_right = 0
            tfb.margin_top = tfb.margin_bottom = 0
            tfb.word_wrap = True
            p = tfb.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.18
            r = p.add_run()
            r.text = f"{start_no + i}.  (    )  {capitalize_names(_clean_text(q.get('q', '')))}"
            r.font.name = FONT
            r.font.size = Pt(tf_pt)
            r.font.color.rgb = BLACK
        return

    list_bottom = CONTENT_Y + CONTENT_H - 0.18  # 贴近白卡底边，让题目区竖向铺满
    avail_h = list_bottom - list_top

    # 布局（用户拍板 2026-06-04/05/06）：横版 worksheet 阅读页【左右两列】，横竖对齐。
    # 偶数题左右等分（2+2 完全对称）；奇数题左列多一题（如 3 左 2 右，对齐官方参考样式），
    # 杜绝"左满右空"。印原文页（L3-6）答案在上方原文里，题目不再标 (P#)/Hint。
    n = len(qs)
    # 排版规则（用户拍板 2026-06-09）：≤3 题单栏横排、竖向左对齐（不左右分块）；
    # 仅 4 题及以上才分左右两块（2×2 对称）。杜绝 3 题时"左 2 右 1"的不平衡空隙。
    two_col = n >= 4 and not force_single_col
    gutter = 0.45
    if two_col:
        mid = (n + 1) // 2                 # 左列取上整：偶数→等分；奇数→左多一题
        col_qs = [qs[:mid], qs[mid:]]
        col_w = (CONTENT_W - 1.10 - gutter) / 2.0
        col_x = [CONTENT_X + 0.55, CONTENT_X + 0.55 + col_w + gutter]
    else:
        col_qs = [qs]
        col_w = CONTENT_W - 1.10
        col_x = [CONTENT_X + 0.55]

    import math as _m2

    # ============================================================
    #  阅读页标准化排版（用户拍板 2026-06-06）：
    #  ① 字号固定梯度：正文(read_pt) > 题干(stem_pt) > 选项(opt_pt)
    #  ② 行对齐 + 等距平铺：左右栏同一行【题干顶端基线齐平】；
    #     原文↔首行、行↔行、末行↔底栏 三段间距统一为同一模数 GAP
    #  ③ 缩进统一：所有题干首字共用列左基线；A/B/C 选项用固定悬挂缩进 OPT_INDENT，
    #     同栏选项首字符落在同一条垂直线上
    # ============================================================
    LS = 1.14                      # 全页统一行距
    OPT_INDENT = 0.26              # 选项悬挂缩进（英寸，固定值 → A/B/C 对齐同一竖线）

    def _cpl(width_in: float, pt: float, base: float = 6.8) -> int:
        return max(8, int(base * width_in * (12.5 / pt)))

    def _wrap(text_len: int, cpl: int) -> int:
        return max(1, _m2.ceil(text_len / max(1, cpl)))

    def _line_h(pt: float) -> float:
        return pt / 72.0 * LS + 0.028

    def _opt_pt_of(sp: float) -> float:
        return max(9.5, sp - 1.0)      # 选项比题干小 1pt（字号梯度）

    def _stem_lines(q: dict, sp: float) -> int:
        extra = 0 if show_passage else 7
        return _wrap(len(str(q.get("q", ""))) + extra, _cpl(col_w, sp))

    def _opt_line_counts(q: dict, op: float) -> list:
        cpl = _cpl(col_w - OPT_INDENT, op)
        kind = q.get("kind", "mc")
        if kind == "mc":
            return [_wrap(len(f"{chr(65 + j)}. {o}"), cpl)
                    for j, o in enumerate((q.get("options") or [])[:3])]
        return [1]                      # tf / short：一行作答
    SPACE_AFTER = 0.08                  # 题干↔选项、选项↔选项之间的段距（Bug4：略加大，缓解拥挤）

    def _q_height(q: dict, sp: float) -> float:
        op = _opt_pt_of(sp)
        h = _stem_lines(q, sp) * _line_h(sp) + SPACE_AFTER
        for n in _opt_line_counts(q, op):
            h += n * _line_h(op) + SPACE_AFTER
        # Bug3：已移除 Hint 行，不再为其预留行高。
        return h

    def _row_heights(sp: float) -> list:
        n_left = len(col_qs[0])
        n_right = len(col_qs[1]) if len(col_qs) > 1 else 0
        rows = max(n_left, n_right)
        rh = []
        for r in range(rows):
            hl = _q_height(col_qs[0][r], sp) if r < n_left else 0.0
            hr = _q_height(col_qs[1][r], sp) if (len(col_qs) > 1 and r < n_right) else 0.0
            rh.append(max(hl, hr))
        return rh

    def _fits(sp: float) -> bool:
        rh = _row_heights(sp)
        return sum(rh) + (len(rh) + 1) * 0.16 <= avail_h

    # 题干字号：固定梯度 → 上限 = 正文 - 1（保证 正文 > 题干），向下自适应直到放得下
    compact_mc_page = bool(
        show_passage
        and force_single_col
        and qs
        and all((q.get("kind", "mc") == "mc") for q in qs)
    )
    stem_min = 8.5 if compact_mc_page else 10.0
    min_keep = 4 if compact_mc_page else 2
    stem_cap = max(10.5, float(read_pt) - 1.0) if show_passage else 12.5
    stem_pt = stem_min
    cand = stem_cap
    while cand >= stem_min:
        if _fits(cand):
            stem_pt = cand
            break
        cand -= 0.5
    # Bug2：即便降到最小字号仍放不下（题太多/太长）→ 从末尾裁题直到放下，杜绝核心内容跑出页面。
    while not _fits(stem_pt) and sum(len(c) for c in col_qs) > min_keep:
        flat = [q for c in col_qs for q in c][:-1]
        if len(flat) >= 2 and two_col:
            mid = (len(flat) + 1) // 2
            col_qs = [flat[:mid], flat[mid:]]
        else:
            col_qs = [flat]
    opt_pt = _opt_pt_of(stem_pt)

    def _render_q(x: float, y: float, w: float, q: dict, qno: int) -> None:
        kind = q.get("kind", "mc")
        h = _q_height(q, stem_pt) + 0.05
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tfb = tb.text_frame
        tfb.vertical_anchor = MSO_ANCHOR.TOP        # 顶端对齐 → 左右栏同行题干基线齐平
        tfb.margin_left = tfb.margin_right = 0
        tfb.margin_top = tfb.margin_bottom = 0
        tfb.word_wrap = True

        def _opt_para(text: str, *, font=FONT, pt=opt_pt, color=BLACK, italic=False):
            """统一的选项/作答段：固定悬挂缩进 OPT_INDENT（marL/indent），A/B/C 首字对齐同一竖线。"""
            po = tfb.add_paragraph()
            po.alignment = PP_ALIGN.LEFT
            po.line_spacing = LS
            po.space_after = Pt(SPACE_AFTER * 72)
            _set_indent(po, OPT_INDENT)
            import re as _re
            for chunk in _re.split(r"(_{3,})", text):
                if not chunk:
                    continue
                ro = po.add_run()
                ro.text = chunk
                # 下划线段统一 Arial（连成一条实线，Poppins 下会断成虚线）
                ro.font.name = FONT_BLANK if chunk.startswith("___") else font
                ro.font.size = Pt(pt)
                ro.font.color.rgb = color
                ro.font.italic = italic

        page_no = None if show_passage else q.get("page")
        para = tfb.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = LS
        para.space_after = Pt(SPACE_AFTER * 72)
        stem_txt = f"{qno}. {_clean_text(q.get('q', ''))}"
        if page_no:
            stem_txt += f"  (P{page_no})"
        import re as _re_stem
        for _chunk in _re_stem.split(r"(_{3,})", stem_txt):
            if not _chunk:
                continue
            run = para.add_run()
            run.text = _chunk
            # 题干内填空下划线用 Arial（实线，不断成虚线）
            run.font.name = FONT_BLANK if _chunk.startswith("___") else FONT
            run.font.size = Pt(stem_pt)          # 题干字号（< 正文）
            run.font.color.rgb = BLACK
        if kind == "mc":
            for j, opt in enumerate((q.get("options") or [])[:3]):
                _opt_para(f"{chr(ord('A') + j)}. {_clean_text(opt)}")
        elif kind == "tf":
            _opt_para("Your answer (T / F): ____")
        else:  # short
            n_us = max(10, int((w - OPT_INDENT) / 0.12))
            _opt_para("_" * n_us, font=FONT_BLANK)
        # Bug3：去掉多余的「Hint: read page N again.」——题干已带 (P#) 定位，重复且占行。

    # —— 等距平铺（用户硬要求）：三段间距统一为同一模数 GAP —— 
    # GAP = (可用高 - 各行内容高之和) / (行数 + 1)；首行上边距 = 行间距 = 末行下边距 = GAP，
    # 全页留白标准化；同时左右栏【同一行共用同一 y】→ 题干顶端基线齐平、横竖对齐。
    n_left = len(col_qs[0])
    n_right = len(col_qs[1]) if len(col_qs) > 1 else 0
    rows = max(n_left, 1)
    row_h = _row_heights(stem_pt)
    gap = max(0.16, (avail_h - sum(row_h)) / (len(row_h) + 1))
    y = list_top + gap
    for r in range(len(row_h)):
        if r < n_left:
            _render_q(col_x[0], y, col_w, col_qs[0][r], start_no + r)
        if len(col_qs) > 1 and r < n_right:
            _render_q(col_x[1], y, col_w, col_qs[1][r], start_no + n_left + r)
        y += row_h[r] + gap


# ============================================================
#  Page 5 — Writing 脚手架
# ============================================================

def _build_p5_writing(slide, brand_rgb: tuple, writing: dict, title: str = "Reading") -> None:
    # 用户拍板：L4-6 第 2 阅读页 = 写作/PBL，但页面大标题仍叫 Reading（写作内嵌）
    # v2.2：写作任务必须【基于本文故事】（复述/分析），不再写空泛抽象题；
    #   副标题改成明确的复述指令，让孩子读完原文就知道怎么写。
    _default_sub = ("Plan with the chart, then write your own short story."
                    if title == "Writing"
                    else "Read the story again, then retell it in your own words.")
    subtitle = writing.get("subtitle") or _default_sub
    _add_title(slide, title, subtitle)

    # 中部黄虚线框 = 5 步骨架
    scaff_top = CONTENT_Y + 1.40
    scaff_h = 2.40
    scaff = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 1.20), Inches(scaff_top),
        Inches(CONTENT_W - 2.40), Inches(scaff_h),
    )
    scaff.adjustments[0] = 0.04
    scaff.fill.solid()
    scaff.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)  # 浅黄
    scaff.line.color.rgb = RGBColor(0xE0, 0xC8, 0x80)
    scaff.line.width = Pt(1.5)
    scaff.shadow.inherit = False

    tf = scaff.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True

    # Title 行
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "Title: "
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x6B, 0x4D, 0xA8)  # 紫
    r2 = p.add_run()
    _wt = writing.get("title", "")
    r2.text = _wt or "_______________"
    r2.font.name = FONT if _wt else FONT_BLANK  # 空标题占位线用 Arial 实线
    r2.font.bold = True
    r2.font.size = Pt(14)
    r2.font.color.rgb = BLACK

    # 5 步
    step_colors = [
        RGBColor(0x4C, 0xA1, 0x65),  # Beginning - 绿
        RGBColor(0x4C, 0x8F, 0xD8),  # First event - 蓝
        RGBColor(0xE3, 0x76, 0x35),  # Second event - 橙
        RGBColor(0xC0, 0x39, 0x6F),  # Funny event - 红
        RGBColor(0x6B, 0x4D, 0xA8),  # Ending - 紫
    ]
    step_labels = writing.get("step_labels") or [
        "Beginning:", "First event:", "Second event:", "Funny event:", "Ending:",
    ]
    step_contents = writing.get("steps") or [""] * 5

    for i in range(5):
        pi = tf.add_paragraph()
        pi.alignment = PP_ALIGN.LEFT
        pi.space_after = Pt(2)
        rn = pi.add_run()
        rn.text = f"{i + 1}. "
        rn.font.name = FONT
        rn.font.bold = True
        rn.font.size = Pt(13)
        rn.font.color.rgb = BLACK
        rl = pi.add_run()
        rl.text = step_labels[i] + "  "
        rl.font.name = FONT
        rl.font.bold = True
        rl.font.size = Pt(13)
        rl.font.color.rgb = step_colors[i % len(step_colors)]
        # step content 里含 ________ 时切出来用 Arial 字体
        _emit_with_underscore_lock(
            pi, step_contents[i] if i < len(step_contents) else "", 13, BLACK
        )

    # 字数提示
    hint_y = scaff_top + scaff_h + 0.10
    hb = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(hint_y),
        Inches(CONTENT_W), Inches(0.25),
    )
    p = hb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"Write {writing.get('min_words', 50)}-{writing.get('max_words', 80)} words."
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0xE6, 0xA8, 0x2B)

    # 蓝色横线写作区
    write_top = hint_y + 0.35
    write_bottom = CONTENT_Y + CONTENT_H - 0.30
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.40), Inches(write_top),
        Inches(CONTENT_W - 0.80), Inches(write_bottom - write_top),
    )
    box.adjustments[0] = 0.03
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0x8C, 0xB6, 0xE6)
    box.line.width = Pt(1.2)
    box.shadow.inherit = False

    # 内部横线（5-6 条）
    line_left = CONTENT_X + 0.70
    line_right = CONTENT_X + CONTENT_W - 0.70
    n_lines = 5
    line_step = (write_bottom - write_top - 0.30) / n_lines
    for k in range(1, n_lines + 1):
        ly = write_top + k * line_step
        ln = slide.shapes.add_connector(
            1,
            Inches(line_left), Inches(ly),
            Inches(line_right), Inches(ly),
        )
        ln.line.color.rgb = RGBColor(0xB7, 0xCD, 0xEC)
        ln.line.width = Pt(0.75)


# ============================================================
#  Page 6 (可选) — PBL 迷你项目（标题统一 Reading）
# ============================================================

def _build_pbl_page(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                    title: str = "Reading") -> None:
    """读后 PBL 迷你项目：项目目标 + 3 步骤脚手架 + 大块创作区 + 几条书写线。

    全程纯纸面、可独立完成；标题统一 Reading（大结构不变）。
    """
    theme = (outline.theme or "the story").strip()
    pbl = (data.get("pbl") if isinstance(data, dict) else None) or {}
    project = pbl.get("project") or f"Make a poster about {theme}."
    steps = pbl.get("steps") or [
        f"Draw the most important part of {theme}.",
        "Add labels for what you drew.",
        "Write one sentence about your picture.",
    ]
    _add_title(slide, title, f"Mini Project: {project}")

    # 步骤脚手架框（浅色底 + 主色边）
    brief_top = CONTENT_Y + 1.40
    brief_h = 1.75
    brief = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.45), Inches(brief_top),
        Inches(CONTENT_W - 0.90), Inches(brief_h),
    )
    brief.adjustments[0] = 0.05
    brief.fill.solid()
    brief.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)
    brief.line.color.rgb = RGBColor(*brand_rgb)
    brief.line.width = Pt(1.5)
    brief.shadow.inherit = False
    tf = brief.text_frame
    tf.margin_left = tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.space_after = Pt(4)
    r0 = p0.add_run()
    r0.text = "Steps:"
    r0.font.name = FONT
    r0.font.bold = True
    r0.font.size = Pt(15)
    r0.font.color.rgb = RGBColor(*brand_rgb)
    for i, step in enumerate(steps[:3]):
        pi = tf.add_paragraph()
        pi.alignment = PP_ALIGN.LEFT
        pi.space_after = Pt(3)
        ri = pi.add_run()
        ri.text = f"{i + 1}. {capitalize_names(str(step))}"
        ri.font.name = FONT
        ri.font.size = Pt(14)
        ri.font.color.rgb = BLACK

    # 大块创作区（白底 + 主色边）
    canvas_top = brief_top + brief_h + 0.25
    canvas_bottom = CONTENT_Y + CONTENT_H - 0.95
    canvas = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.45), Inches(canvas_top),
        Inches(CONTENT_W - 0.90), Inches(canvas_bottom - canvas_top),
    )
    canvas.adjustments[0] = 0.03
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = WHITE
    canvas.line.color.rgb = RGBColor(*brand_rgb)
    canvas.line.width = Pt(1.2)
    canvas.shadow.inherit = False
    ctf = canvas.text_frame
    ctf.margin_left = ctf.margin_right = Inches(0.20)
    ctf.margin_top = Inches(0.10)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    cr = cp.add_run()
    cr.text = "Draw and make it here:"
    cr.font.name = FONT
    cr.font.italic = True
    cr.font.size = Pt(12)
    cr.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)

    # 底部 2 条书写线
    line_top = canvas_bottom + 0.30
    line_left = CONTENT_X + 0.70
    line_right = CONTENT_X + CONTENT_W - 0.70
    for k in range(2):
        ly = line_top + k * 0.32
        ln = slide.shapes.add_connector(
            1, Inches(line_left), Inches(ly), Inches(line_right), Inches(ly),
        )
        ln.line.color.rgb = LIGHT_GRAY
        ln.line.width = Pt(0.9)


# ============================================================
#  Page 6 (L0-2) — 涂色线稿 + 一句话（标题统一 Reading）
# ============================================================

def _color_say_stem(outline: BookOutline) -> str:
    """从故事里取一句最简单的第一人称句作填空母句，挖掉末词（贴本书内容/已学知识）。

    例：'I have a pair of red shoes.' → 'I have a pair of red'（末词 shoes 留空让孩子填）。
    没有第一人称句则回退 'I like'。
    """
    import re as _re2
    for pg in (getattr(outline, "pages", []) or []):
        t = (getattr(pg, "text", "") or "").strip()
        if not t or getattr(pg, "page_type", "") != "story":
            continue
        if _re2.match(r"(?i)^i\b", t):
            words = _re2.sub(r"[.!?]+$", "", t).split()
            if len(words) >= 2:
                return " ".join(words[:-1])
    return "I like"


def _build_l3_summary_page(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                           *, title: str = "Reading") -> None:
    """L3 第②阅读页（对齐用户标杆 L3-30）：引导式小结——
    L3-4(A1+) 不做长篇自由写作（太难），改成【补全句子】的引导式小结：
    副标题下方先放一条【词库】（置顶，给孩子选词支撑，对齐标杆版式）；中部 3 条贴本书的
    句子框（句首引导 + Arial 实线作答空 + 整宽书写线）。整体留白均衡、适合低龄独立完成。"""
    _add_title(slide, title,
               "Read the book again. Finish each sentence about the story.")

    is_nonfic = "non" in (getattr(outline, "fiction_type", "") or "").lower()
    if is_nonfic:
        starters = ["This book is about", "One thing I learned is", "I can"]
    else:
        starters = ["This story is about", "First,", "In the end,"]

    # —— 词库置顶（用户标杆 L3-30）：副标题正下方一条词库框，供孩子选词补全 —— 
    words = [str(w).strip() for w in (outline.vocabulary_for_display or []) if str(w).strip()][:6]
    prompts_top = CONTENT_Y + 1.55
    if words:
        wb_top = CONTENT_Y + 1.50
        wb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(CONTENT_X + 0.90), Inches(wb_top), Inches(CONTENT_W - 1.80), Inches(0.60))
        wb.adjustments[0] = 0.30
        wb.fill.solid()
        wb.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)
        wb.line.color.rgb = RGBColor(*brand_rgb)
        wb.line.width = Pt(1.0)
        wb.shadow.inherit = False
        wtf = wb.text_frame
        wtf.margin_left = wtf.margin_right = Inches(0.18)
        wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        wtf.word_wrap = True
        wp = wtf.paragraphs[0]
        wp.alignment = PP_ALIGN.CENTER
        wl = wp.add_run()
        wl.text = "Word bank:  "
        wl.font.name = FONT
        wl.font.bold = True
        wl.font.size = Pt(18)
        wl.font.color.rgb = RGBColor(*brand_rgb)
        ww = wp.add_run()
        ww.text = "    ".join(words)
        ww.font.name = FONT
        ww.font.size = Pt(20)
        ww.font.color.rgb = BLACK
        prompts_top = wb_top + 0.95

    # —— 句子补全（词库下方铺满整页）——
    area_top = prompts_top
    area_bottom = CONTENT_Y + CONTENT_H - 0.35
    x = CONTENT_X + 0.90
    box_w = CONTENT_W - 1.80
    n = len(starters)
    slot_h = (area_bottom - area_top) / n
    for i, st in enumerate(starters):
        y = area_top + i * slot_h
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(box_w), Inches(0.50))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _emit_with_underscore_lock(p, f"{i + 1}. {st} ____________", 22, BLACK, bold=True)
        # 整宽书写线（多写一行的空间）
        _draw_writing_line(slide, x + 0.30, y + slot_h - 0.34, box_w - 0.30, LIGHT_GRAY, 1.3)


def _timeline_events(outline: BookOutline, max_n: int = 4) -> list[str]:
    sents = _story_sentences_for_grammar(outline)
    cleaned: list[str] = []
    seen: set[str] = set()
    for s in sents:
        text = capitalize_names(_clean_text(s)).strip()
        if not text:
            continue
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        if len(text) > 72:
            words = text.split()
            text = " ".join(words[:10]).rstrip(",.;:") + "."
        cleaned.append(text)
    if len(cleaned) <= max_n:
        return cleaned[:max_n]
    idxs = [0, max(1, len(cleaned) // 3), max(2, (len(cleaned) * 2) // 3), len(cleaned) - 1]
    out: list[str] = []
    for idx in idxs:
        if cleaned[idx] not in out:
            out.append(cleaned[idx])
    return out[:max_n]


def _build_timeline_page(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                         *, title: str = "Reading") -> None:
    """SOP-style GO page for sequence / timeline books.

    A1/A1+ support: four nodes, first and last mostly prefilled, middle nodes
    write-in. This keeps the worksheet printable and avoids turning every book
    into a generic comprehension-question page.
    """
    _add_title(slide, title, "Fill in the timeline. Write what happened in order.")
    events = _timeline_events(outline, 4)
    while len(events) < 4:
        events.append("")

    labels = ["First", "Next", "Then", "Finally"]
    bc = RGBColor(*brand_rgb)
    area_top = CONTENT_Y + 1.55
    card_w = 2.05
    card_h = 2.18
    gap = (CONTENT_W - 0.90 - card_w * 4) / 3
    x0 = CONTENT_X + 0.45
    y = area_top + 0.48
    fill_rgb = RGBColor(0xFF, 0xF7, 0xE8)
    prompt_rgb = RGBColor(0x6B, 0x72, 0x80)

    for i, label in enumerate(labels):
        x = x0 + i * (card_w + gap)
        if i:
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x - gap + 0.12), Inches(y + card_h / 2),
                Inches(x - 0.12), Inches(y + card_h / 2),
            )
            ln.line.color.rgb = bc
            ln.line.width = Pt(2.0)
            ln.shadow.inherit = False

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.adjustments[0] = 0.12
        card.fill.solid()
        card.fill.fore_color.rgb = fill_rgb
        card.line.color.rgb = bc
        card.line.width = Pt(1.5)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.12)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(18)
        r.font.color.rgb = bc

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(6)
        if i in (0, 3) and events[i]:
            txt = events[i]
        else:
            txt = "Look back at the story.\nWrite one short sentence."
        r2 = p2.add_run()
        r2.text = txt
        r2.font.name = FONT
        r2.font.size = Pt(12.5 if len(txt) > 52 else 14)
        r2.font.color.rgb = prompt_rgb if i not in (0, 3) else BLACK

        if i not in (0, 3):
            _draw_writing_line(slide, x + 0.18, y + card_h - 0.58, card_w - 0.36, LIGHT_GRAY, 1.2)
            _draw_writing_line(slide, x + 0.18, y + card_h - 0.34, card_w - 0.36, LIGHT_GRAY, 1.2)

    words = [str(w).strip() for w in (outline.vocabulary_for_display or data.get("word_bank") or []) if str(w).strip()]
    if words:
        wb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(CONTENT_X + 1.15), Inches(CONTENT_Y + CONTENT_H - 0.92),
            Inches(CONTENT_W - 2.30), Inches(0.46)
        )
        wb.adjustments[0] = 0.30
        wb.fill.solid()
        wb.fill.fore_color.rgb = WHITE
        wb.line.color.rgb = bc
        wb.line.width = Pt(1.0)
        wb.shadow.inherit = False
        wtf = wb.text_frame
        wtf.margin_left = wtf.margin_right = Inches(0.16)
        wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        wp = wtf.paragraphs[0]
        wp.alignment = PP_ALIGN.CENTER
        rw = wp.add_run()
        rw.text = "Word bank: " + "   ".join(words[:8])
        rw.font.name = FONT
        rw.font.size = Pt(13.5)
        rw.font.color.rgb = BLACK


def _plan_chart_rows(outline: BookOutline) -> list[dict]:
    story = " ".join(_story_sentences_for_grammar(outline)).lower()
    rows = [
        {"clue": "homework", "action": "do homework", "time": "first"},
        {"clue": "room", "action": "clean her room", "time": "on Tuesday"},
        {"clue": "piano", "action": "practice the piano", "time": "every day"},
        {"clue": "piano", "action": "play for one hour", "time": "every day"},
    ]
    if not any(k in story for k in ("homework", "room", "piano", "sunday")):
        sents = [capitalize_names(_clean_text(s)) for s in _story_sentences_for_grammar(outline)]
        simple = [s for s in sents if 12 <= len(s) <= 74][:4]
        rows = [{"clue": f"part {i + 1}", "action": s, "time": ""} for i, s in enumerate(simple)]
    while len(rows) < 4:
        rows.append({"clue": f"part {len(rows) + 1}", "action": "", "time": ""})
    return rows[:4]


def _build_plan_chart_page(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                           *, title: str = "Reading") -> None:
    """A1-supported GO for plan/schedule stories.

    This differs from a pure timeline: students match what the character will do
    with when it happens, which fits plan stories such as Mia better.
    """
    _add_title(slide, title, "Complete the plan chart. Use the clues and word bank to help.")
    bc = RGBColor(*brand_rgb)
    rows = _plan_chart_rows(outline)

    x0 = CONTENT_X + 0.45
    y0 = CONTENT_Y + 1.38
    col_w = [2.05, 4.10, 2.35]
    row_h = 0.74
    head_h = 0.48
    headers = ["Clue", "What will Mia do?", "Time / Order"]
    fill_head = RGBColor(0xE7, 0xF0, 0xFF)
    fill_cell = RGBColor(0xFF, 0xFB, 0xF2)
    muted = RGBColor(0x6B, 0x72, 0x80)

    x = x0
    for i, h in enumerate(headers):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y0),
            Inches(col_w[i]), Inches(head_h)
        )
        sh.adjustments[0] = 0.12
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_head
        sh.line.color.rgb = bc
        sh.line.width = Pt(1.1)
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = Inches(0.10)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = bc
        x += col_w[i]

    for r_idx, row in enumerate(rows):
        y = y0 + head_h + r_idx * row_h
        x = x0
        values = [
            row.get("clue", ""),
            row.get("action", "") if r_idx in (0, 2) else "",
            row.get("time", "") if r_idx in (1, 3) else "",
        ]
        placeholders = ["look", "write the action", "write the time"]
        for c_idx, value in enumerate(values):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                Inches(col_w[c_idx]), Inches(row_h)
            )
            cell.adjustments[0] = 0.10
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_cell if c_idx else WHITE
            cell.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
            cell.line.width = Pt(1.0)
            cell.shadow.inherit = False
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = capitalize_names(value) if value else placeholders[c_idx]
            run.font.name = FONT
            run.font.size = Pt(12.5 if len(run.text) > 36 else 14)
            run.font.color.rgb = BLACK if value else muted
            if not value and c_idx:
                _draw_writing_line(slide, x + 0.18, y + row_h - 0.18, col_w[c_idx] - 0.36, LIGHT_GRAY, 1.1)
            x += col_w[c_idx]

    bank = ["first", "clean her room", "every day", "play for one hour"]
    if not any(k in " ".join((r.get("clue", "") for r in rows)).lower() for k in ("homework", "piano")):
        bank = [str(w).strip() for w in (outline.vocabulary_for_display or data.get("word_bank") or []) if str(w).strip()][:6]
    if bank:
        wb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(CONTENT_X + 0.38), Inches(CONTENT_Y + CONTENT_H - 1.02),
            Inches(CONTENT_W - 0.76), Inches(0.68)
        )
        wb.adjustments[0] = 0.30
        wb.fill.solid()
        wb.fill.fore_color.rgb = WHITE
        wb.line.color.rgb = bc
        wb.line.width = Pt(1.0)
        wb.shadow.inherit = False
        tf = wb.text_frame
        tf.margin_left = tf.margin_right = Inches(0.16)
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        rr = p.add_run()
        rr.text = "Word bank: " + "   ".join(bank)
        rr.font.name = FONT
        rr.font.size = Pt(14.0)
        rr.font.color.rgb = BLACK


def _build_l3_bubble_map(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                         *, title: str = "Reading") -> None:
    """L3 非虚构第②阅读页（重新设计 2026-06-09）：代码原生【关键词词网 / Word Web】。

    问题诊断（旧版"写出学到的 4 件事 + 空白气泡 + 底部词库"）：
      ① 题干空泛、作答无支撑——4 个只标 1./2./3./4. 的空气泡，A1+ 孩子要凭空写 4 条
         "学到的事"太抽象、太难；
      ② 词库与题干不匹配——底部 3 个词汇名词 vs 4 个开放气泡，词不够、且名词无法直接
         拼成"学到的事"，孩子不知怎么用；
      ③ 与正文脱节——气泡只有序号，不指向本书任何具体内容；
      ④ 与"小结页/前面词汇页"概念重复。
    新设计：把每个分支气泡【预填一个本书关键词】（作答支撑就在气泡里），孩子读后
      针对每个给定词写【一句话/一个想法】→ 题干清晰、有支撑、紧扣正文；中心 = 主旨。
      关键词已在气泡内 → 去掉底部那条不匹配的词库，版面更干净。
    """
    bc = RGBColor(*brand_rgb)

    # —— 分支标签：优先本书关键词（给定即支撑）；词不足时退回 4 条可答的引导提示 —— 
    seen: set = set()
    kw: list[str] = []
    for w in (outline.vocabulary_for_display or []):
        w = str(w).strip()
        lw = w.lower()
        if w and lw not in seen:
            seen.add(lw)
            kw.append(w)
    if len(kw) >= 3:
        labels = kw[:4]
        word_mode = True
        _add_title(slide, title,
                   "Read the book again. Write one idea about each key word.")
    else:
        labels = ["A new word I learned", "One important fact",
                  "Something I can do", "I want to learn more about"]
        word_mode = False
        _add_title(slide, title,
                   "Read the book again. Finish each bubble about the book.")
    n = min(len(labels), 4)

    # 中心主旨：主题→书名，过滤非英文/过长，否则用 Main Idea
    center_label = (getattr(outline, "theme", "") or "").strip()
    if not center_label or not center_label.isascii():
        center_label = (outline.title or "").strip()
    if not center_label or not center_label.isascii():
        center_label = "Main Idea"
    if len(center_label) > 22:
        center_label = center_label[:22].rstrip() + "…"

    area_top = CONTENT_Y + 1.45
    area_bottom = CONTENT_Y + CONTENT_H - 0.40   # 词库已去掉 → 气泡区可铺满
    area_h = area_bottom - area_top
    cx = CONTENT_X + CONTENT_W / 2.0
    cy = area_top + area_h / 2.0
    cw, ch = 2.95, 1.10   # 中心气泡

    bw, bh = 3.05, 1.20
    margin_x = 0.22
    left_x = CONTENT_X + margin_x
    right_x = CONTENT_X + CONTENT_W - margin_x - bw
    center_x = cx - bw / 2.0
    top_y = area_top
    bot_y = area_bottom - bh
    if n == 3:
        # 3 个词 → 上 2 下 1（底部居中），杜绝右下角空缺导致的不平衡
        spots = [(left_x, top_y), (right_x, top_y), (center_x, bot_y)]
    else:
        spots = [(left_x, top_y), (right_x, top_y), (left_x, bot_y), (right_x, bot_y)]

    for i in range(n):
        bx, by = spots[i]
        label = labels[i]
        # 连接线：中心 → 气泡中心
        bcx, bcy = bx + bw / 2, by + bh / 2
        ln = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(cy), Inches(bcx), Inches(bcy))
        ln.line.color.rgb = bc
        ln.line.width = Pt(1.4)
        ln.shadow.inherit = False
        # 气泡
        bub = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by), Inches(bw), Inches(bh))
        bub.adjustments[0] = 0.16
        bub.fill.solid()
        bub.fill.fore_color.rgb = WHITE
        bub.line.color.rgb = bc
        bub.line.width = Pt(1.6)
        bub.shadow.inherit = False
        btf = bub.text_frame
        btf.margin_left = btf.margin_right = Inches(0.16)
        btf.margin_top = Inches(0.10)
        btf.vertical_anchor = MSO_ANCHOR.TOP
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        # 序号 + 给定关键词（加粗品牌色 = 作答支撑），或引导提示
        rn = bp.add_run()
        rn.text = f"{i + 1}. "
        rn.font.name = FONT
        rn.font.bold = True
        rn.font.size = Pt(14 if word_mode else 12)
        rn.font.color.rgb = bc
        rl = bp.add_run()
        rl.text = capitalize_names(label) if not word_mode else label
        rl.font.name = FONT
        rl.font.bold = True
        rl.font.size = Pt(15 if word_mode else 12)
        rl.font.color.rgb = bc if word_mode else BLACK
        # 一条作答横线（写一句话/一个想法）
        _draw_writing_line(slide, bx + 0.34, by + bh - 0.34, bw - 0.60, LIGHT_GRAY, 1.2)

    # 中心气泡（最后绘制 → 盖在连线之上，避免连线穿过中心文字）
    center = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - cw / 2), Inches(cy - ch / 2), Inches(cw), Inches(ch))
    center.adjustments[0] = 0.5
    center.fill.solid()
    center.fill.fore_color.rgb = bc
    center.line.fill.background()
    center.shadow.inherit = False
    ctf = center.text_frame
    ctf.margin_left = ctf.margin_right = Inches(0.10)
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = center_label
    cr.font.name = FONT
    cr.font.bold = True
    cr.font.size = Pt(17)
    cr.font.color.rgb = _readable_text_rgb(bc)


def _l34_clean_story_bits(outline: BookOutline, max_n: int = 4) -> list[str]:
    bits: list[str] = []
    seen: set[str] = set()
    for sent in _story_sentences_for_grammar(outline):
        text = capitalize_names(_clean_text(sent)).strip()
        if not text:
            continue
        if len(text) > 76:
            words = text.split()
            text = " ".join(words[:11]).rstrip(",.;:") + "."
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        bits.append(text)
    return bits[:max_n]


def _l34_short_topic(outline: BookOutline) -> str:
    for value in (getattr(outline, "theme", ""), getattr(outline, "title", "")):
        text = _clean_text(str(value or "")).strip()
        if text and text.isascii():
            return text[:34]
    return "the story"


def _build_l34_graphic_organizer_page(slide, brand_rgb: tuple, data: dict,
                                      outline: BookOutline, go_mode: str) -> None:
    """L3/L4 SOP GO page: clear chart, strong support, only 3-6 blanks."""
    _add_title(slide, "Graphic Organizer", "Fill in the graphic organizer.")
    bc = RGBColor(*brand_rgb)
    soft = RGBColor(0xFF, 0xF7, 0xE8)
    pale = RGBColor(0xF7, 0xFB, 0xFF)
    muted = RGBColor(0x6B, 0x72, 0x80)
    is_nonfic = "non" in (getattr(outline, "fiction_type", "") or "").lower()
    mode = (go_mode or "").lower()
    go_code = _l34_go_activity_code(outline, go_mode)
    explicit_go = (
        str(getattr(outline, "graphic_organizer", "") or "") + " "
        + str(getattr(outline, "graphic_organizer_desc", "") or "")
    ).lower()
    words = [str(w).strip() for w in (outline.vocabulary_for_display or data.get("word_bank") or []) if str(w).strip()]
    story_bits = _l34_clean_story_bits(outline, 4)

    if go_code == "go_compare_chart":
        center = "Compare"
        labels = ["Same", "Different 1", "Different 2", "Evidence"]
        prompts = ["", "", "", ""]
        blank_slots = {0, 1, 2, 3}
    elif go_code == "go_fact_web":
        center = "Facts"
        if any(k in explicit_go for k in ("classification", "classify", "category", "categories", "habitat")):
            labels = ["Topic", "Group 1", "Group 2", "Example"]
        elif any(k in explicit_go for k in ("main idea", "details")):
            labels = ["Main Idea", "Detail 1", "Detail 2", "Example"]
        else:
            labels = ["Topic", "Fact 1", "Fact 2", "Example"]
        prompts = [
            _l34_short_topic(outline),
            "",
            "",
            "",
        ]
        blank_slots = {1, 2, 3}
    elif mode == "planchart" or ("plan" in explicit_go and go_code == "go_sequence_chart"):
        center = "Plan"
        labels = ["First", "Next", "Time", "Result"]
        rows = _plan_chart_rows(outline)
        prompts = [
            rows[0].get("action", "") or (story_bits[0] if story_bits else ""),
            "",
            "",
            "",
        ]
        blank_slots = {1, 2, 3}
        bank = ["first", "clean her room", "every day", "play for one hour"]
        if any(k in " ".join((r.get("clue", "") for r in rows)).lower() for k in ("homework", "piano")):
            words = bank
    elif go_code == "go_sequence_chart":
        center = "Sequence"
        labels = ["First", "Next", "Then", "Finally"]
        prompts = [
            story_bits[0] if story_bits else "",
            "",
            "",
            "",
        ]
        blank_slots = {1, 2, 3}
    else:
        center = "Story"
        if "solution" in explicit_go:
            labels = ["Problem", "Action", "Solution", "Result"]
        else:
            labels = ["Beginning", "Middle", "Action", "End"]
        prompts = [
            story_bits[0] if story_bits else "",
            "",
            "",
            "",
        ]
        blank_slots = {1, 2, 3}

    x0 = CONTENT_X + 0.62
    y0 = CONTENT_Y + 1.42
    card_w = 2.05
    card_h = 2.28
    gap = (CONTENT_W - 1.24 - card_w * 4) / 3
    for i, label in enumerate(labels):
        x = x0 + i * (card_w + gap)
        if i:
            ln = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x - gap + 0.15), Inches(y0 + 1.10),
                Inches(x - 0.15), Inches(y0 + 1.10),
            )
            ln.line.color.rgb = bc
            ln.line.width = Pt(1.6)
            ln.shadow.inherit = False
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y0), Inches(card_w), Inches(card_h)
        )
        card.adjustments[0] = 0.12
        card.fill.solid()
        card.fill.fore_color.rgb = soft if i % 2 == 0 else pale
        card.line.color.rgb = bc
        card.line.width = Pt(1.2)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.margin_left = tf.margin_right = Inches(0.13)
        tf.margin_top = Inches(0.10)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(15)
        r.font.color.rgb = bc
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(8)
        p2.text = prompts[i] if i < len(prompts) and prompts[i] else ""
        for rr in p2.runs:
            rr.font.name = FONT
            rr.font.size = Pt(12.5)
            rr.font.color.rgb = muted if i in blank_slots else BLACK
        if i in blank_slots:
            _draw_writing_line(slide, x + 0.22, y0 + card_h - 0.64, card_w - 0.44, LIGHT_GRAY, 1.2)
            _draw_writing_line(slide, x + 0.22, y0 + card_h - 0.36, card_w - 0.44, LIGHT_GRAY, 1.2)

    center_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 3.80), Inches(CONTENT_Y + 4.08),
        Inches(2.74), Inches(0.58),
    )
    center_box.adjustments[0] = 0.30
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = bc
    center_box.line.fill.background()
    center_box.shadow.inherit = False
    ctf = center_box.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = center
    cr.font.name = FONT
    cr.font.bold = True
    cr.font.size = Pt(17)
    cr.font.color.rgb = _readable_text_rgb(bc)

    if words:
        wb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(CONTENT_X + 0.80), Inches(CONTENT_Y + CONTENT_H - 0.96),
            Inches(CONTENT_W - 1.60), Inches(0.58),
        )
        wb.adjustments[0] = 0.26
        wb.fill.solid()
        wb.fill.fore_color.rgb = WHITE
        wb.line.color.rgb = bc
        wb.line.width = Pt(1.0)
        wb.shadow.inherit = False
        tf = wb.text_frame
        tf.margin_left = tf.margin_right = Inches(0.16)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Word bank: " + "   ".join(words[:6])
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.color.rgb = BLACK


def _build_l34_writing_page(slide, brand_rgb: tuple, outline: BookOutline) -> None:
    """L3/L4 SOP writing page: starters on top, independent writing lines below."""
    topic = _l34_short_topic(outline)
    _add_title(slide, "Writing", f"Write about {topic}.")
    bc = RGBColor(*brand_rgb)
    lvl = _level_num(getattr(outline, "level", "") or "")
    example = _display_sentence_frame(outline)
    sig = _frame_signature(example)
    is_nonfic = "non" in (getattr(outline, "fiction_type", "") or "").lower()
    go_mode = str(getattr(outline, "_worksheet_second_reading_mode", "") or "")
    go_code = _l34_go_activity_code(outline, go_mode)
    explicit_go = (
        str(getattr(outline, "graphic_organizer", "") or "") + " "
        + str(getattr(outline, "graphic_organizer_desc", "") or "")
    ).lower()
    if go_code == "go_compare_chart":
        starters = [
            f"They are the same because {ANSWER_BLANK}.",
            f"One difference is {ANSWER_BLANK}.",
            f"The text shows {ANSWER_BLANK}.",
        ]
    elif go_code == "go_fact_web":
        if any(k in explicit_go for k in ("classification", "classify", "category", "categories", "habitat")):
            starters = [
                f"The topic is {ANSWER_BLANK}.",
                f"One group is {ANSWER_BLANK}.",
                f"Another group is {ANSWER_BLANK}.",
            ]
        else:
            starters = [
                f"The topic is {ANSWER_BLANK}.",
                f"One fact is {ANSWER_BLANK}.",
                f"Another fact is {ANSWER_BLANK}.",
            ]
    elif go_code == "go_sequence_chart":
        starters = [
            f"First, {ANSWER_BLANK}.",
            f"Next, {ANSWER_BLANK}.",
            f"At the end, {ANSWER_BLANK}.",
        ]
    elif is_nonfic:
        starters = [
            f"The topic is {ANSWER_BLANK}.",
            f"One fact is {ANSWER_BLANK}.",
            f"Another fact is {ANSWER_BLANK}.",
        ]
    elif sig == "will":
        starters = [
            f"I will {ANSWER_BLANK}.",
            f"I will {ANSWER_BLANK} first.",
            f"I will {ANSWER_BLANK} every day.",
        ]
    elif sig == "because":
        starters = [
            f"I think {ANSWER_BLANK}.",
            f"It is {ANSWER_BLANK} because {ANSWER_BLANK}.",
            f"I feel {ANSWER_BLANK} because {ANSWER_BLANK}.",
        ]
    elif sig == "there":
        starters = [
            f"There is {ANSWER_BLANK}.",
            f"There are {ANSWER_BLANK}.",
            f"I can see {ANSWER_BLANK}.",
        ]
    else:
        starters = [
            f"First, {ANSWER_BLANK}.",
            f"Then, {ANSWER_BLANK}.",
            f"At the end, {ANSWER_BLANK}.",
        ]
    if lvl >= 4 and not is_nonfic and sig not in {"will", "there"}:
        starters = starters[:2] + [f"At the end, {ANSWER_BLANK} because {ANSWER_BLANK}."]

    top = CONTENT_Y + 1.36
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.70), Inches(top),
        Inches(CONTENT_W - 1.40), Inches(2.18),
    )
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE8)
    box.line.color.rgb = bc
    box.line.width = Pt(1.2)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.14)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Use these sentence starters."
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = bc
    if example:
        p_ex = tf.add_paragraph()
        p_ex.space_before = Pt(2)
        p_ex.alignment = PP_ALIGN.LEFT
        rr = p_ex.add_run()
        rr.text = f"Example: {example}"
        rr.font.name = FONT
        rr.font.size = Pt(12.5)
        rr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    for i, starter in enumerate(starters[:3]):
        pi = tf.add_paragraph()
        pi.space_before = Pt(5)
        pi.alignment = PP_ALIGN.LEFT
        _emit_with_underscore_lock(pi, f"{i + 1}. {starter}", 15, BLACK, bold=False)

    hint = slide.shapes.add_textbox(
        Inches(CONTENT_X + 0.78), Inches(top + 2.46),
        Inches(CONTENT_W - 1.56), Inches(0.34),
    )
    hp = hint.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run()
    hr.text = "Write 2-3 sentences."
    hr.font.name = FONT
    hr.font.size = Pt(14)
    hr.font.bold = True
    hr.font.color.rgb = bc

    write_top = top + 2.92
    write_bottom = CONTENT_Y + CONTENT_H - 0.44
    write_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(CONTENT_X + 0.70), Inches(write_top),
        Inches(CONTENT_W - 1.40), Inches(write_bottom - write_top),
    )
    write_box.adjustments[0] = 0.04
    write_box.fill.solid()
    write_box.fill.fore_color.rgb = WHITE
    write_box.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    write_box.line.width = Pt(1.0)
    write_box.shadow.inherit = False
    n_lines = 4 if lvl <= 3 else 5
    line_left = CONTENT_X + 1.02
    line_w = CONTENT_W - 2.04
    usable_h = write_bottom - write_top - 0.40
    step = usable_h / max(n_lines, 1)
    for i in range(n_lines):
        y = write_top + 0.34 + (i + 1) * step
        _draw_writing_line(slide, line_left, y, line_w, RGBColor(0xAE, 0xBD, 0xD0), 1.1)


def _build_color_say_page(slide, brand_rgb: tuple, data: dict, outline: BookOutline,
                          *, coloring_image: Optional[Path] = None,
                          title: str = "Reading") -> None:
    """最低龄读后页（用户拍板 2026-06-08，替代 L0-2 偏难的 PBL；2026-06-09 简化居中）：

    上半：一张【线稿涂色图】(batch 按故事第一句生成；缺省则留白画框，孩子自己画)；
    下半：用本书第一人称句做一句话填空（贴已学知识，如 pairs），词库辅助。
    版式整体居中、尽量简单，便于低龄学生独立完成；标题统一 Reading。
    """
    _add_title(slide, title, "Color the picture. Then finish the sentence.")

    # 画框：收窄并整体居中（用户反馈：原版式贴边、要更居中更简单）
    frame_w = 7.40
    frame_left = CONTENT_X + (CONTENT_W - frame_w) / 2.0
    frame_top = CONTENT_Y + 1.40
    frame_h = 3.25

    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(frame_left), Inches(frame_top), Inches(frame_w), Inches(frame_h),
    )
    frame.adjustments[0] = 0.03
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(*brand_rgb)
    frame.line.width = Pt(1.2)
    frame.shadow.inherit = False
    ftf = frame.text_frame
    ftf.margin_left = ftf.margin_right = Inches(0.16)
    ftf.margin_top = Inches(0.08)
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "Color the picture"
    fr.font.name = FONT
    fr.font.italic = True
    fr.font.size = Pt(12)
    fr.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)

    # 线稿图（有则等比放入画框居中，留出顶部标签空间）
    if coloring_image is not None:
        try:
            cpath = Path(coloring_image)
            if cpath.exists() and cpath.stat().st_size > 0:
                from PIL import Image as _PILImg
                with _PILImg.open(str(cpath)) as _pim:
                    iw, ih = _pim.size
                aspect = (iw / ih) if ih else 1.0
                pad = 0.22
                avail_w = frame_w - 2 * pad
                avail_h = frame_h - 0.62
                fit_w = avail_w
                fit_h = fit_w / aspect
                if fit_h > avail_h:
                    fit_h = avail_h
                    fit_w = fit_h * aspect
                off_x = frame_left + (frame_w - fit_w) / 2
                off_y = frame_top + 0.46 + (avail_h - fit_h) / 2
                slide.shapes.add_picture(
                    str(cpath), Inches(off_x), Inches(off_y),
                    width=Inches(fit_w), height=Inches(fit_h),
                )
        except Exception:
            pass

    # 一句话填空（整行居中，单 textbox：母句 + 下划线空 + 句号）——贴本书已学句型。
    stem = _color_say_stem(outline)
    sent_top = frame_top + frame_h + 0.30
    sent = slide.shapes.add_textbox(
        Inches(CONTENT_X), Inches(sent_top), Inches(CONTENT_W), Inches(0.62))
    stf = sent.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_right = 0
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    r1 = sp.add_run()
    r1.text = stem + " "
    r1.font.name = FONT
    r1.font.bold = True
    r1.font.size = Pt(24)
    r1.font.color.rgb = BLACK
    r2 = sp.add_run()
    r2.text = "____________"
    r2.font.name = FONT_BLANK  # Arial：下划线连成一条实线（Poppins 下会断成虚线）
    r2.font.size = Pt(24)
    r2.font.color.rgb = RGBColor(0x9A, 0x9A, 0x9A)
    r3 = sp.add_run()
    r3.text = " ."
    r3.font.name = FONT
    r3.font.bold = True
    r3.font.size = Pt(24)
    r3.font.color.rgb = BLACK

    # 词库（从绘本词表取，居中，给孩子挑词填空）
    words = [str(w).strip() for w in (outline.vocabulary_for_display or []) if str(w).strip()][:6]
    if words:
        wb_w = frame_w
        wb_left = CONTENT_X + (CONTENT_W - wb_w) / 2.0
        wb_top = sent_top + 0.80
        wb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(wb_left), Inches(wb_top), Inches(wb_w), Inches(0.58),
        )
        wb.adjustments[0] = 0.24
        wb.fill.solid()
        wb.fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE0)
        wb.line.color.rgb = RGBColor(*brand_rgb)
        wb.line.width = Pt(1.0)
        wb.shadow.inherit = False
        wtf = wb.text_frame
        wtf.margin_left = wtf.margin_right = Inches(0.18)
        wtf.margin_top = Inches(0.05)
        wtf.word_wrap = True
        wp = wtf.paragraphs[0]
        wp.alignment = PP_ALIGN.CENTER
        wl = wp.add_run()
        wl.text = "Word bank:  "
        wl.font.name = FONT
        wl.font.bold = True
        wl.font.size = Pt(13)
        wl.font.color.rgb = RGBColor(*brand_rgb)
        ww = wp.add_run()
        ww.text = "    ".join(words)
        ww.font.name = FONT
        ww.font.size = Pt(14)
        ww.font.color.rgb = BLACK


# ============================================================
#  Page 6 — Mind Map
# ============================================================

def _build_p6_mindmap(slide, rows: list[dict], title: str = "Reading") -> None:
    """故事复述思维导图 —— SWBST 框架（Somebody/Wanted/But/So/Then）。

    教学目的：用国际通行的『五步复述法』训练学生抓人物、目标、冲突、行动、结局，
    把读到的故事用自己的话有逻辑地概括出来（读后输出 / 写作前的结构脚手架）。
    rows 若带有 AI 提示，会作为浅灰提示词写入右侧引导问题后。
    """
    _add_title(slide, title,
               "Retell the story in five steps: Somebody \u2013 Wanted \u2013 But \u2013 So \u2013 Then.")

    # 五步：关键词 + 引导问题 + 颜色
    steps = [
        ("Somebody", "Who is the story mainly about?", MM_PURPLE),
        ("Wanted",   "What did they want to do?",      RGBColor(0x4C, 0x8F, 0xD8)),
        ("But",      "What was the problem?",          RGBColor(0xE3, 0x76, 0x35)),
        ("So",       "What did they do about it?",     MM_GREEN),
        ("Then",     "How did it end? What did we learn?", RGBColor(0xC0, 0x39, 0x6F)),
    ]

    area_top = CONTENT_Y + 1.40
    area_bottom = CONTENT_Y + CONTENT_H - 0.35
    area_h = area_bottom - area_top
    gap = 0.16
    band_h = (area_h - gap * (len(steps) - 1)) / len(steps)

    left_x = CONTENT_X + 0.45
    label_w = 2.05
    right_x = left_x + label_w + 0.25
    right_w = CONTENT_X + CONTENT_W - 0.45 - right_x

    for i, (kw, question, color) in enumerate(steps):
        y = area_top + i * (band_h + gap)

        # 左：彩色关键词卡
        lab = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left_x), Inches(y), Inches(label_w), Inches(band_h),
        )
        lab.adjustments[0] = 0.18
        lab.fill.solid()
        lab.fill.fore_color.rgb = color
        lab.line.fill.background()
        lab.shadow.inherit = False
        tf = lab.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i + 1}. {kw}"
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(18)
        r.font.color.rgb = _readable_text_rgb(color)  # Bug5：pastel 浅底用深字

        # 右：白底浅边书写区（引导问题 + 作答横线）
        ans = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(right_x), Inches(y), Inches(right_w), Inches(band_h),
        )
        ans.adjustments[0] = 0.10
        ans.fill.solid()
        ans.fill.fore_color.rgb = WHITE
        ans.line.color.rgb = color
        ans.line.width = Pt(1.5)
        ans.shadow.inherit = False
        tf2 = ans.text_frame
        tf2.margin_left = tf2.margin_right = Inches(0.16)
        tf2.margin_top = Inches(0.06)
        tf2.vertical_anchor = MSO_ANCHOR.TOP
        tf2.word_wrap = True
        pq = tf2.paragraphs[0]
        pq.alignment = PP_ALIGN.LEFT
        rq = pq.add_run()
        rq.text = question
        rq.font.name = FONT
        rq.font.size = Pt(12)
        rq.font.italic = True
        rq.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)  # B&W 打印可读
        # 作答横线
        pl = tf2.add_paragraph()
        pl.alignment = PP_ALIGN.LEFT
        pl.space_before = Pt(2)
        rl = pl.add_run()
        rl.text = "_" * 60
        rl.font.name = FONT_BLANK
        rl.font.size = Pt(13)
        rl.font.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)


# ============================================================
#  默认数据 fallback（AI 失败时也能出像样的 worksheet）
# ============================================================

def _resolve_worksheet_data(outline: BookOutline) -> dict:
    """优先用 outline._worksheet_data；否则用 outline 字段构造默认。
    在返回前统一做 v1.8 文本规整：美式拼写、答案格式、难度排序、剔除无效题型。"""
    data = getattr(outline, "_worksheet_data", None)
    if isinstance(data, dict) and data.get("match_pairs"):
        norm = _normalize_worksheet_data(data)
    else:
        norm = _normalize_worksheet_data(_build_default_data(outline))
    norm = _apply_official_vocab_to_worksheet_data(outline, norm)
    # L3+ 词汇页（猜词/填缺字母）依赖真实释义 → 保证补全，杜绝 "meaning of X" 占位
    if _level_num(getattr(outline, "level", "") or "") >= 3 and norm.get("match_pairs"):
        norm["match_pairs"] = _ensure_real_vocab_defs(norm["match_pairs"], outline)
    return norm


def _apply_official_vocab_to_worksheet_data(outline: BookOutline, data: dict) -> dict:
    official_words = _verbatim_vocab(outline, 6)
    if not official_words:
        return data
    out = dict(data or {})
    old_defs = {
        str(p.get("word", "")).strip().lower(): str(p.get("def", "")).strip()
        for p in (out.get("match_pairs") or [])
        if str(p.get("word", "")).strip()
    }
    out["match_pairs"] = [
        {"word": w, "def": old_defs.get(w.lower(), f"word from the story: {w}")}
        for w in official_words[:5]
    ]
    out["word_bank"] = official_words[:6]
    cloze = _story_cloze_fills(out.get("reading_text", ""), official_words, 4)
    if len(cloze) >= 2:
        out["fill_blanks"] = [
            {"sentence": format_sentence_answer(c["sentence"]),
             "answer": format_word_answer(c["answer"])}
            for c in cloze
        ]
    return out


# ----- v1.8 文本规整：所有英文走美式 + 答案格式 + 难度排序 -----

# 不允许的活动类型关键词（color / colour 只是涂色，不算 reading 输出）
_BANNED_PROMPT_PATTERNS = [
    "color the ", "colour the ",
    "color in ", "colour in ",
    "circle the picture",  # 仅涂圈没有输出
]


# v2.0 vocab 兜底真定义词典 — 防止 AI 输出 "meaning of X" 占位
# 数据源：常见低龄启蒙词汇 + L5-1 What Makes a Good Friend 词汇 + 高频感觉/动作/状态词
_KID_DICT: dict[str, str] = {
    # L5-1 What Makes a Good Friend
    "nervous":   "feeling worried and not calm",
    "shake":     "to move quickly from side to side or up and down",
    "recess":    "a short break at school for students to play outside",
    "a pile of": "many things lying one on top of another",
    "wooden":    "made of wood",
    # 高频形容词
    "excited":   "feeling very happy and full of energy",
    "amazed":    "feeling very surprised in a good way",
    "worried":   "feeling that something bad might happen",
    "happy":     "feeling glad and full of joy",
    "sad":       "feeling unhappy",
    "scared":    "feeling afraid",
    "tired":     "feeling like you need to rest or sleep",
    "kind":      "nice and caring to others",
    "quiet":     "making very little sound",
    "loud":      "making a lot of sound",
    "friendly":  "kind and nice to other people",
    "smart":     "able to think and learn quickly",
    "brave":     "not afraid of dangerous or difficult things",
    "shy":       "feeling not comfortable talking to new people",
    "gentle":    "kind and soft, not rough",
    # 高频动作
    "share":     "to give part of what you have to someone else",
    "help":      "to do something nice for someone who needs it",
    "listen":    "to pay attention to a sound or a person",
    "smile":     "to make a happy face with your mouth turned up",
    "laugh":     "to make a happy sound when something is funny",
    "look":      "to use your eyes to see something",
    "run":       "to move very fast on your feet",
    "walk":      "to move on your feet at a normal speed",
    "jump":      "to push yourself up into the air with your legs",
    "grab":      "to take hold of something quickly",
    "drop":      "to let something fall to the ground",
    "pick up":   "to lift something from the ground or a surface",
    "give":      "to let someone have something",
    # 高频名词
    "friend":    "someone you like and spend time with",
    "friendship": "the feeling of being friends with someone",
    "classmate": "a person in the same class as you at school",
    "kindness":  "the quality of being nice and caring",
    "teacher":   "a person whose job is to teach in a school",
    "student":   "a person who is learning at a school",
    "class":     "a group of students who learn together",
    "school":    "a place where children go to learn",
    "desk":      "a kind of table you sit at when you read or write",
    "chair":     "something you sit on",
    "book":      "pages with words and pictures joined together to read",
    "pencil":    "a thin stick used to write or draw",
    "eraser":    "a small piece of rubber used to remove pencil marks",
    "hamster":   "a small soft animal with short legs that people keep as a pet",
    # L1-L4 主题词
    "culture":   "the customs, beliefs, and ways of life of a group of people",
    "castle":    "a large old strong building where kings or queens used to live",
    "bagpipes":  "a musical instrument with a bag and pipes, played in Scotland",
    "journey":   "a trip from one place to another",
    # L3 健康主题（Book30 Stay Strong and Healthy）
    "healthy":   "strong and well, not sick",
    "protein":   "food that helps your body grow and get strong",
    "vitamin":   "something in food that keeps your body healthy",
    "vitamins":  "things in food that keep your body healthy",
    "exercise":  "moving your body to make it strong and fit",
    "energy":    "the power your body needs to move and play",
    "grain":     "a food like rice, bread, or oats",
    "grains":    "foods like rice, bread, and oats",
    "dairy":     "foods made from milk, like cheese and yogurt",
    "body":      "all the parts of a person from head to toe",
    "bodies":    "the whole of people's physical parts",
    "fruit":     "a sweet food that grows on plants, like apples",
    "vegetable": "a plant part we eat, like carrots or beans",
    "vegetables": "plant foods we eat, like carrots and beans",
    "rest":      "to stop and relax so your body feels better",
    "strong":    "having a lot of power in your body",
    "growth":    "getting bigger and taller",
    # L3 农场主题（Book42 From Farm to Fork）
    "farm":      "a place where people grow food and raise animals",
    "farmer":    "a person who grows food or raises animals on a farm",
    "crop":      "a plant that farmers grow for food",
    "crops":     "plants that farmers grow for food",
    "soil":      "the dirt that plants grow in",
    "harvest":   "to pick or gather the food when it is ready",
    "fresh":     "newly made or picked, not old",
    "market":    "a place where people buy and sell food",
    "plant":     "to put seeds in the ground so they grow",
    "seed":      "a small thing that grows into a plant",
    "seeds":     "small things that grow into plants",
    "deliver":   "to take something to where it needs to go",
    "ripe":      "ready to be picked and eaten",
    "fork":      "a tool with points used to eat food",
    # L3 seasons / weather
    "autumn":    "the season after summer, when leaves may turn yellow or red",
    "hot days":  "days when the weather is very warm",
    "cold days": "days when the weather is chilly or freezing",
    "parts of the world": "different places on Earth",
    "season":    "one part of the year, like spring or summer",
    "seasons":   "the four parts of the year: spring, summer, autumn, and winter",
    # L3-L4 common planning, places, and work words.
    "week":      "seven days from Monday to Sunday",
    "homework":  "school work that students do at home",
    "plan":      "an idea about what to do and when to do it",
    "practice":  "to do something many times to get better",
    "ocean":     "a very large area of salt water",
    "mountain":  "a very high piece of land",
    "desert":    "a dry place with very little rain",
    "travel":    "to go from one place to another",
    "doctor":    "a person who helps sick people get better",
    "chef":      "a person whose job is to cook food",
    "firefighter": "a person who helps stop fires",
    "builder":   "a person who builds or fixes buildings",
    "uniform":   "special clothes worn by people in the same group or job",
    "suit":      "a jacket and pants or skirt worn for smart clothes",
    "helmet":    "a hard hat that protects your head",
    "apron":     "clothing worn over the front of the body to keep clean",
    "boots":     "strong shoes that cover the feet and ankles",
    "hard hat":  "a strong hat that protects a worker's head",
    "kelp":      "a large brown sea plant that grows in the ocean",
    "seal":      "an ocean animal with flippers that can swim well",
    "tropical":  "from a warm place near the equator",
    "equator":   "an imaginary line around the middle of the Earth",
    "look different": "to not look the same",
    "go on walks": "to walk outside for fun or exercise",
    "far apart": "with a long distance between people or places",
    "video call": "a call where people can see and hear each other",
}


def _is_placeholder_def(def_text: str, word: str = "") -> bool:
    s = (def_text or "").strip().lower()
    return (
        not s
        or s.startswith("meaning of ")
        or s.startswith("definition of ")
        or s.startswith("word from the story")
        or s.startswith("see story")
        or (bool(word) and s == word.strip().lower())
    )


def _fix_vocab_def(word: str, def_text: str) -> str:
    """如果 def 是占位（meaning of X / definition of X / 空），用 _KID_DICT 兜底。
    仍无真实释义时返回空串（""），交由 _ensure_real_vocab_defs 走 AI 补全/留空，
    绝不把"meaning of X"占位带进交付物。"""
    if not _is_placeholder_def(def_text, word):
        return def_text
    return _KID_DICT.get(word.strip().lower(), "")


def _ensure_real_vocab_defs(pairs: list[dict], outline: BookOutline) -> list[dict]:
    """保证 match_pairs 每个词都有真实释义（L3 词汇页"猜词/填字母"依赖）：
    ① 先用 _KID_DICT 兜底；② 仍缺的批量找 AI 定义（非 mock）；③ 还缺则留空，绝不写占位。"""
    need = [p for p in pairs if _is_placeholder_def(p.get("def", ""), p.get("word", ""))]
    if not need:
        return pairs
    # ① 离线词典
    for p in need:
        kd = _KID_DICT.get(str(p.get("word", "")).strip().lower())
        if kd:
            p["def"] = kd
    still = [p for p in pairs if _is_placeholder_def(p.get("def", ""), p.get("word", ""))]
    # ② AI 批量定义
    if still:
        try:
            from ai_extractor import ai_define_words
            story = " ".join((pg.text or "").strip()
                             for pg in (getattr(outline, "pages", []) or [])
                             if (pg.text or "").strip())
            ai_defs = ai_define_words([p.get("word", "") for p in still],
                                      story, str(getattr(outline, "level", "3")))
            for p in still:
                d = ai_defs.get(str(p.get("word", "")).strip().lower())
                if d:
                    p["def"] = d
        except Exception as e:  # noqa: BLE001
            print(f"[worksheet] 词义 AI 补全失败：{e}")
    # ③ 仍缺 → 留空（不写占位）
    for p in pairs:
        if _is_placeholder_def(p.get("def", ""), p.get("word", "")):
            p["def"] = ""
    return pairs


# 词库里不该出现的"指令词"（说明这条不是单词、是任务说明，需剔除）
_BANK_STOPWORDS = {
    "sort", "into", "each", "circle", "match", "fill", "write", "draw",
    "choose", "feelings", "actions", "blank", "blanks", "category",
    "categories", "group", "groups", "using", "use", "them", "below",
    "word", "words",
}


def _clean_word_bank(words: list) -> list[str]:
    """清洗词库：只保留真正的单词（1-2 个 token、不过长、不含指令性停用词）。

    修复线上 bug：AI 偶尔把任务说明（如 'sort each word into feelings or actions'）
    塞进 word_bank，导致填空页词库里出现一整句指令、且后续填空题兜底只剩 1 题。
    """
    cleaned: list[str] = []
    seen: set[str] = set()
    for w in words or []:
        s = str(w or "").strip().strip(",.;:").strip()
        if not s:
            continue
        toks = s.split()
        low = s.lower()
        if len(toks) > 4 or len(s) > 32:
            continue
        if any(t.strip(",.").lower() in _BANK_STOPWORDS for t in toks):
            continue
        if low in seen:
            continue
        seen.add(low)
        cleaned.append(s)
    return cleaned


def _story_cloze_fills(reading_text: str, words: list, max_n: int = 4) -> list[dict]:
    """从本文里给每个目标词挖一句话做完形填空（扣住该词，挖空 → ____）。

    用户反馈：原来的兜底填空题 'I feel ____ when I see this.' 既不标准、又只有 1 题、
    版面大片留白。改成从故事原句里挖空，题目扣住本文、标准且能填满页面。
    答案取句中实际出现的词形（如 share→shared），保证句子语法通顺。
    """
    import re as _re
    sents = [s.strip() for s in _re.split(r"(?<=[.!?])\s+", reading_text or "") if s.strip()]
    out: list[dict] = []
    used: set[str] = set()
    for w in words:
        wl = str(w or "").strip().lower()
        if not wl or wl in used:
            continue
        if " " in wl:
            pat = _re.compile(r"\b(" + _re.escape(wl).replace(r"\ ", r"\s+") + r")\b", _re.I)
        else:
            pat = _re.compile(r"\b(" + _re.escape(wl) + r"(?:s|es|ed|ing|d)?)\b", _re.I)
        for s in sents:
            if len(s) > 115:   # 太长的句子不适合做填空
                continue
            m = pat.search(s)
            if m:
                blanked = s[:m.start()] + "____" + s[m.end():]
                out.append({"sentence": blanked, "answer": m.group(1).lower()})
                used.add(wl)
                break
        if len(out) >= max_n:
            break
    return out


def _fill_answer_allowed(answer: str, allowed_words: list[str]) -> bool:
    ans = format_word_answer(answer).lower()
    if not ans:
        return False
    allowed = {format_word_answer(w).lower() for w in (allowed_words or []) if str(w or "").strip()}
    if ans in allowed:
        return True
    # Allow simple story inflections when the bank word is singular/base.
    for word in allowed:
        if not word:
            continue
        if ans in {word + "s", word + "es", word + "ed", word + "ing", word + "d"}:
            return True
    return False


def _is_generic_fill(f: dict) -> bool:
    s = (f.get("sentence") or "").strip().lower()
    return ("when i see" in s) or ("see story" in s) or ("goes here" in s)


def _definition_fill_items(pairs: list[dict], words: list, max_n: int = 4) -> list[dict]:
    """词义→写词 兜底填空（扣住词义，绝不再用 'I feel ____ when I see this' 占位）。

    用真实释义造句：'<释义>: ____'，答案=目标词。优先用 match_pairs 自带释义，
    没有释义的词再用 _KID_DICT 兜底；都没有则跳过。保证每条都是有意义的词汇练习，
    且与第①页"连线匹配"版式不同（这里是看义写词），多样且不空洞。"""
    out: list[dict] = []
    used: set[str] = set()
    # 词→释义 速查（match_pairs 优先，其次离线词典）
    def_map: dict[str, str] = {}
    for p in (pairs or []):
        w = str(p.get("word", "")).strip()
        d = str(p.get("def", "")).strip()
        if w and d and not _is_placeholder_def(d, w):
            def_map[w.lower()] = d
    # 候选词序：先 word_bank/words，再 pairs
    cand: list[str] = []
    for w in list(words or []) + [p.get("word", "") for p in (pairs or [])]:
        ww = str(w or "").strip()
        if ww and ww.lower() not in used:
            cand.append(ww)
            used.add(ww.lower())
    out_seen: set[str] = set()
    for w in cand:
        wl = w.lower()
        d = def_map.get(wl) or _KID_DICT.get(wl, "")
        if not d or wl in out_seen:
            continue
        clue = d.strip().rstrip(".")
        clue = clue[0].upper() + clue[1:] if clue else clue
        out.append({"sentence": f"{clue}: ____", "answer": w})
        out_seen.add(wl)
        if len(out) >= max_n:
            break
    return out


def _normalize_worksheet_data(data: dict) -> dict:
    out = dict(data)
    # v2.0：reading_q_count 在 data 上可配置（4/6/8），默认 4
    reading_q_count = max(4, min(8, int(data.get("_reading_q_count") or 4)))

    # 1) match_pairs：word 小写无标点；def 美式 + 首字母小写（更像词典定义）+ v2.0 兜底真定义
    pairs = []
    for p in (data.get("match_pairs") or []):
        word = format_word_answer(p.get("word", ""))
        definition = _to_us_spelling(str(p.get("def", "")).strip().rstrip("."))
        definition = _fix_vocab_def(word, definition)
        if word:
            pairs.append({"word": word, "def": definition, "_len": len(word)})
    # 用户拍板（2026-06-04）：词汇匹配要显示【全部】目标词（每级 4-6 个），上限 6。
    pairs.sort(key=lambda x: (x["_len"], x["word"]))
    for p in pairs:
        p.pop("_len", None)
    out["match_pairs"] = pairs[:6]

    # 2) word_bank：小写美式
    out["word_bank"] = [format_word_answer(w) for w in (data.get("word_bank") or []) if w]

    # 3) fill_blanks：句子格式 + 答案小写
    fills = []
    for q in (data.get("fill_blanks") or []):
        sent = format_sentence_answer(q.get("sentence", ""))
        ans = format_word_answer(q.get("answer", ""))
        # 难度启发：句长升序
        fills.append({"sentence": sent, "answer": ans, "_len": len(sent)})
    # v2.0：上限 4 题（对齐官方每页 3-4 题）
    fills.sort(key=lambda x: x["_len"])
    for x in fills:
        x.pop("_len", None)
    out["fill_blanks"] = fills[:4]

    # 4) sentence_mcs：每个 option 走句子格式；过滤 color-only 等无效题
    mcs = []
    for q in (data.get("sentence_mcs") or []):
        opts = [format_sentence_answer(o) for o in (q.get("options") or []) if o]
        # 过滤无效
        if any(any(bad in o.lower() for bad in _BANNED_PROMPT_PATTERNS) for o in opts):
            continue
        if len(opts) < 2:
            continue
        mcs.append({
            "options": opts[:2],
            "correct": q.get("correct", 0),
            "_len": sum(len(o) for o in opts),
        })
    mcs.sort(key=lambda x: x["_len"])
    for x in mcs:
        x.pop("_len", None)
    out["sentence_mcs"] = mcs[:4]

    # 5) reading_text：美式
    out["reading_text"] = _to_us_spelling(str(data.get("reading_text", "")).strip())

    # 5b) v2.3：词库清洗 + 填空题改为本文完形（修复"词库混入指令串/只有1题/中间大片留白"）
    out["word_bank"] = _clean_word_bank(out.get("word_bank") or [])
    fill_words = out["word_bank"] or [p.get("word", "") for p in out.get("match_pairs") or []]
    allowed_fill_words = fill_words or [p.get("word", "") for p in out.get("match_pairs") or []]
    invalid_fills = [
        f for f in (out.get("fill_blanks") or [])
        if "____" not in (f.get("sentence") or "")
        or not _fill_answer_allowed(f.get("answer", ""), allowed_fill_words)
    ]
    need_rebuild = (len(out.get("fill_blanks") or []) < 3
                    or any(_is_generic_fill(f) for f in (out.get("fill_blanks") or []))
                    or bool(invalid_fills))
    if need_rebuild:
        cloze = _story_cloze_fills(out["reading_text"], fill_words, 4)
        if len(cloze) >= 2:
            out["fill_blanks"] = [
                {"sentence": format_sentence_answer(c["sentence"]),
                 "answer": format_word_answer(c["answer"])}
                for c in cloze
            ]
            # 词库直接用完形答案，避免占位/指令串混入
            out["word_bank"] = _clean_word_bank(
                [format_word_answer(c["answer"]) for c in cloze]
            )

    # 题量硬保证（修复 Book06/18/27/33/39/42 第②页只有 1 题）：
    # 故事完形对抽象非虚构常常只挖到 0-1 句（目标词不在原文/句子过长），导致 fill_blanks
    # 仍 <2 或残留 "I feel ____ when I see this" 占位。此处统一用【看义写词】兜底补齐到 ≥3，
    # 既杜绝单题/占位，又让填空题始终扣住本书词汇（任何级别、任何书都生效）。
    _good = [f for f in (out.get("fill_blanks") or [])
             if (f.get("sentence") or "").strip() and not _is_generic_fill(f)]
    if len(_good) < 3:
        def_items = _definition_fill_items(out.get("match_pairs") or [], fill_words, 4)
        seen = {(f.get("answer") or "").strip().lower() for f in _good}
        for it in def_items:
            ans = (it.get("answer") or "").strip().lower()
            if ans and ans not in seen:
                _good.append({"sentence": format_sentence_answer(it["sentence"]),
                              "answer": format_word_answer(it["answer"])})
                seen.add(ans)
            if len(_good) >= 4:
                break
        # 仅在能凑到 ≥2 条"有意义"题时才覆盖（否则保留原数据，绝不让题量比原来更少）。
        if len(_good) >= 2:
            out["fill_blanks"] = _good
            if not out["word_bank"]:
                out["word_bank"] = _clean_word_bank([f.get("answer") for f in _good])
    # 词库兜底：用填空答案凑齐（清洗后为空时），并打乱顺序避免与题目逐行对应
    if not out["word_bank"] and out.get("fill_blanks"):
        out["word_bank"] = _clean_word_bank([f.get("answer") for f in out["fill_blanks"]])
    if out["word_bank"]:
        import random as _rnd
        out["word_bank"] = list(dict.fromkeys(out["word_bank"]))[:6]
        _rnd.shuffle(out["word_bank"])

    # 6) reading_mcs：题干补问号 + 美式 + 首字母大写；选项 smart_format
    rmcs = []
    for q in (data.get("reading_mcs") or []):
        stem = _to_us_spelling(str(q.get("q", "")).strip().rstrip(".!?"))
        # 独立 i 大写
        import re as _re
        stem = _re.sub(r"\bi\b", "I", stem)
        stem = _re.sub(r"\bi(['\u2019])", r"I\1", stem)
        # 首字母大写
        if stem:
            stem = stem[0].upper() + stem[1:]
        if stem and stem[-1] not in "?？":
            stem += "?"
        opts = [smart_format_answer(o) for o in (q.get("options") or []) if o]
        # 过滤 color-only 活动
        stem_low = stem.lower()
        if any(bad in stem_low for bad in _BANNED_PROMPT_PATTERNS):
            continue
        if not stem or len(opts) < 2:
            continue
        rmcs.append({
            "q": stem, "options": opts[:3],
            "correct": q.get("correct", 0),
            "_len": len(stem),
        })
    # 题干越短越简单，放顶上；上限按 _reading_q_count 配置（4/6/8）
    rmcs.sort(key=lambda x: x["_len"])
    for x in rmcs:
        x.pop("_len", None)
    out["reading_mcs"] = rmcs[:reading_q_count]

    # 7) writing：theme/title/steps 美式；步骤句子格式
    writing = dict(data.get("writing") or {})
    writing["theme"] = _to_us_spelling(str(writing.get("theme", "the story")))
    writing["title"] = _to_us_spelling(str(writing.get("title", "")))
    steps_raw = writing.get("steps") or [""] * 5
    writing["steps"] = [
        format_sentence_answer(s) if s.strip() else "" for s in steps_raw
    ]
    out["writing"] = writing

    # 8) mind_map：character 短语形式（首字母大写、无句号）；problem/solution 句子形式
    mm = []
    for r in (data.get("mind_map") or [])[:5]:
        ch_raw = str(r.get("character", "")).strip().rstrip(".")
        ch = _to_us_spelling(ch_raw) if ch_raw else ""
        if ch:
            ch = ch[0].upper() + ch[1:]
        mm.append({
            "character": ch,
            "problem": format_sentence_answer(r.get("problem", "")),
            "solution": format_sentence_answer(r.get("solution", "")),
        })
    out["mind_map"] = mm

    return out


def _build_default_data(outline: BookOutline) -> dict:
    words = (
        _verbatim_vocab(outline, 5)
        or ["nervous", "shake", "recess", "wooden", "a pile of"]
    )[:5]
    words = list(words) + ["word"] * max(0, 5 - len(words))
    words = words[:5]

    pages = outline.pages or []
    story_text = capitalize_names(" ".join(
        (p.text or "").strip() for p in pages if (p.text or "").strip()
    ).strip())
    if not story_text:
        story_text = "Story text goes here."

    story_sents: list[str] = []
    for p in pages:
        if (p.text or "").strip():
            story_sents.append(capitalize_names(p.text.strip()))
    if not story_sents:
        story_sents = ["Story sentence goes here."]

    def _pair_sent(idx: int) -> dict:
        """从故事页取真句子做选项，避免 'Not X' 这种生硬兜底。"""
        correct = story_sents[idx % len(story_sents)]
        distractor_idx = (idx + 1) % len(story_sents)
        if len(story_sents) > 1 and distractor_idx == idx % len(story_sents):
            distractor_idx = (idx + 2) % len(story_sents)
        distractor = story_sents[distractor_idx]
        return {"options": [correct, distractor], "correct": 0}

    return {
        "match_pairs": [
            {"word": w, "def": f"definition of {w}"} for w in words
        ],
        "word_bank": list(words),
        "fill_blanks": [
            {"sentence": f"I _______ when I see {w}.", "answer": w} for w in words
        ],
        "sentence_mcs": [_pair_sent(i) for i in range(min(4, len(story_sents)))],
        "reading_text": story_text,
        "reading_mcs": [
            {
                "q": f"Question {i + 1}?",
                "options": ["Option A", "Option B", "Option C"],
                "correct": 0,
            }
            for i in range(8)
        ],
        "writing": _retell_writing_scaffold(outline),
        "mind_map": [
            {"character": "Main character", "problem": "Problem statement.",
             "solution": "Solution statement."},
        ] * 5,
    }


# ============================================================
#  工具
# ============================================================

def _level_label(level: str) -> str:
    s = (level or "").strip()
    if not s:
        return "Level 1"
    if s.lower().startswith("smart"):
        return "Smart"
    if s.lower().startswith("level"):
        return s
    digits = "".join(ch for ch in s if ch.isdigit())
    return f"Level {digits or '1'}"


def safe_filename(name: str) -> str:
    """与 ppt_builder.safe_filename 对齐，方便复用。"""
    import re
    safe = re.sub(r"[^\w\u4e00-\u9fff -]+", "_", name).strip("_ ")
    return safe + ".pptx"
