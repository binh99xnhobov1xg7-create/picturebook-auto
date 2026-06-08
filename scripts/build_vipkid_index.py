# -*- coding: utf-8 -*-
"""VIPKID 权威资料全量索引构建器（RAG 第一步）。

扫描 `下载/VIPKID`，把所有权威资料编成一份可检索索引：
  - docx / pdf / xlsx → 抽取文本（截断）
  - 图片(png/jpg/jpeg) → 编目录（按所在文件夹打标签，便于按主题检索官方参考图）
  - 字体(ttf) → 记录为字体资产
产物：references/vipkid_index.json（供 scripts/vipkid_rag.py 检索）

跑：py -u scripts/build_vipkid_index.py
"""
from __future__ import annotations

import json
import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "references" / "vipkid_index.json"

# 资料根目录（用户 2026-06-07 提供的权威资料）
VIPKID_DIR = Path(r"C:\Users\Jered\下载\VIPKID")

IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
DOC_EXT = {".docx", ".pdf", ".xlsx"}
FONT_EXT = {".ttf", ".otf"}
SKIP_NAMES = {"thumbs.db", ".ds_store"}
TEXT_LIMIT = 4000          # 单文档抽取文本上限
PPTX_SKIP_MB = 40          # 超过此大小的 pptx 只编目录、不抽文本

# 顶层文件夹 → 类别（按前缀匹配，宽松）
CATEGORY_BY_PREFIX = [
    ("00", "tg_rr"),                 # Teacher's Guide & Reading Report
    ("01. 字体", "font"),
    ("01.水彩2D人物", "character_art"),
    ("02. 版本PPT排版", "layout_spec"),
    ("02.Worksheet", "worksheet_sample"),
    ("03. 编辑标准", "editing_standard"),
    ("大纲", "syllabus"),
]


def _category(rel: Path) -> str:
    top = rel.parts[0] if rel.parts else ""
    for pref, cat in CATEGORY_BY_PREFIX:
        if top.startswith(pref):
            return cat
    return "misc"


def _tags(rel: Path) -> list[str]:
    """用路径里的每一段（去扩展名）作为标签，便于按主题（家庭场景/学校/Mia/Tommy…）检索。"""
    segs: list[str] = []
    for part in rel.parts:
        name = Path(part).stem
        for tok in name.replace("_", " ").replace("-", " ").split():
            tok = tok.strip()
            if tok and tok.lower() not in SKIP_NAMES:
                segs.append(tok)
    # 去重保序
    seen, out = set(), []
    for t in segs:
        if t not in seen:
            seen.add(t)
            out.append(t)
    return out[:24]


def _extract_docx(p: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(p))
        parts = [para.text for para in doc.paragraphs if para.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)[:TEXT_LIMIT]
    except Exception as e:  # noqa: BLE001
        return f"[docx 抽取失败: {e}]"


def _extract_pdf(p: Path) -> str:
    try:
        import fitz  # PyMuPDF
        text = []
        with fitz.open(str(p)) as d:
            for i, page in enumerate(d):
                if i >= 12:
                    break
                text.append(page.get_text("text"))
        return "\n".join(text)[:TEXT_LIMIT]
    except Exception as e:  # noqa: BLE001
        return f"[pdf 抽取失败: {e}]"


def _extract_xlsx(p: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(p), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# Sheet: {ws.title}")
            for ri, row in enumerate(ws.iter_rows(values_only=True)):
                if ri >= 6:
                    break
                vals = [str(c) for c in row if c not in (None, "")]
                if vals:
                    out.append(" | ".join(vals))
        wb.close()
        return "\n".join(out)[:TEXT_LIMIT]
    except Exception as e:  # noqa: BLE001
        return f"[xlsx 抽取失败: {e}]"


def _extract_pptx(p: Path) -> str:
    if p.stat().st_size > PPTX_SKIP_MB * 1024 * 1024:
        return ""  # 太大，只编目录
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        out = []
        for si, slide in enumerate(prs.slides):
            if si >= 30:
                break
            for shp in slide.shapes:
                if shp.has_text_frame and shp.text_frame.text.strip():
                    out.append(shp.text_frame.text.strip())
        return "\n".join(out)[:TEXT_LIMIT]
    except Exception:  # noqa: BLE001
        return ""


def main() -> None:
    if not VIPKID_DIR.exists():
        print(f"!! 找不到资料目录: {VIPKID_DIR}")
        sys.exit(1)

    entries: list[dict] = []
    n_img = n_doc = n_font = 0
    for dirpath, _dirs, files in os.walk(VIPKID_DIR):
        for fn in files:
            if fn.lower() in SKIP_NAMES:
                continue
            abs = Path(dirpath) / fn
            rel = abs.relative_to(VIPKID_DIR)
            ext = abs.suffix.lower()
            cat = _category(rel)
            tags = _tags(rel)
            base = {
                "rel_path": str(rel).replace("\\", "/"),
                "abs_path": str(abs),
                "category": cat,
                "ext": ext,
                "title": abs.stem,
                "tags": tags,
            }
            if ext in IMG_EXT:
                base["kind"] = "image"
                base["text"] = ""
                n_img += 1
            elif ext in FONT_EXT:
                base["kind"] = "font"
                base["text"] = ""
                n_font += 1
            elif ext == ".docx":
                base["kind"] = "doc"
                base["text"] = _extract_docx(abs)
                n_doc += 1
            elif ext == ".pdf":
                base["kind"] = "doc"
                base["text"] = _extract_pdf(abs)
                n_doc += 1
            elif ext == ".xlsx":
                base["kind"] = "doc"
                base["text"] = _extract_xlsx(abs)
                n_doc += 1
            elif ext == ".pptx":
                base["kind"] = "doc"
                base["text"] = _extract_pptx(abs)
                n_doc += 1
            else:
                continue
            entries.append(base)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"[OK] 索引 {len(entries)} 条 → {OUT}")
    print(f"     图片 {n_img} · 文档 {n_doc} · 字体 {n_font}")
    # 类别分布
    from collections import Counter
    for c, n in Counter(e["category"] for e in entries).most_common():
        print(f"     - {c}: {n}")


if __name__ == "__main__":
    main()
