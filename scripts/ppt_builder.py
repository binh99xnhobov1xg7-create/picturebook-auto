"""组装 9 页绘本 PPT：封面 + 7 故事 + 元信息页，统一 Poppins Bold。"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from config import (
    BLACK, FONT_BOLD, FONT_FAMILY, FONT_SIZE_BADGE, FONT_SIZE_BODY,
    FONT_SIZE_META_BODY, FONT_SIZE_META_HEAD, FONT_SIZE_PAGE_NUM,
    FONT_SIZE_TITLE, LIGHT_GRAY_BORDER, ORANGE_BADGE,
    PAGE_NUM_DIAMETER_IN, PAGE_NUM_MARGIN_IN,
    SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN, TEXT_BOX_PADDING_IN, TEXT_BOX_WIDTH_RATIO,
    WHITE,
)

try:  # 图像分析用于智能避让主角；缺 PIL 时优雅回退到固定角位
    from PIL import Image, ImageFilter, ImageStat
    _PIL_OK = True
except Exception:  # pragma: no cover
    _PIL_OK = False
from parser import BookOutline, PageSpec
from text_format import capitalize_names


def build_picturebook_pptx(
    outline: BookOutline,
    image_paths: list[Path],
    out_path: Path,
) -> Path:
    """images: 0=cover, 1-7=story, 8 ignored (元信息页用文字)."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    if len(image_paths) < 8:
        raise ValueError(f"需要至少 8 张图（封面+7故事），实际 {len(image_paths)}")

    # p1 封面
    _build_cover(prs.slides.add_slide(blank), outline, image_paths[0])

    # p2–p8 故事
    for i in range(1, 8):
        slide = prs.slides.add_slide(blank)
        _build_story(slide, outline.pages[i], image_paths[i], page_number=i + 1)

    # p9 元信息页
    _build_metadata(prs.slides.add_slide(blank), outline)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---------- 封面 ----------
def _build_cover(slide, outline: BookOutline, image_path: Path) -> None:
    sw, sh = Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN)

    if Path(image_path).exists():
        pic = slide.shapes.add_picture(str(image_path), 0, 0, width=sw, height=sh)
        _send_to_back(slide, pic)

    # 书名（上方居中，靠上）
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.45), Inches(SLIDE_WIDTH_IN - 1.6), Inches(1.4)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run() if p.runs else _ensure_run(p), capitalize_names(outline.title), Pt(FONT_SIZE_TITLE), BLACK)

    # 右上双徽章 Level X / Book N
    badge_w, badge_h = Inches(1.5), Inches(0.5)
    badge_x = Inches(SLIDE_WIDTH_IN - 1.5 - 0.4)
    _add_badge(slide, badge_x, Inches(0.35), badge_w, badge_h, f"Level {_clean_num(outline.level)}")
    _add_badge(slide, badge_x, Inches(0.95), badge_w, badge_h, f"Book {_clean_num(outline.book_number)}")


# ---------- 智能文字位置：分析图像，挑"最空 + 偏亮"的角位贴白底框，避开主角 ----------
_BOX_H_IN = 1.6


def _estimate_text_box_height(text: str, box_w_in: float, font_pt: float) -> int:
    """按文字实际长度估算白底框需要的高度（英寸×100→整数，便于比较）。
    白底要"贴合实际内容"：行数 = 文本宽度 / 每行可容字符数，高度 = 行数×行高 + 上下内边距。
    """
    inner_w = max(0.5, box_w_in - 2 * TEXT_BOX_PADDING_IN)
    # Poppins 估算：每字符宽 ≈ font_pt/72 * 0.52 英寸
    char_w = font_pt / 72.0 * 0.52
    cpl = max(8, int(inner_w / char_w))
    n_lines = max(1, -(-len(str(text or "").strip()) // cpl))  # 向上取整
    line_h = font_pt / 72.0 * 1.22
    return n_lines * line_h + 2 * TEXT_BOX_PADDING_IN


def _corner_xy(corner: str, box_w_in: float, box_h_in: float) -> tuple[float, float]:
    """按角位 + 实际框尺寸返回 (left_in, top_in)；底部角位用实际高度定位，不写死 1.6。"""
    margin = 0.35
    right = SLIDE_WIDTH_IN - box_w_in - margin
    bottom = SLIDE_HEIGHT_IN - box_h_in - margin
    return {
        "top-left": (margin, margin),
        "top-right": (right, margin),
        "bottom-left": (margin, bottom),
        "bottom-right": (right, bottom),
    }.get(corner, (margin, margin))


def _smart_text_corner(
    image_path: Path, box_w_in: float, box_h_in: float, fallback: str
) -> str:
    """对整张图做轻量分析，在 4 个标准角位里选出主体最少（边缘能量最低）、
    且较亮的区域来放白底文字框，从而尽量不挡住主角 / 关键事物。

    评分：score = 边缘能量(越低=越空) - 0.08*亮度(越亮越优)；取最小者。
    缺 PIL 或读图失败 → 回退到 fallback（page.text_corner）。
    """
    if not (_PIL_OK and Path(image_path).exists()):
        return fallback
    try:
        gray = Image.open(image_path).convert("L")
    except Exception:
        return fallback
    W, H = gray.size
    if W < 4 or H < 4:
        return fallback
    edges = gray.filter(ImageFilter.FIND_EDGES)

    best_corner, best_score = fallback, None
    for corner in ("top-left", "top-right", "bottom-left", "bottom-right"):
        left_in, top_in = _corner_xy(corner, box_w_in, box_h_in)
        x0 = max(0, int(left_in / SLIDE_WIDTH_IN * W))
        x1 = min(W, int((left_in + box_w_in) / SLIDE_WIDTH_IN * W))
        y0 = max(0, int(top_in / SLIDE_HEIGHT_IN * H))
        y1 = min(H, int((top_in + box_h_in) / SLIDE_HEIGHT_IN * H))
        if x1 <= x0 or y1 <= y0:
            continue
        edge_energy = ImageStat.Stat(edges.crop((x0, y0, x1, y1))).mean[0]
        brightness = ImageStat.Stat(gray.crop((x0, y0, x1, y1))).mean[0]
        score = edge_energy - 0.08 * brightness
        if best_score is None or score < best_score:
            best_corner, best_score = corner, score
    return best_corner


# ---------- 故事页 ----------
def _build_story(slide, page: PageSpec, image_path: Path, page_number: int) -> None:
    sw, sh = Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN)

    if Path(image_path).exists():
        pic = slide.shapes.add_picture(str(image_path), 0, 0, width=sw, height=sh)
        _send_to_back(slide, pic)

    # 文字框（白底圆角 + 黑色 Poppins Bold）
    # 用户拍板（2026-06-05）：白底必须"贴合实际文字"——框高随文本长度自适应，不留大片空白；
    #   并落在画面最空的角位（避开主角/关键事物）。
    txt = capitalize_names(page.text)
    box_w_in = SLIDE_WIDTH_IN * TEXT_BOX_WIDTH_RATIO
    box_h_in = _estimate_text_box_height(txt, box_w_in, float(FONT_SIZE_BODY))
    box_h_in = max(0.6, min(box_h_in, SLIDE_HEIGHT_IN - 1.2))   # 合理夹取，避免越界

    fallback_corner = page.text_corner or "top-left"
    corner = _smart_text_corner(image_path, box_w_in, box_h_in, fallback_corner)
    left_in, top_in = _corner_xy(corner, box_w_in, box_h_in)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(box_w_in), Inches(box_h_in),
    )
    try:
        box.adjustments[0] = 0.06   # 轻微圆角，更精致
    except Exception:
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*WHITE)
    box.line.color.rgb = RGBColor(*WHITE)
    box.line.width = Pt(0.0)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    # 白底随文字自适应高度：PowerPoint 里会精确贴合；预览(LibreOffice)用上面的估算高度。
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(TEXT_BOX_PADDING_IN)
    tf.margin_top = tf.margin_bottom = Inches(TEXT_BOX_PADDING_IN * 0.7)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_run(_ensure_run(p), txt, Pt(FONT_SIZE_BODY), BLACK)

    # 页码：偶数页左下，奇数页右下
    _add_page_number(slide, page_number)


# ---------- 元信息页（p9 封底）----------
# 标准模板（c:\Users\Jered\Desktop\sample 尾页.pptx）规范：
#   细深色边框白底框 ≈4.40"×2.37"，位于 (0.98, 1.62)，文字垂直居中；
#   每行【标签 Poppins SemiBold 15pt + 数值 Poppins 14pt】同行（标签加粗、数值不加粗）；
#   Vocabulary 标签单独一行，词表换行用 Poppins 14pt；字段顺序 Level/Book/CEFR/Lexile/Word count/Vocabulary。
_META_LABEL_FONT = "Poppins SemiBold"
_META_VALUE_FONT = "Poppins"
_META_LABEL_PT = 15
_META_VALUE_PT = 14


def _build_metadata(slide, outline: BookOutline) -> None:
    # L0-2 词表分 2 行（每行 4 词）→ 比 L3-6 多一行，框略加高；并开启自适应让 PPT 精确贴合。
    box_h = 2.72 if outline.is_dual_vocab_level else 2.37
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.98), Inches(1.62), Inches(4.40), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*WHITE)
    box.line.color.rgb = RGBColor(*LIGHT_GRAY_BORDER)
    box.line.width = Pt(1.0)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.06)

    def _para():
        p0 = tf.paragraphs[0]
        return p0 if (not p0.runs and not p0.text) else tf.add_paragraph()

    def _run(p, text: str, *, label: bool) -> None:
        r = p.add_run()
        r.text = text
        r.font.name = _META_LABEL_FONT if label else _META_VALUE_FONT
        r.font.bold = False
        r.font.size = Pt(_META_LABEL_PT if label else _META_VALUE_PT)
        r.font.color.rgb = RGBColor(*BLACK)

    def kv(label: str, value: str) -> None:
        p = _para()
        p.alignment = PP_ALIGN.LEFT
        _run(p, label, label=True)
        if value:
            _run(p, value, label=False)

    def value_line(text: str) -> None:
        p = _para()
        p.alignment = PP_ALIGN.LEFT
        _run(p, text, label=False)

    # Level：0 / Smart 级别显示 "Smart"，其余 1-6 显示数字（用户拍板 2026-06-08）
    level_disp = "Smart" if outline.level_key in ("smart", "0") else (_clean_num(outline.level) or "Smart")
    kv("Level: ", level_disp)
    kv("Book: ", _clean_num(outline.book_number))
    kv("CEFR: ", outline.cefr or "-")
    kv("Lexile: ", outline.lexile or "-")
    kv("Word count: ", str(outline.total_words))
    kv("Vocabulary: ", "")

    # 词表：L0-2 → 每行 4 词、分多行（通常 8 词 → 2 行）；L3-6 → 单行
    words = [w for w in outline.vocabulary_for_display if w]
    if outline.is_dual_vocab_level:
        if words:
            for i in range(0, len(words), 4):
                value_line(", ".join(words[i:i + 4]))
        else:
            value_line("-")
    else:
        value_line(", ".join(words) or "-")


# ---------- 通用 ----------
def _add_badge(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5  # 完全圆角
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*ORANGE_BADGE)
    shape.line.color.rgb = RGBColor(*BLACK)
    shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(_ensure_run(p), text, Pt(FONT_SIZE_BADGE), WHITE)


def _add_page_number(slide, page_number: int) -> None:
    d = PAGE_NUM_DIAMETER_IN
    m = PAGE_NUM_MARGIN_IN
    top = SLIDE_HEIGHT_IN - m - d
    if page_number % 2 == 0:  # 偶数页 左下
        left = m
    else:                      # 奇数页 右下
        left = SLIDE_WIDTH_IN - m - d

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(d), Inches(d)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*WHITE)
    circle.line.color.rgb = RGBColor(*BLACK)
    circle.line.width = Pt(0.75)
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = Inches(0.0)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(_ensure_run(p), str(page_number), Pt(FONT_SIZE_PAGE_NUM), BLACK)


def _ensure_run(p):
    if p.runs:
        for r in p.runs[1:]:
            r.text = ""
        p.runs[0].text = ""
        return p.runs[0]
    return p.add_run()


def _set_run(run, text: str, size, rgb_tuple: tuple[int, int, int]) -> None:
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.bold = FONT_BOLD
    run.font.size = size
    run.font.color.rgb = RGBColor(*rgb_tuple)


def _send_to_back(slide, shape) -> None:
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _clean_num(s: str) -> str:
    if not s:
        return "-"
    m = re.search(r"\d+", s)
    return m.group(0) if m else s


def safe_filename(title: str) -> str:
    name = re.sub(r"[^\w\s-]", "", title, flags=re.UNICODE)
    name = re.sub(r"\s+", "_", name.strip()) or "PictureBook"
    return f"{name}.pptx"


# ---------- 精修后原地重出 Reader.pptx（只换图、不动文字/版式/outline） ----------
def _largest_picture(slide):
    """返回该 slide 上面积最大的图片 shape（即全幅背景图），无则 None。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        return None
    return max(pics, key=lambda p: (p.width or 0) * (p.height or 0))


def swap_reader_images(
    reader_pptx: Path | str,
    img_dir: Path | str,
    out_path: Path | str | None = None,
) -> Path:
    """把现有 Reader.pptx 每页的全幅背景图替换为 img_dir 下精修后的 page_XX.png。

    slide 0=封面(page_00) … slide 7=正文(page_07)；元信息页(slide 8)无背景图，跳过。
    只替换图片二进制，文字/版式/页码/徽章全部保留，无需 outline。

    Returns: 写出的 pptx 路径（默认覆盖原文件）。
    """
    reader_pptx = Path(reader_pptx)
    img_dir = Path(img_dir)
    out_path = Path(out_path) if out_path else reader_pptx

    prs = Presentation(str(reader_pptx))
    swapped = 0
    for sidx, slide in enumerate(prs.slides):
        if sidx > 7:
            break
        new_img = img_dir / f"page_{sidx:02d}.png"
        if not new_img.exists():
            continue
        pic = _largest_picture(slide)
        if pic is None:
            continue
        embed = pic._element.blip_rId  # rId
        if not embed:
            continue
        img_part = slide.part.related_part(embed)
        img_part._blob = new_img.read_bytes()
        swapped += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
