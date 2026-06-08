# -*- coding: utf-8 -*-
"""生成「VIPKID Dino 绘本 4 件套自动化」内部汇报 PPT（v2 精简美化版，11 页）。

叙事：封面 → 全景 → 任务&要求 → 卡点 → 痛点⇄解法 → 网页方案 →
      三大能力 → Before/After → 价值 → 路线图 → 结尾
素材：outputs/_ppt_assets（4 件套真实成品 + 网页截图 + Dino 抠图）
运行：.venv\\Scripts\\python.exe scripts\\_build_share_deck.py
"""
from __future__ import annotations

from pathlib import Path
from datetime import date

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(r"C:\Users\Jered\picturebook-auto")
A = ROOT / "outputs" / "_ppt_assets"
BRAND = ROOT / "assets" / "brand"
REFSHEET = (ROOT / "outputs" / "What_Makes_a_Good_Friend_20260604_202039"
            / "images" / "_refsheets" / "sheet_p02.png")
OUT_PPTX = ROOT / "outputs" / "VIPKID_Dino_绘本自动化_汇报版.pptx"

# ---- 配色（VIPKID 橙色品牌；变量名沿用 PINK，统一改色即可） ----
PINK = RGBColor(0xF7, 0x6B, 0x1C)
PINK_DK = RGBColor(0xD4, 0x55, 0x10)
PINK_LT = RGBColor(0xFF, 0xF2, 0xE8)
INK = RGBColor(0x23, 0x27, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LGRAY = RGBColor(0xA8, 0xAE, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x12, 0xB5, 0xA8)
TEAL_DK = RGBColor(0x0B, 0x82, 0x79)
TEAL_LT = RGBColor(0xE5, 0xF7, 0xF5)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
RED = RGBColor(0xE3, 0x4D, 0x4D)
RED_LT = RGBColor(0xFC, 0xEC, 0xEC)
CARD = RGBColor(0xF6, 0xF7, 0xF9)
LINE = RGBColor(0xE7, 0xE9, 0xEE)
INK_BG = RGBColor(0x1E, 0x22, 0x2B)

CN = "Microsoft YaHei"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)


def _shadow(sp, blur=110000, dist=46000, alpha=30000, color='9CA3AF'):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad': str(blur), 'dist': str(dist),
                           'dir': '5400000', 'rotWithShape': '0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val': color})
    a = spPr.makeelement(qn('a:alpha'), {'val': str(alpha)})
    clr.append(a); sh.append(clr); el.append(sh); spPr.append(el)


def rect(s, x, y, w, h, fill=None, line=None, lw=1, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, round_=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if round_ is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = round_
        except Exception:
            pass
    if shadow:
        _shadow(sp)
    return sp


def _ea(run, name):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)


def txt(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        sa=6, ls=1.0, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0); p.line_spacing = ls
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, col, b, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = b; r.font.name = fnt
            _ea(r, fnt)
    return tb


def R(t, sz, col=INK, b=False, fnt=CN):
    return (t, sz, col, b, fnt)


def pic_fit(s, path, bx, by, bw, bh, align="center", valign="middle",
            frame=False, shadow=False):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    if ar > bw / bh:
        w = bw; h = bw / ar
    else:
        h = bh; w = bh * ar
    x = bx + (bw - w) / 2 if align == "center" else (bx if align == "left" else bx + bw - w)
    y = by + (bh - h) / 2 if valign == "middle" else (by if valign == "top" else by + bh - h)
    if frame:
        bg = rect(s, x - 0.05, y - 0.05, w + 0.1, h + 0.1, fill=WHITE)
        if shadow:
            _shadow(bg, blur=120000, alpha=34000)
    return s.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def dino(s, name, x, y, h):
    p = A / f"{name}.png"
    with Image.open(p) as im:
        iw, ih = im.size
    w = h * iw / ih
    return s.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))


def dino_corner(s):
    s.shapes.add_picture(str(BRAND / "dino_head_icon.png"),
                         Inches(SW - 0.74), Inches(0.28), height=Inches(0.46))


def header(s, kicker, title, accent=PINK, corner=True):
    rect(s, 0, 0, 0.2, SH, fill=accent)
    txt(s, 0.62, 0.42, 11.0, 0.32, [R(kicker, 12.5, accent, True)])
    txt(s, 0.6, 0.74, 11.4, 0.8, [R(title, 27, INK, True)])
    rect(s, 0.64, 1.5, 0.78, 0.055, fill=accent)
    if corner:
        dino_corner(s)


def foot(s, n):
    txt(s, 0.62, SH - 0.42, 6.0, 0.3, [R("VIPKID Dino · 绘本 4 件套自动化", 9.5, LGRAY)])
    txt(s, SW - 1.1, SH - 0.42, 0.8, 0.3, [R(f"{n:02d} / 11", 9.5, LGRAY)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, label, fill, tcol=WHITE, size=13, b=True):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
    txt(s, x, y, w, h, [R(label, size, tcol, b)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    """把逐段讲稿写进该页的演讲者备注（放映时演讲者可见，观众不可见）。"""
    s.notes_slide.notes_text_frame.text = text.strip()


def _prep_ip_proof():
    dest = A / "ip_proof.png"
    with Image.open(REFSHEET) as im:
        w, h = im.size
        crop = im.crop((0, 0, int(w * 0.30), h))
        bbox = crop.getbbox()
        if bbox:
            crop = crop.crop(bbox)
        crop.save(dest)
    return dest


# ============================================================
# 1 封面
# ============================================================
def s_cover():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    rect(s, 0, 0, 5.95, SH, fill=PINK)
    rect(s, 5.95, 0, 0.10, SH, fill=PINK_DK)
    s.shapes.add_picture(str(BRAND / "dino_reading_logo.png"), Inches(0.6), Inches(0.62),
                         height=Inches(0.62))
    txt(s, 0.62, 2.0, 5.0, 0.5, [R("AI 自动化工具 · 内部汇报", 14, PINK_LT, True)])
    txt(s, 0.6, 2.5, 5.2, 2.0,
        [R("绘本 4 件套", 40, WHITE, True), R("一键自动化", 40, WHITE, True)], ls=1.05)
    rect(s, 0.66, 4.35, 1.0, 0.06, fill=WHITE)
    txt(s, 0.62, 4.6, 5.1, 1.4,
        [R("一段故事原文 → AI 抽取 → 老师微调 → 一键 4 件套", 13.5, PINK_LT),
         R("把 IC 的「人肉生图流水线」升级为", 13.5, WHITE),
         R("可批量、标准统一、随点随出的工具。", 13.5, WHITE)], ls=1.35, sa=4)
    txt(s, 0.62, SH - 0.85, 5.0, 0.4,
        [[R("分享人 Jered", 12, PINK_LT, True), R(f"   ·   {date.today().isoformat()}", 12, PINK_LT)]])
    pic_fit(s, A / "reader_p00.png", 6.5, 1.45, 6.3, 4.6, frame=True, shadow=True)
    txt(s, 6.5, 6.2, 6.3, 0.35,
        [R("↑ 工具产出的真实绘本封面（L5《What Makes a Good Friend?》）", 10.5, LGRAY)],
        align=PP_ALIGN.CENTER)
    dino(s, "dino_wave", 5.25, 5.0, 1.95)
    notes(s, """
【0–2′ 开场 Hook】
• 做过教研的同事都知道：一本绘本全套 4 件套（绘本/Worksheet/Reading Report/Teacher's Guide），IC 配合教研复查要大半天到一天，还常返工改图。
• 今天我想用约 3 分钟，当着大家的面，从一段故事原文现场生成这一整套。
• 留一拍——如果能成，意味着我们 0–6 级绘本第一次能规模化、标准统一地铺出来。
（语气：先共鸣，再用「大半天 vs 3 分钟」制造反差，敢现场赌一把。）
""")
    return s


# ============================================================
# 2 全景一图看懂
# ============================================================
def s_overview():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "OVERVIEW", "一图看懂：它是什么")
    steps = [("①", "输入原文", "贴一段故事", PINK),
             ("②", "AI 抽取", "拆页/抽词/出题", PINK),
             ("③", "老师微调", "确认内容画面", PINK),
             ("④", "一键产出", "4 件套 ZIP", TEAL)]
    n = 4; w = 2.55; gap = 0.65; x = 0.7; y = 2.0
    for i, (num, t, d, col) in enumerate(steps):
        rect(s, x, y, w, 1.5, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1, shadow=True)
        txt(s, x + 0.25, y + 0.22, w - 0.5, 0.5, [[R(num + "  ", 18, col, True), R(t, 16, INK, True)]])
        txt(s, x + 0.27, y + 0.85, w - 0.5, 0.5, [R(d, 12, GRAY)])
        if i < n - 1:
            txt(s, x + w, y, gap, 1.5, [R("→", 22, PINK, True)], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
        x += w + gap
    txt(s, 0.7, 3.95, 11, 0.4, [R("产出「4 件套」——一次生成全套教学物料：", 14, INK, True)])
    cards = [("reader_p00.png", "Picture Book"), ("worksheet_p00.png", "Worksheet"),
             ("rr_p00.png", "Reading Report"), ("tg_p00.png", "Teacher's Guide")]
    w = 2.85; gap = 0.22; x = 0.7; y = 4.5
    for img, t in cards:
        rect(s, x, y, w, 2.25, fill=WHITE, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
        pic_fit(s, A / img, x + 0.12, y + 0.12, w - 0.24, 1.6, frame=True)
        txt(s, x, y + 1.82, w, 0.4, [R(t, 12.5, PINK_DK, True)], align=PP_ALIGN.CENTER)
        x += w + gap
    foot(s, 2)
    notes(s, """
【承接·一句话讲清是什么】
• 先用一张图让大家知道它是什么：原文 → AI 抽取 → 老师微调 → 一键 4 件套。
• 一次产出 4 件「连续」物料：绘本 / Worksheet / Reading Report / Teacher's Guide，内容互相对得上、不矛盾。
• 这条主线后面会反复出现，请大家记住它。
""")
    return s


# ============================================================
# 3 任务 & 要求
# ============================================================
def s_task():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "01 · 背景", "任务 & 要求")
    txt(s, 0.7, 1.85, 11.8, 0.6,
        [[R("基于现有故事大纲，按 ", 18, INK), R("L0 – L6", 20, PINK, True),
          R(" 七个级别，自动生成符合品牌规范的「4 件套」。", 18, INK)]])
    levels = [("Smart", RGBColor(0x5E,0x9F,0x49)), ("L1", RGBColor(0xF1,0x82,0x00)),
              ("L2", RGBColor(0x54,0xC2,0xF0)), ("L3", RGBColor(0xE9,0x46,0x53)),
              ("L4", RGBColor(0x00,0xB0,0xC4)), ("L5", RGBColor(0xE9,0x52,0x83)),
              ("L6", RGBColor(0x06,0x77,0xB7))]
    x = 0.7
    for name, col in levels:
        chip(s, x, 2.75, 1.5, 0.52, name, col, size=14)
        x += 1.65
    cons = [("🎨 统一画风", "全系列共用一套水彩童书画风"),
            ("🧍 固定 IP", "Mia / Tommy 形象跨页一致"),
            ("📐 统一规范", "VIPKID Dino 品牌母版 + 题量梯度")]
    w = 3.85; gap = 0.27; x = 0.7; y = 3.75
    for t, d in cons:
        rect(s, x, y, w, 1.4, fill=PINK_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
        txt(s, x + 0.3, y + 0.25, w - 0.6, 0.5, [R(t, 16, PINK_DK, True)])
        txt(s, x + 0.3, y + 0.78, w - 0.6, 0.5, [R(d, 12.5, GRAY)], ls=1.1)
        x += w + gap
    rect(s, 0.7, 5.5, 11.9, 1.0, fill=INK_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
    txt(s, 1.05, 5.5, 11.2, 1.0,
        [[R("核心难点 = ", 15, AMBER, True),
          R("不是「画一本」，而是「成体系、标准化地批量画 7 个级别」。", 15, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    foot(s, 3)
    notes(s, """
【2–7′ · SCQA 之 S 背景，慢讲】
• 背景：我们 Dino 线下绘本馆，每本书都要配齐 4 件套，覆盖 0–6 级，量很大。
• 难点不是「画一本」，是「成体系、标准化地批量画 7 个级别」。
• 画风统一、IP 固定（Mia/Tommy 等家族）、品牌母版规范——全都要守住，且跨册一致。
""")
    return s


# ============================================================
# 4 卡点 Before
# ============================================================
def s_pain():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "02 · 原先的卡点", "过去：全靠 IC 一步步手工做", accent=RED)
    steps = ["写提示词", "调 AI 生图", "反复返工", "对规范出题", "手工排版"]
    x = 0.7
    for i, st in enumerate(steps):
        rect(s, x, 1.85, 1.9, 0.62, fill=RED_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.25)
        txt(s, x, 1.85, 1.9, 0.62, [R(st, 13, RED, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            txt(s, x + 1.9, 1.85, 0.45, 0.62, [R("→", 16, LGRAY, True)], align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
        x += 2.35
    pains = [("⏱️", "效率低 / 产能小", "一本 4 件套全手工拼装，耗时大半天"),
             ("💰", "费用高", "每页多候选 + 反复重生，图越调越烧钱"),
             ("🎲", "质量不稳", "IP 形象漂移、画风跳变，质量随人波动"),
             ("🧩", "难复用", "题目 / 版式靠人记规范，换人就走样")]
    y = 2.85
    for i, (e, t, d) in enumerate(pains):
        col = i % 2; row = i // 2
        x = 0.7 + col * 5.55
        yy = y + row * 1.62
        rect(s, x, yy, 5.35, 1.42, fill=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.08)
        rect(s, x, yy, 0.12, 1.42, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
        txt(s, x + 0.35, yy + 0.22, 5.0, 0.55, [[R(e + "  ", 18, INK), R(t, 18, INK, True)]])
        txt(s, x + 0.35, yy + 0.82, 5.0, 0.5, [R(d, 12.5, GRAY)], ls=1.1)
    dino(s, "dino_think", 11.0, 2.95, 2.5)
    rect(s, 0.7, 6.05, 11.9, 0.66, fill=INK_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.14)
    txt(s, 1.05, 6.05, 11.2, 0.66,
        [[R("结论 = ", 14, AMBER, True),
          R("质量靠人盯、费用随返工涨 —— 几乎无法规模化。", 14, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
    foot(s, 4)
    notes(s, """
【SCQA 之 C 冲突 · 核心共鸣段，最慢讲】
• ① 人工重、周期长：全靠 IC 手工 + 教研复查，一本大半天。
• ② 修图费劲、画风/人物不稳：AI 图反复调，Mia/Tommy 跨页「跳脸」。
• ③ 出题不匹配、格式因人而异：即使有统一 SOP 文档，每个人理解不同，产出还是不一样。
• ④ 核心矛盾：时间紧，又要保质保量——靠堆人堆时间，规模化做不到。
• 关键话术：「问题不在某个人不努力，而在于质量挂在『人』身上——人一换、活一多，标准就守不住。」
""")
    return s


# ============================================================
# 5 痛点 ⇄ 解法 对照（逻辑核心）
# ============================================================
def s_bridge():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "02 → 03 · 因果", "每个卡点，都有对应的解法")
    rect(s, 0.7, 1.8, 5.4, 0.55, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.2)
    txt(s, 0.7, 1.8, 5.4, 0.55, [R("原先的卡点", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 7.2, 1.8, 5.4, 0.55, fill=TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.2)
    txt(s, 7.2, 1.8, 5.4, 0.55, [R("网页的解法", 15, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("慢：一本书耗大半天", "分钟级生成，老师只做微调"),
        ("IP 漂移、画风跳变", "参考图锁定 + 画风注入，跨页稳定"),
        ("返工多、提示词靠经验", "3 候选挑选 + 单页重生，提示词系统生成"),
        ("规范靠人记、易走样", "规范沉淀进系统，一处维护、可批量"),
    ]
    y = 2.6; rh = 0.8; gap = 0.13
    for a, b in rows:
        rect(s, 0.7, y, 5.4, rh, fill=RED_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
        txt(s, 1.0, y, 4.8, rh, [R(a, 13.5, INK, True)], anchor=MSO_ANCHOR.MIDDLE, ls=1.1)
        txt(s, 6.12, y, 1.06, rh, [R("➜", 20, PINK, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, 7.2, y, 5.4, rh, fill=TEAL_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
        txt(s, 7.5, y, 4.8, rh, [R(b, 13.5, TEAL_DK, True)], anchor=MSO_ANCHOR.MIDDLE, ls=1.1)
        y += rh + gap
    txt(s, 0.7, y + 0.02, 11.9, 0.45,
        [[R("一句话：", 14, PINK_DK, True),
          R("把规范/IP/画风/题型都沉淀进工具，让人从「执行」回到「判断」。", 14, GRAY)]])
    foot(s, 5)
    notes(s, """
【7–9′ · SCQA 之 Q 提问 + A 一句话回答】
• Q：能不能让产出「不依赖个人手感」，又快、又稳、又标准？
• A（一句话）：能。把规范 / IP / 画风 / 题型全部沉淀进一个网页工具——一段原文进去，AI 抽取、老师微调、一键出 4 件套。
• 逐行念左右对照，强调「每个卡点都有对应解法」的因果关系。
""")
    return s


# ============================================================
# 6 网页方案
# ============================================================
def s_solution():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "03 · 解决方案", "我的网页：把流程全部装进去")
    pic_fit(s, A / "webapp_landing.png", 6.5, 1.85, 6.35, 4.6, align="center", valign="top", frame=True, shadow=True)
    txt(s, 6.5, 6.5, 6.35, 0.3, [R("↑ 网页内置「底层逻辑」规则面板（运行实拍）", 10.5, LGRAY)], align=PP_ALIGN.CENTER)
    feats = [("AI 自动抽取", "贴原文 → 自动拆页/抽词/出题/写画面"),
             ("底层逻辑内置", "规范/字体/题量/品牌色写进系统"),
             ("IP + 画风锁定", "参考图锁形象，每页只改表情动作"),
             ("一键 4 件套", "微调确认 → 打包 ZIP 下载")]
    y = 1.95
    for t, d in feats:
        rect(s, 0.7, y, 5.5, 1.02, fill=PINK_LT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.1)
        txt(s, 1.0, y + 0.15, 5.0, 0.5, [R("✓  " + t, 15.5, PINK_DK, True)])
        txt(s, 1.35, y + 0.6, 4.7, 0.4, [R(d, 11.5, GRAY)])
        y += 1.16
    dino(s, "dino_point", 5.0, 4.55, 1.75)
    foot(s, 6)
    notes(s, """
【9–11′ · 方案概览】
• 网页把整条流程装进去：AI 自动抽取 / 底层逻辑内置 / IP+画风锁定 / 一键 4 件套。
• 强调：右图是运行实拍，右侧那块是内置的「底层逻辑」规则面板——规范是写进系统的，不靠人记。
• 过渡到下一页：「具体怎么保证又快又稳又准？靠三大能力。」
""")
    return s


# ============================================================
# 7 三大核心能力
# ============================================================
def s_caps():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "03 · 核心能力", "三大能力，保证「又快又稳又准」")
    caps = [
        ("ip_proof.png", "① 锁死 IP", "定妆参考图锁形象", "固定 IP 家族 · 按级别配年龄，每页只改表情动作", PINK),
        ("reader_p05.png", "② 统一画风", "gpt-image-2 · 暖米水彩", "画风注入每页提示词，低随机不乱发挥", TEAL),
        ("worksheet_p00.png", "③ 规范落地", "题量/词汇/品牌色自动套", "换人不走样，质量有下限", AMBER),
    ]
    w = 3.85; gap = 0.27; x = 0.7; y = 1.95
    for img, t, sub, d, col in caps:
        rect(s, x, y, w, 4.55, fill=WHITE, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05, shadow=True)
        rect(s, x, y, w, 0.12, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
        pic_fit(s, A / img, x + 0.2, y + 0.32, w - 0.4, 2.35, frame=True)
        txt(s, x + 0.3, y + 2.85, w - 0.6, 0.5, [R(t, 18, col, True)])
        txt(s, x + 0.3, y + 3.35, w - 0.6, 0.45, [R(sub, 13.5, INK, True)])
        txt(s, x + 0.3, y + 3.8, w - 0.6, 0.7, [R(d, 12, GRAY)], ls=1.2)
        x += w + gap
    foot(s, 7)
    notes(s, """
【11–14′ · 三大能力】
• ① 锁死 IP：定妆参考图，固定 IP 家族 + 按级别配年龄（L0–3/4–5/6 三套年龄），每页只改表情动作。
• ② 统一画风：gpt-image-2 注入暖米水彩，低随机、不乱发挥。
• ③ 规范落地：题量 / 词汇 / 品牌色自动套，换人不走样、质量有下限。

——————————————————————
【14–22′ · 切到网页现场 Demo（全场高潮，真机演示）】
1) 贴一段故事原文 → 2) 展示 AI 抽取结果（拆页 / 抽词 / 出题 / 画面提示词）→ 3) 现场微调一处 → 4) 点生成 → 5) 出 4 件套。
• 边演边强调：「注意——我没写代码、没调参数，只是确认和微调。」
• ⚠️ 防翻车（重要）：开场前先跑通一次、留一份已生成好的成品；若网络/接口慢，立刻切到成品继续讲，节奏不停。提前确认网页服务已启动（localhost:8501）、网络可用、浏览器已打开到首页。
""")
    return s


# ============================================================
# 8 Before / After
# ============================================================
def s_compare():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "04 · 收益", "Before / After 对比")
    rows = [
        ("维度", "原先（IC 手工）", "现在（网页工具）"),
        ("出一本 4 件套", "大半天 ～ 1 天", "分钟级生成 + 微调"),
        ("IP 一致性", "靠人盯，易漂移", "参考图锁定，跨页稳定"),
        ("画风稳定性", "忽明忽暗、随机大", "统一注入、低随机"),
        ("规范落地", "靠人记、易走样", "系统内置、一处维护"),
        ("可批量 / 复用", "几乎不可", "支持批量生产"),
    ]
    x0 = 0.7; y0 = 1.95; w = 11.9
    c0, c1, c2 = 3.2, 4.35, 4.35
    rh = 0.82
    for ri, (a, b, c) in enumerate(rows):
        y = y0 + ri * rh
        if ri == 0:
            rect(s, x0, y, w, rh, fill=PINK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.04)
            cols = [(a, c0, WHITE, True), (b, c1, WHITE, True), (c, c2, WHITE, True)]; sz = 14
        else:
            rect(s, x0, y + 0.02, w, rh - 0.04, fill=CARD if ri % 2 else WHITE)
            cols = [(a, c0, INK, True), (b, c1, GRAY, False), (c, c2, TEAL_DK, True)]; sz = 13.5
        cx = x0
        for (t, cw, col, bb) in cols:
            txt(s, cx + 0.3, y, cw - 0.4, rh, [R(t, sz, col, bb)], anchor=MSO_ANCHOR.MIDDLE)
            cx += cw
    foot(s, 8)
    notes(s, """
【22–24′ · Demo 回来收口 · Before/After】
• 用这页一眼对比：大半天 → 分钟级；易漂移 → 锁定；忽明忽暗 → 统一低随机；靠人记 → 系统内置；几乎不可批量 → 可批量。
• 话术：「同样一件事，区别不在更努力，而在标准从『人』搬到了『系统』。」
""")
    return s


# ============================================================
# 9 价值
# ============================================================
def s_value():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "04 · 价值", "它真正解决了什么")
    cards = [("⚡", "提效", "一天的活压到分钟级", "IC 从「画图工」回到「内容判断者」", PINK),
             ("🎯", "标准", "规范沉淀进系统", "换人不走样，质量有下限保证", TEAL),
             ("📈", "规模", "支持批量生产", "L0–L6 成体系铺量成为可能", AMBER)]
    w = 3.7; gap = 0.3; x = 0.7; y = 2.1
    for e, t, big, d, col in cards:
        rect(s, x, y, w, 3.5, fill=WHITE, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05, shadow=True)
        rect(s, x, y, w, 0.95, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
        rect(s, x, y + 0.5, w, 0.45, fill=col)
        txt(s, x, y + 0.08, w, 0.85, [[R(e + "  ", 22, WHITE), R(t, 22, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + 0.25, y + 1.25, w - 0.5, 0.9, [R(big, 17, INK, True)], align=PP_ALIGN.CENTER, ls=1.15)
        txt(s, x + 0.3, y + 2.35, w - 0.6, 1.0, [R(d, 12.5, GRAY)], align=PP_ALIGN.CENTER, ls=1.3)
        x += w + gap
    dino(s, "dino_cheer", 11.05, 5.45, 1.7)
    foot(s, 9)
    notes(s, """
【24–26′ · 价值 + 诚实复盘（最赢领导信任的一段）】
• 三支柱：提效（一天→分钟级）/ 标准（沉淀进系统，换人不走样）/ 规模（L0–L6 成体系铺量成为可能）。
• 做到了：一次产出 4 件连续物料、格式规范统一、文本内容稳定 85%+（目标冲 90%）。
• 坦诚还没完全解决：① 高级别复杂画面偶有偏差，需人工微调；② 文本正确率还在冲 90%+；③ 批量稳定性待大规模验证。
• 收尾话术：「我不夸大——它现在是『好用的助手』，还不是『全自动』，但方向已经跑通。」
""")
    return s


# ============================================================
# 10 路线图
# ============================================================
def s_roadmap():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "05 · 下一步", "下一步计划")
    txt(s, 0.7, 1.7, 11.8, 0.4,
        [R("聚焦三个可验证的目标，先把「能用、稳定、可放量」跑通：", 14.5, GRAY)])
    cards = [("🎯", "正确率", "≥ 85%", "内容 / 版式 / 题目一次产出即正确", PINK),
             ("🔒", "稳定率", "≥ 85%", "同一输入稳定复现，IP / 画风不跳变", TEAL),
             ("🏭", "批量生产", "可测试", "跑通批量流程，验证规模化可行性", AMBER)]
    w = 3.7; gap = 0.3; x = 0.7; y = 2.35
    for e, t, big, d, col in cards:
        rect(s, x, y, w, 3.7, fill=WHITE, line=LINE, lw=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.05, shadow=True)
        rect(s, x, y, w, 1.0, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.06)
        rect(s, x, y + 0.55, w, 0.45, fill=col)
        txt(s, x, y + 0.1, w, 0.85, [[R(e + "  ", 22, WHITE), R(t, 22, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x, y + 1.3, w, 0.9, [R(big, 38, INK, True)], align=PP_ALIGN.CENTER)
        rect(s, x + w / 2 - 0.4, y + 2.45, 0.8, 0.03, fill=col)
        txt(s, x + 0.35, y + 2.65, w - 0.7, 0.9, [R(d, 13, GRAY)], align=PP_ALIGN.CENTER, ls=1.25)
        x += w + gap
    dino(s, "dino_cheer", 11.05, 5.5, 1.65)
    foot(s, 10)
    notes(s, """
【26–28′ · 下一步目标】
• 三个可验证目标：正确率 ≥85% · 稳定率 ≥85% · 跑通批量生产测试。
• 口径保守：当前文本稳定 85%+，下一步目标冲 90%；先把「能用、稳定、可放量」跑通。
""")
    return s


# ============================================================
# 11 结尾
# ============================================================
def s_end():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=PINK)
    rect(s, 0, SH - 0.45, SW, 0.45, fill=PINK_DK)
    s.shapes.add_picture(str(BRAND / "dino_reading_logo.png"), Inches(0.6), Inches(0.6), height=Inches(0.6))
    txt(s, 0, 2.4, SW, 1.2, [R("谢谢 · Q & A", 42, WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, 0, 3.75, SW, 0.6, [R("让 IC 从「人肉生图工具」回到「内容判断者」", 18, WHITE, True)], align=PP_ALIGN.CENTER)
    txt(s, 0, 4.45, SW, 0.6, [R("工具已开放给大家试用 · 欢迎随时反馈，我会持续打磨", 15, PINK_LT)], align=PP_ALIGN.CENTER)
    dino(s, "dino_cheer", 6.0, 5.0, 1.95)
    notes(s, """
【28–30′ · 结尾】
• 定位：这工具是给大家用的——同事们用起来，随时反馈问题，我来持续优化，一起把它打磨好。（不向领导提资源诉求，强调赋能、共建。）
• 金句：让 IC 从「人肉生图工具」回到「内容判断者」，让标准长在系统里、不长在某个人身上。
• 谢谢，进入 Q&A。
""")
    return s


def main():
    _prep_ip_proof()
    s_cover(); s_overview(); s_task(); s_pain(); s_bridge(); s_solution()
    s_caps(); s_compare(); s_value(); s_roadmap(); s_end()
    prs.save(str(OUT_PPTX))
    print(f"SAVED {OUT_PPTX}  slides={len(prs.slides._sldIdLst)}")


if __name__ == "__main__":
    main()
